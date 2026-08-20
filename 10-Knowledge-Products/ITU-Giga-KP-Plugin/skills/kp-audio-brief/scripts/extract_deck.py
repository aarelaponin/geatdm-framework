#!/usr/bin/env python3
"""Extract a KP video deck into the shape the audio brief needs.

Usage:
    python3 extract_deck.py <deck.pptx> [--json] [--budget SECONDS]

Prints, per slide: the visible text frames (in reading order) and the speaker
notes (which in the KP decks carry the voice-over and the production notes).
With --budget it also proposes a per-slide time allocation, weighted by the
amount of narration each slide actually carries, so the brief starts from a
defensible split instead of an even one.

Pure python-pptx; no LibreOffice needed.
"""

import argparse
import json
import re
import sys

try:
    from pptx import Presentation
except ImportError:
    sys.exit("python-pptx not installed:  pip install python-pptx --break-system-packages")

# Title and Sources slides are fixed-cost bookends, not narration-weighted.
TITLE_SECONDS = 15
SOURCES_SECONDS = 10


def shape_text(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append(sh.text_frame.text.strip())
    return out


def notes_text(slide):
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def vo_words(notes):
    """Word count of the VO portion of the notes, ignoring production directions."""
    if not notes:
        return 0
    # The KP convention marks narration with "VO:" / "VO, slide N:".
    m = re.search(r"\bVO\b[^:]*:(.*)", notes, re.S)
    body = m.group(1) if m else notes
    # Drop bracketed stage directions and the RETRIEVAL MOMENT block header.
    body = re.sub(r"RETRIEVAL MOMENT[^\n]*", " ", body)
    return len(body.split())


def is_sources(texts):
    return bool(texts) and texts[0].strip().lower().startswith("sources")


def budget(slides, total_seconds):
    n = len(slides)
    alloc = [None] * n
    fixed = 0
    weighted = []
    for i, s in enumerate(slides):
        if i == 0:
            alloc[i] = TITLE_SECONDS
            fixed += TITLE_SECONDS
        elif is_sources(s["texts"]):
            alloc[i] = SOURCES_SECONDS
            fixed += SOURCES_SECONDS
        else:
            weighted.append(i)
    remaining = max(total_seconds - fixed, 0)
    weights = [max(slides[i]["vo_words"], 1) for i in weighted]
    tw = sum(weights)
    for i, w in zip(weighted, weights):
        alloc[i] = int(round(remaining * w / tw / 5.0)) * 5  # round to 5s
    # Reconcile rounding drift onto the largest content slide.
    drift = total_seconds - sum(alloc)
    if weighted and drift:
        biggest = max(weighted, key=lambda i: alloc[i])
        alloc[biggest] += drift
    return alloc


def mmss(sec):
    return f"{sec // 60}:{sec % 60:02d}"


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--budget", type=int, default=None,
                    help="total target runtime in seconds (e.g. 240)")
    args = ap.parse_args()

    prs = Presentation(args.deck)
    slides = []
    for i, s in enumerate(prs.slides, 1):
        texts = shape_text(s)
        notes = notes_text(s)
        slides.append({
            "n": i,
            "title": texts[0] if texts else f"(slide {i})",
            "texts": texts,
            "notes": notes,
            "vo_words": vo_words(notes),
        })

    if args.budget:
        alloc = budget(slides, args.budget)
        clock = 0
        for s, a in zip(slides, alloc):
            s["seconds"] = a
            s["start"] = clock
            clock += a

    if args.json:
        print(json.dumps({"deck": args.deck, "slides": slides}, indent=2))
        return

    print(f"# {args.deck}\n# {len(slides)} slides")
    if args.budget:
        print(f"# target runtime {mmss(args.budget)}")
    for s in slides:
        print("=" * 70)
        head = f"SLIDE {s['n']} — {s['title']}"
        if args.budget:
            head += (f"   [{mmss(s['start'])}–{mmss(s['start'] + s['seconds'])}"
                     f" · {s['seconds']}s · {s['vo_words']} VO words]")
        print(head)
        for t in s["texts"][1:]:
            print("  · " + t.replace("\n", "\n    "))
        if s["notes"]:
            print("  [NOTES] " + s["notes"].replace("\n", "\n            "))
    if args.budget:
        print("=" * 70)
        print("CUE FILE DRAFT (paste into KP*_Cues_*.txt, then verify against the take):")
        for s in slides:
            print(f"{mmss(s['start'])}   # slide {s['n']} — {s['title']}")


if __name__ == "__main__":
    main()
