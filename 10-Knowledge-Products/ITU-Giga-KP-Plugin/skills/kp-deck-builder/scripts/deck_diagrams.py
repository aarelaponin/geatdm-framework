#!/usr/bin/env python3
"""Structural diagram helpers for KP module decks — a sibling of deck_lib.py.

Shape-only visuals permitted by the ITU text-only rule: proportional bars, entity graphs,
flows with a ruling diamond, stacks, wave timelines, play cards. No icons, no imagery.
Added after the M1–M6 v0.1 review; kept in their own module so deck_lib.py's slide
furniture stays small. Import alongside deck_lib:

    from deck_lib import ...
    from deck_diagrams import bars_slide, entity_graph, gate_flow, play_card, ...
"""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from deck_lib import (
    GREY, INK, ITU_BLUE, ITU_BLUE_DARK, LIGHT, MIDGREY, PANEL_GREY, WHITE, LAYOUT_WHITE,
    add_slide, box, footer, hline, notes, num_chip, panel, set_text, solid, title)
from lxml import etree
from pptx.enum.shapes import MSO_CONNECTOR


def arrow(slide, x1, y1, x2, y2, color=ITU_BLUE_DARK, weight_pt=1.5, head=True):
    """Straight connector with an optional triangle arrowhead at (x2, y2)."""
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(weight_pt)
    if head:
        ln = c.line._get_or_add_ln()
        etree.SubElement(ln, qn('a:tailEnd'), type='triangle', w='med', len='med')
    return c


def label(slide, x, y, w, h, text, size=11, bold=False, color=GREY, align=PP_ALIGN.CENTER,
          fill=None, anchor=MSO_ANCHOR.MIDDLE, italic=False):
    """Small caption box; give it a fill to sit on top of a connector line."""
    if fill is not None:
        p = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        solid(p, fill)
        tf = p.text_frame
    else:
        tf = box(slide, x, y, w, h).text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    set_text(tf, [[(text, size, bold, color, italic)]], align=align)
    tf.vertical_anchor = anchor
    return tf


def node(slide, x, y, w, h, heading, lines=(), fill=LIGHT, ink=INK, head_ink=ITU_BLUE_DARK,
         head_size=13, size=12.5, radius=0.08, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line=None):
    """A labelled box for flow / graph diagrams."""
    p = panel(slide, x, y, w, h, fill, line=line, radius=radius)
    tf = p.text_frame
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.06)
    paras = [[(heading, head_size, True, head_ink, False)]]
    paras += [[(ln, size, False, ink, False)] for ln in lines]
    set_text(tf, paras, align=align, space_after=Pt(3))
    tf.vertical_anchor = anchor
    return p


def bar_rows(slide, rows, x=0.72, y=1.7, w=11.9, label_w=2.9, row_h=0.62, gap=0.42,
             group_gap=0.55, caption_w=0.0):
    """Proportional bars — length IS the argument. No figures, only ratios.
    rows: list of items; each item is either ('GROUP', text) for a group heading, or
    (row_label, segments, caption) with segments = [(seg_text, frac, fill, ink), ...] where
    frac is a fraction of the full bar width and seg_text may be ''."""
    bar_x = x + label_w
    bar_w = w - label_w - caption_w
    cy = y
    for item in rows:
        if item[0] == 'GROUP':
            tb = box(slide, x, cy, w, 0.36)
            set_text(tb.text_frame, [[(item[1], 12, True, GREY, False)]])
            cy += 0.42
            continue
        rl, segs, caption = item
        tb = box(slide, x, cy, label_w - 0.15, row_h)
        set_text(tb.text_frame, [[(rl, 14, True, INK, False)]])
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        sx = bar_x
        for (st, frac, fill, ink) in segs:
            sw = max(0.02, bar_w * frac)
            r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sx), Inches(cy + 0.06),
                                       Inches(sw - 0.03), Inches(row_h - 0.12))
            solid(r, fill, line_color=WHITE)
            if st:
                tf = r.text_frame
                tf.margin_left = tf.margin_right = Inches(0.06)
                tf.margin_top = tf.margin_bottom = 0
                set_text(tf, [[(st, 11, True, ink, False)]], align=PP_ALIGN.CENTER)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.word_wrap = True
            sx += sw
        if caption:
            tb = box(slide, sx + 0.1, cy, max(1.0, x + w - sx - 0.1), row_h)
            set_text(tb.text_frame, [[(caption, 12, False, GREY, True)]])
            tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cy += row_h + gap
    return cy


def bars_slide(prs, head, rows, closing, tag, note, **kw):
    """Headline + proportional bars + closing line. See bar_rows for `rows`."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    bar_rows(s, rows, **kw)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def entity_graph(prs, head, closing, tag, note,
                 names=('Capability', 'Service', 'Application', 'Data Domain',
                        'Technology Component', 'Organisation'),
                 defs=('what a body can do', 'how it reaches a citizen or body',
                       'the software that supports it', 'a kind of information, one owner',
                       'the infrastructure underneath', 'owns each of them'),
                 rels=('delivered by', 'supported by', 'uses', 'runs on', 'owns')):
    """The PAERA metamodel as a graph: four entity nodes with labelled arrows, Technology as
    the base everything runs on, Organisation as the owner bracketing them all."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    cap, svc, app, dom, tech, org = names
    nw, nh = 2.9, 1.05
    # column for Organisation on the left, then the 2x2 of entities, Technology underneath
    ox, oy, ow, oh = 0.72, 1.75, 2.0, 4.2
    node(s, ox, oy, ow, oh, org, [defs[5]], fill=PANEL_GREY, head_size=15, size=12.5)
    gx = 3.55
    x1, x2 = gx, gx + 5.4
    y1, y2 = 1.75, 3.55
    node(s, x1, y1, nw, nh, cap, [defs[0]], head_size=15)
    node(s, x2, y1, nw, nh, svc, [defs[1]], head_size=15)
    node(s, x1, y2, nw, nh, app, [defs[2]], head_size=15)
    node(s, x2, y2, nw, nh, dom, [defs[3]], head_size=15)
    # technology base under the 2x2
    ty = 5.05
    node(s, x1, ty, x2 + nw - x1, 0.9, tech, [defs[4]], fill=ITU_BLUE_DARK, ink=WHITE,
         head_ink=WHITE, head_size=15)
    # arrows + relationship labels
    arrow(s, x1 + nw, y1 + nh / 2, x2, y1 + nh / 2)
    label(s, x1 + nw + 0.55, y1 + nh / 2 - 0.2, 1.4, 0.4, rels[0], size=12, bold=True,
          color=ITU_BLUE_DARK, fill=WHITE)
    arrow(s, x1 + nw / 2, y1 + nh, x1 + nw / 2, y2)
    label(s, x1 + nw / 2 + 0.12, y1 + nh + 0.17, 1.5, 0.4, rels[1], size=12, bold=True,
          color=ITU_BLUE_DARK, align=PP_ALIGN.LEFT, fill=WHITE)
    arrow(s, x1 + nw, y2 + nh / 2, x2, y2 + nh / 2)
    label(s, x1 + nw + 0.9, y2 + nh / 2 - 0.2, 0.9, 0.4, rels[2], size=12, bold=True,
          color=ITU_BLUE_DARK, fill=WHITE)
    # everything runs on technology: two short arrows down into the base
    for xx in (x1 + nw / 2, x2 + nw / 2):
        arrow(s, xx, y2 + nh, xx, ty)
    label(s, x1 + nw + 0.3, y2 + nh + 0.05, 2.1, 0.36, 'everything ' + rels[3], size=12, bold=True,
          color=ITU_BLUE_DARK, fill=WHITE)
    # organisation owns each: one arrow into the grid with a label
    arrow(s, ox + ow, oy + oh / 2, x1, oy + oh / 2)
    label(s, ox + ow + 0.02, oy + oh / 2 - 0.42, 0.8, 0.36, rels[4], size=12, bold=True,
          color=ITU_BLUE_DARK, fill=WHITE)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def prompt_anatomy(prs, head, parts, closing, tag, note, panel_heading='THE PROMPT, AS PASTED'):
    """One prompt shown as one block, its four parts bracketed and labelled.
    parts: list of (label, example_line)."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    n = len(parts)
    top, bottom = 1.65, 6.15
    px, pw = 6.55, 6.07
    panel(s, px, top, pw, bottom - top, PANEL_GREY, radius=0.03)
    tb = box(s, px + 0.25, top + 0.12, pw - 0.5, 0.35)
    set_text(tb.text_frame, [[(panel_heading, 10.5, True, GREY, False)]])
    rows_top = top + 0.55
    rh = (bottom - rows_top) / n
    for i, (lab, example) in enumerate(parts):
        y = rows_top + i * rh
        num_chip(s, 0.72, y + 0.08, i + 1)
        tb = box(s, 1.22, y, 4.9, rh - 0.1)
        set_text(tb.text_frame, [[(lab, 16, True, INK, False)]])
        # bracket: a short blue bar on the panel's left edge for this part
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px + 0.12), Inches(y + 0.1), Inches(0.06), Inches(rh - 0.25))
        solid(r, ITU_BLUE)
        hline(s, 6.2, y + 0.24, 0.3, color=ITU_BLUE, weight_pt=1.5)
        tb = box(s, px + 0.32, y + 0.02, pw - 0.55, rh - 0.1)
        set_text(tb.text_frame, [[(example, 15, False, INK, True)]])
        if i < n - 1:
            hline(s, 0.72, y + rh - 0.05, 5.4)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def gate_flow(prs, head, questions, closing, tag, note, project=('A project proposal', 'before funding'),
              pass_box=('PASS', 'Funded — consume the shared blocks'),
              exc_box=('EXCEPTION', 'Granted in writing, with an expiry date and a reason'),
              log_box=('DECISION LOG', 'Every ruling and every exception, with its reason — into the repository')):
    """The architecture review gate as one mechanism: proposal → the questions → the Board's
    ruling → pass or time-boxed exception → the decision log."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    y0, h = 1.75, 4.3
    mid = y0 + h / 2
    # project box
    node(s, 0.72, mid - 0.7, 1.75, 1.4, project[0], [project[1]], fill=PANEL_GREY, head_size=13,
         anchor=MSO_ANCHOR.MIDDLE)
    # gate panel with numbered questions
    gx, gw = 2.95, 4.55
    p = panel(s, gx, y0, gw, h, LIGHT, radius=0.05)
    tb = box(s, gx + 0.2, y0 + 0.1, gw - 0.4, 0.35)
    set_text(tb.text_frame, [[('THE GATE — THE SAME QUESTIONS, EVERY TIME', 10.5, True, ITU_BLUE_DARK, False)]])
    qy = y0 + 0.55
    qh = (h - 0.65) / len(questions)
    for i, q in enumerate(questions):
        num_chip(s, gx + 0.2, qy + i * qh + 0.04, i + 1, d=0.3)
        tb = box(s, gx + 0.6, qy + i * qh - 0.02, gw - 0.75, qh)
        set_text(tb.text_frame, [[(q, 12.5, False, INK, False)]])
    arrow(s, 2.47, mid, gx, mid)
    # ruling diamond
    dx, dw = 7.9, 1.35
    d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(dx), Inches(mid - 0.75), Inches(dw), Inches(1.5))
    solid(d, ITU_BLUE_DARK)
    tf = d.text_frame
    tf.margin_left = tf.margin_right = 0
    set_text(tf, [[('The Board', 11, True, WHITE, False)], [('rules', 11, True, WHITE, False)]], align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    arrow(s, gx + gw, mid, dx, mid)
    # outcomes
    ox, ow, oh = 9.6, 1.55, 1.45
    node(s, ox, y0 + 0.15, ow, oh, pass_box[0], [pass_box[1]], fill=LIGHT, head_size=12, size=11)
    node(s, ox, y0 + h - oh - 0.15, ow, oh, exc_box[0], [exc_box[1]], fill=PANEL_GREY, head_size=12, size=11)
    arrow(s, dx + dw, mid, ox, y0 + 0.15 + oh / 2)
    arrow(s, dx + dw, mid, ox, y0 + h - 0.15 - oh / 2)
    # decision log
    lx, lw = 11.45, 1.17
    node(s, lx, mid - 1.1, lw, 2.2, log_box[0], [log_box[1]], fill=ITU_BLUE_DARK, ink=WHITE, head_ink=WHITE,
         head_size=11, size=10.5)
    arrow(s, ox + ow, y0 + 0.15 + oh / 2, lx, mid - 0.4)
    arrow(s, ox + ow, y0 + h - 0.15 - oh / 2, lx, mid + 0.4)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def stack_slide(prs, head, top_block, base_block, side_note, closing, tag, note):
    """One block standing on a wider base — 'what Wave 1 really builds'.
    top_block / base_block: (heading, [lines])."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    node(s, 2.9, 1.9, 5.2, 1.6, top_block[0], top_block[1], fill=LIGHT, head_size=15, size=13)
    node(s, 0.72, 3.7, 9.55, 1.75, base_block[0], base_block[1], fill=ITU_BLUE_DARK, ink=WHITE,
         head_ink=WHITE, head_size=15, size=13)
    # bracket note on the right
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.55), Inches(1.9), Inches(0.05), Inches(3.55))
    solid(r, ITU_BLUE)
    tb = box(s, 10.75, 1.9, 1.9, 3.55)
    set_text(tb.text_frame, [[(ln, 12.5, False, GREY, True)] for ln in side_note], space_after=Pt(8))
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def wave_timeline(prs, head, waves, board_label, platform_label, closing, tag, note):
    """Rollout as a wave roadmap: each wave a block whose width and height shrink as the shared
    platforms carry more of the work; a Board bar governs the whole pipeline; a platform bar
    grows underneath. waves: list of (heading, [lines], width_frac, height_frac)."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    x0, w = 0.72, 11.9
    # governance bar across the top
    b = panel(s, x0, 1.7, w, 0.5, ITU_BLUE_DARK, radius=0.3)
    tf = b.text_frame
    set_text(tf, [[(board_label, 12, True, WHITE, False)]], align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # waves
    base_y, max_h = 5.35, 2.85
    gap = 0.18
    total = sum(wf for _, _, wf, _ in waves)
    scale = (w - gap * (len(waves) - 1)) / total
    cx = x0
    for i, (hd, lines, wf, hf) in enumerate(waves):
        ww = wf * scale
        hh = max_h * hf
        node(s, cx, base_y - hh, ww, hh, hd, lines, fill=LIGHT if i else ITU_BLUE, ink=INK if i else WHITE,
             head_ink=ITU_BLUE_DARK if i else WHITE, head_size=13, size=11.5)
        if i:
            arrow(s, cx - gap + 0.02, base_y - 0.45, cx - 0.02, base_y - 0.45, color=ITU_BLUE_DARK, weight_pt=2)
        cx += ww + gap
    # platform bar: grows with each wave (tapered look via stepped rectangles)
    py = base_y + 0.15
    cx = x0
    for i, (hd, lines, wf, hf) in enumerate(waves):
        ww = wf * scale
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(py + 0.32 - 0.08 * (i + 1)),
                               Inches(ww + (gap if i < len(waves) - 1 else 0)), Inches(0.08 * (i + 1) + 0.12))
        solid(r, ITU_BLUE_DARK if i == 0 else ITU_BLUE)
        cx += ww + gap
    tb = box(s, x0, py + 0.5, w, 0.35)
    set_text(tb.text_frame, [[(platform_label, 11.5, True, ITU_BLUE_DARK, False)]])
    if closing:
        tb = box(s, 0.72, 6.35, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def play_card(prs, head, paste, returns, you, closing, tag, note,
              labels=('YOU PASTE', 'THE AI RETURNS', 'YOU DECIDE')):
    """The spine of an AI play: input → draft → the human step. Three panels, two arrows; the
    human step is the only dark panel so the eye lands on it. paste/returns/you: [lines]."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    y, h = 1.8, 4.2
    pw, gap = 3.55, 0.62
    xs = [0.72, 0.72 + pw + gap, 0.72 + 2 * (pw + gap)]
    fills = (LIGHT, PANEL_GREY, ITU_BLUE_DARK)
    inks = (INK, INK, WHITE)
    heads = (ITU_BLUE_DARK, ITU_BLUE_DARK, WHITE)
    for x, lab, lines, f, ink, hk in zip(xs, labels, (paste, returns, you), fills, inks, heads):
        node(s, x, y, pw, h, lab, lines, fill=f, ink=ink, head_ink=hk, head_size=12, size=14.5)
    for x in xs[1:]:
        arrow(s, x - gap + 0.08, y + h / 2, x - 0.08, y + h / 2, weight_pt=2)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s
