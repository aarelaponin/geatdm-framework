#!/usr/bin/env python3
"""Helper library for building KP module video decks on the ITU template.

Extracted from the worked, QA'd generator for KP1 Module 1
(KP1-GEA/build_kp1_module1_deck_v01.py). Import these helpers when building a
new module's deck; the worked example remains the reference for composition.

Usage sketch:

    import deck_lib as dl
    prs = dl.open_template()          # the ITU template shipped next to this file
    dl.edit_cover(prs, title='...', kicker='KP1 · Government Enterprise Architecture · Module 2',
                  blurb='...', length='~30 mins across 8 videos', audience='...',
                  panel_heading='THE LIFECYCLE THIS MODULE TEACHES',
                  panel_items=['Discover', 'Assess', 'Adapt', 'Plan', 'Execute & Govern'],
                  panel_footer='4 sign-offs · 6 months to a roadmap · then ongoing')
    dl.edit_agenda(prs, header='Module 2 — eight videos', items=[('2.1  ...', '~4 min'), ...],
                   message_paras=['...', '...'])
    dl.delete_template_slides(prs, keep=2)
    s = dl.add_slide(prs, dl.LAYOUT_BLUE)   # then compose with the helpers below
    prs.save('out.pptx')

All sizes in inches unless noted. Slide canvas: 13.333 x 7.5 in.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- ITU branding constants (from the template + KP video guide §3.i) ----
ITU_BLUE = RGBColor(0x00, 0x9C, 0xD6)        # ITU Blue
ITU_BLUE_DARK = RGBColor(0x00, 0x6E, 0x96)   # darker accent for text on white
INK = RGBColor(0x1A, 0x1A, 0x1A)             # near-black body text
GREY = RGBColor(0x59, 0x59, 0x59)            # captions, kickers, footers
LIGHT = RGBColor(0xE5, 0xF5, 0xFB)           # light-blue tint panels (guide bg colour)
PANEL_GREY = RGBColor(0xF2, 0xF2, 0xF2)      # neutral contrast panel
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIDGREY = RGBColor(0xBF, 0xBF, 0xBF)         # connector / timeline lines
SEPARATOR = RGBColor(0xDD, 0xEC, 0xF4)       # thin row separators

# ---- ITU_ppt_template.pptx layout indices (verified against the template) ----
LAYOUT_WHITE = 11        # 'Blank - Footer (white bg)' — white, www.itu.int + page no. baked in
LAYOUT_BLANK_WHITE = 12  # 'Blank (white bg)'
LAYOUT_BLUE = 13         # 'Blank (blue bg)' — full-bleed F5FAFC rectangle
LAYOUT_THANKS = 14       # '2_Thank you Slide' — "Thank you!" text baked into the layout


TEMPLATE = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ITU_ppt_template.pptx')


def open_template(path=None):
    """Open the ITU template. Defaults to the copy shipped with this skill."""
    return Presentation(path or TEMPLATE)


def add_slide(prs, layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def delete_template_slides(prs, keep=2):
    """Drop every template slide after the first `keep` (cover, agenda).
    python-pptx discards the orphaned parts on save."""
    sldIdLst = prs.slides._sldIdLst
    for sld in list(sldIdLst)[keep:]:
        prs.part.drop_rel(sld.get(qn('r:id')))
        sldIdLst.remove(sld)


# ---------------------------------------------------------------- text/shape primitives
def set_text(tf, runs_spec, align=None, space_after=None):
    """runs_spec: list of paragraphs; each = list of (text, size_pt, bold, color, italic)."""
    tf.word_wrap = True
    first = True
    for para in runs_spec:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if align is not None:
            p.alignment = align
        if space_after is not None:
            p.space_after = space_after
        for (txt, size, bold, color, italic) in para:
            r = p.add_run()
            r.text = txt
            r.font.name = 'Arial'
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color


def retext(tf, paras):
    """Replace an existing text frame's content in place (keeps the box, kills old runs)."""
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    first = True
    for para in paras:
        p = p0 if first else tf.add_paragraph()
        first = False
        for (txt, size, bold, color, italic) in para:
            r = p.add_run()
            r.text = txt
            r.font.name = 'Arial'
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb


def solid(shape, color, line_color=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def panel(slide, x, y, w, h, fill, line=None, radius=0.045):
    p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    p.adjustments[0] = radius
    return solid(p, fill, line)


def panel_text(slide, x, y, w, h, paras, m=0.22):
    tb = box(slide, x + m, y + 0.16, w - 2 * m, h - 0.3)
    set_text(tb.text_frame, paras, space_after=Pt(7))
    return tb


def centered_panel(slide, x, y, w, h, fill, text, size, bold, color):
    p = panel(slide, x, y, w, h, fill)
    tf = p.text_frame
    set_text(tf, [[(text, size, bold, color, False)]], align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return p


def hline(slide, x, y, w, color=SEPARATOR, weight_pt=1):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(weight_pt))
    return solid(ln, color)


def num_chip(slide, x, y, n, d=0.34):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    solid(c, ITU_BLUE)
    tf = c.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    set_text(tf, [[(str(n), 13, True, WHITE, False)]], align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return c


def rows_block(prs, head, rows, closing, tag, note, numbered=True, top=1.6, bottom=6.3,
               head_size=19):
    """`rows_slide` with a headline, a footer tag and a closing line under the rows. The shape
    to reach for wherever a bundle cue says 'N text rows'; the closing line is where the list
    lands."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    rows_slide(s, rows, top=top, bottom=bottom, numbered=numbered, head_size=head_size)
    if closing:
        tb = box(s, 0.72, bottom + 0.12, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


# ---------------------------------------------------------------- slide furniture
def tick(slide):
    """The small ITU-blue tick top-left of content slides (mirrors the template look)."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.42), Inches(0.05), Inches(0.28))
    return solid(r, ITU_BLUE)


def title(slide, text, size=28, width=12.3):
    """Assertion headline. Guide branding: Arial Bold 28pt."""
    tick(slide)
    tb = box(slide, 0.68, 0.30, width, 1.0)
    set_text(tb.text_frame, [[(text, size, True, INK, False)]])
    return tb


def footer(slide, tag, itu=False):
    """Left wayfinding tag ('1.3 · Why projects can't…'). LAYOUT_WHITE already bakes
    www.itu.int + page number — set itu=True ONLY on blue-bg slides (no baked footer)."""
    tb = box(slide, 0.5, 7.08, 8.5, 0.32)
    set_text(tb.text_frame, [[(tag, 9, False, GREY, False)]])
    if itu:
        tb2 = box(slide, 11.0, 7.08, 1.83, 0.32)
        set_text(tb2.text_frame, [[('www.itu.int', 9, False, GREY, False)]], align=PP_ALIGN.RIGHT)


def notes(slide, text):
    """Voice-over + production cues. Notes are the spoken words, never a slide paraphrase."""
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------- composite slides
def rows_slide(slide, rows, top=1.55, bottom=6.85, numbered=True, head_size=19, sub_size=15.5):
    """Stacked rows: bold headline + concrete example line, thin separators, optional chips.
    rows: list of (headline, sub_or_empty). Density is allowed to be lumpy."""
    n = len(rows)
    rh = (bottom - top) / n
    for i, (head, sub) in enumerate(rows):
        y = top + i * rh
        if numbered:
            num_chip(slide, 0.68, y + 0.05, i + 1)
            tx = 1.22
        else:
            tx = 0.72
        tb = box(slide, tx, y, 13.333 - tx - 0.6, rh - 0.05)
        paras = [[(head, head_size, True, INK, False)]]
        if sub:
            paras.append([(sub, sub_size, False, GREY, False)])
        set_text(tb.text_frame, paras, space_after=Pt(3))
        if i < n - 1:
            hline(slide, 0.68, top + (i + 1) * rh - 0.045, 11.9)


def section_slide(prs, kicker, code, name, message, runtime_line, note_text):
    """Blue-bg divider that doubles as the standalone video's opening slide.
    kicker e.g. 'KP1 · MODULE 1 · VIDEO 1.3' — module-scoped, never 'of N'."""
    s = add_slide(prs, LAYOUT_BLUE)
    tb = box(s, 0.9, 0.9, 8, 0.5)
    set_text(tb.text_frame, [[(kicker, 12, True, GREY, False)]])
    tb = box(s, 0.9, 2.0, 2.2, 1.2)
    set_text(tb.text_frame, [[(code, 60, True, ITU_BLUE, False)]])
    tb = box(s, 0.9, 3.15, 10.5, 1.5)
    set_text(tb.text_frame, [[(name, 34, True, INK, False)]])
    tb = box(s, 0.9, 4.75, 10.8, 1.7)
    set_text(tb.text_frame, [[(message, 17, False, GREY, True)]])
    tb = box(s, 0.9, 6.55, 8, 0.4)
    set_text(tb.text_frame, [[(runtime_line, 12, False, GREY, False)]])
    notes(s, note_text)
    return s


def big_slide(prs, text, tag, note_text, sub=None, label='IN ONE SENTENCE'):
    """The quotable, screenshot-ready climax slide that closes a video."""
    s = add_slide(prs, LAYOUT_BLUE)
    tb = box(s, 1.1, 2.3, 11.1, 2.6)
    set_text(tb.text_frame, [[(text, 30, True, INK, False)]])
    if sub:
        tb = box(s, 1.1, 5.1, 11.1, 1.0)
        set_text(tb.text_frame, [[(sub, 16, False, GREY, True)]])
    tb = box(s, 1.1, 1.7, 6, 0.4)
    set_text(tb.text_frame, [[(label, 12, True, ITU_BLUE_DARK, False)]])
    footer(s, tag, itu=True)
    notes(s, note_text)
    return s


def sources_slide(prs, tag, items, note_text=''):
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, 'Sources')
    tb = box(s, 0.72, 1.7, 11.8, 4.2)
    set_text(tb.text_frame, [[('•  ' + it, 17, False, INK, False)] for it in items], space_after=Pt(12))
    tb = box(s, 0.72, 6.35, 11.8, 0.5)
    set_text(tb.text_frame, [[('Find the link in the description.', 14, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note_text or ('Sources slide, held ~5 seconds at the end of the video. Links are compiled '
                           'into the YouTube description per ITU convention; no URLs read aloud.'))
    return s


def block_slide(prs, head, lead, punch, tag, note, punch_fill=LIGHT, punch_ink=ITU_BLUE_DARK,
                lead_size=20, punch_y=4.85):
    """A paragraph of argument, then the line it lands on. The shape to reach for wherever a
    bundle cue says 'single text block'. `lead` is a list of paragraphs.

    `punch_y` raises the landing panel toward a short lead. Leave it at the default and the
    slide is laid out exactly as Modules 1-4 render it; pass a smaller value (see
    `punch_after` in the Module 5 build script) when a two-sentence lead would otherwise
    leave a blank band above the panel."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    # round(): keeps the default height bit-identical to the literal 2.9 Modules 1-4 built with
    tb = box(s, 0.72, 1.75, 11.9, round(punch_y - 1.95, 4))
    set_text(tb.text_frame, [[(p, lead_size, False, INK, False)] for p in lead], space_after=Pt(12))
    panel(s, 0.72, punch_y, 11.9, 1.45, punch_fill)
    panel_text(s, 0.72, punch_y + 0.05, 11.9, 1.45, [[(punch, 23, True, punch_ink, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def two_panel(prs, head, left, right, closing, tag, note, right_fill=PANEL_GREY, height=4.35):
    """Comparison columns: (heading, [lines]) each side, optional closing line beneath.

    `height` shrinks the columns toward short content. Leave it at the default and the slide
    is laid out exactly as Modules 1-4 render it; pass a smaller value (see `panel_height` in
    the Module 5 build script) when three short lines would otherwise sit in a tall box with
    an empty lower half."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    cw, ch, cy = 5.9, height, 1.7
    for i, ((ph, lines), fill) in enumerate(zip((left, right), (LIGHT, right_fill))):
        x = 0.72 + i * (cw + 0.23)
        panel(s, x, cy, cw, ch, fill)
        paras = [[(ph, 19, True, ITU_BLUE_DARK, False)]]
        paras += [[(ln, 16, False, INK, False)] for ln in lines]
        panel_text(s, x, cy, cw, ch, paras)
    if closing:
        tb = box(s, 0.72, round(cy + ch + 0.2, 4), 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def mini_strip(slide, n, active_idx, x=8.05, y=0.42, w=0.92, gap=0.06, h=0.34):
    """Progress strip top-right on per-phase / per-step slides."""
    for j in range(n):
        r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x + j * (w + gap)), Inches(y), Inches(w), Inches(h))
        r.adjustments[0] = 0.35
        solid(r, ITU_BLUE if j == active_idx else LIGHT)


# ---------------------------------------------------------------- template slide 1 & 2 editing
def edit_cover(prs, title_text, kicker, blurb, length, audience,
               panel_heading, panel_items, panel_footer, note_text=''):
    """Rewrite the template's video-title cover (slide 1) in place and replace the grey
    [add image] block with a flat ITU-blue motif panel listing `panel_items`."""
    cover = prs.slides[0]
    for sh in list(cover.shapes):
        if (sh.shape_type == 6 and sh.name == 'Group 19') or sh.name == 'Picture Placeholder 6':
            sh._element.getparent().remove(sh._element)
    for sh in cover.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text
        if 'Title of the Video' in t:
            retext(sh.text_frame, [[(title_text, 30, True, INK, False)]])
        elif 'Subtitle' in t:
            retext(sh.text_frame, [
                [(kicker, 16, True, ITU_BLUE_DARK, False)],
                [(blurb, 14, False, GREY, False)],
            ])
        elif 'Lenght' in t or t.startswith('Length'):
            retext(sh.text_frame, [
                [('Length: ', 14, True, INK, False), (length, 14, False, INK, False)],
                [('Target audience: ', 14, True, INK, False), (audience, 14, False, INK, False)],
            ])
    from pptx.util import Emu
    rp = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.68), 0, Inches(5.653), Inches(7.5))
    solid(rp, ITU_BLUE)
    tbh = box(cover, 8.15, 0.5, 4.7, 0.5)
    set_text(tbh.text_frame, [[(panel_heading, 11, True, LIGHT, False)]])
    py = 1.05
    step = min(1.02, 5.2 / max(1, len(panel_items)))
    for i, ph in enumerate(panel_items):
        last = i == len(panel_items) - 1
        bxs = cover.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(py), Inches(4.2), Inches(0.72))
        bxs.adjustments[0] = 0.14
        solid(bxs, ITU_BLUE_DARK if last else WHITE)
        tf = bxs.text_frame
        tf.margin_left = Inches(0.18)
        set_text(tf, [[(ph, 15, True, WHITE if last else INK, False)]])
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        py += step
    tbf = box(cover, 8.15, 6.35, 4.7, 0.8)
    set_text(tbf.text_frame, [[(panel_footer, 12, False, WHITE, False)]])
    if note_text:
        notes(cover, note_text)
    return cover


def edit_agenda(prs, header, items, message_paras, note_text=''):
    """Rewrite the template agenda (slide 2): left list of videos, right italic core message.
    items: list of (label, runtime). The list placeholder inherits RIGHT alignment from the
    template — this forces LEFT."""
    agenda = prs.slides[1]
    for sh in agenda.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text
        if t.strip() == 'Agenda':
            retext(sh.text_frame, [[(header, 14, True, GREY, False)]])
        elif 'Subtopic 1' in t:
            paras = [[(name, 15, True, INK, False), ('   ' + rt, 12, False, GREY, False)]
                     for name, rt in items]
            retext(sh.text_frame, paras)
            for p in sh.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(9)
        elif 'core message' in t:
            retext(sh.text_frame, [[(m, 19, False, INK, True)] for m in message_paras])
            for p in sh.text_frame.paragraphs:
                p.space_after = Pt(10)
    if note_text:
        notes(agenda, note_text)
    return agenda


# ================================================================ structural diagrams
# Added after the M1–M6 v0.1 review: shape-only visuals (bars, flows, graphs, timelines)
# permitted by the ITU text-only rule. No icons, no imagery — text panels, lines, arrows.
from lxml import etree
from pptx.enum.shapes import MSO_CONNECTOR


def arrow(slide, x1, y1, x2, y2, color=ITU_BLUE_DARK, weight_pt=1.5, head=True):
    """Straight connector with an optional triangle arrowhead at (x2, y2)."""
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(weight_pt)
    if head:
        ln = c.line._get_or_add_ln()
        etree.SubElement(ln, qn('a:tailEnd'), type='triangle', w='med', len='med')
    return c


def label(slide, x, y, w, h, text, size=11, bold=False, color=GREY, align=PP_ALIGN.CENTER,
          fill=None, anchor=MSO_ANCHOR.MIDDLE, italic=False):
    """Small caption box; give it a fill to sit on top of a connector line."""
    if fill is not None:
        p = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        solid(p, fill)
        tf = p.text_frame
    else:
        tf = box(slide, x, y, w, h).text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    set_text(tf, [[(text, size, bold, color, italic)]], align=align)
    tf.vertical_anchor = anchor
    return tf


def node(slide, x, y, w, h, heading, lines=(), fill=LIGHT, ink=INK, head_ink=ITU_BLUE_DARK,
         head_size=13, size=12.5, radius=0.08, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line=None):
    """A labelled box for flow / graph diagrams."""
    p = panel(slide, x, y, w, h, fill, line=line, radius=radius)
    tf = p.text_frame
    tf.margin_left = tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.06)
    paras = [[(heading, head_size, True, head_ink, False)]]
    paras += [[(ln, size, False, ink, False)] for ln in lines]
    set_text(tf, paras, align=align, space_after=Pt(3))
    tf.vertical_anchor = anchor
    return p


def bar_rows(slide, rows, x=0.72, y=1.7, w=11.9, label_w=2.9, row_h=0.62, gap=0.42,
             group_gap=0.55, caption_w=0.0):
    """Proportional bars — length IS the argument. No figures, only ratios.
    rows: list of items; each item is either ('GROUP', text) for a group heading, or
    (row_label, segments, caption) with segments = [(seg_text, frac, fill, ink), ...] where
    frac is a fraction of the full bar width and seg_text may be ''."""
    bar_x = x + label_w
    bar_w = w - label_w - caption_w
    cy = y
    for item in rows:
        if item[0] == 'GROUP':
            tb = box(slide, x, cy, w, 0.36)
            set_text(tb.text_frame, [[(item[1], 12, True, GREY, False)]])
            cy += 0.42
            continue
        rl, segs, caption = item
        tb = box(slide, x, cy, label_w - 0.15, row_h)
        set_text(tb.text_frame, [[(rl, 14, True, INK, False)]])
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        sx = bar_x
        for (st, frac, fill, ink) in segs:
            sw = max(0.02, bar_w * frac)
            r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sx), Inches(cy + 0.06),
                                       Inches(sw - 0.03), Inches(row_h - 0.12))
            solid(r, fill, line=WHITE)
            if st:
                tf = r.text_frame
                tf.margin_left = tf.margin_right = Inches(0.06)
                tf.margin_top = tf.margin_bottom = 0
                set_text(tf, [[(st, 11, True, ink, False)]], align=PP_ALIGN.CENTER)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.word_wrap = True
            sx += sw
        if caption:
            tb = box(slide, sx + 0.1, cy, max(1.0, x + w - sx - 0.1), row_h)
            set_text(tb.text_frame, [[(caption, 12, False, GREY, True)]])
            tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cy += row_h + gap
    return cy


def bars_slide(prs, head, rows, closing, tag, note, **kw):
    """Headline + proportional bars + closing line. See bar_rows for `rows`."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    bar_rows(s, rows, **kw)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def entity_graph(prs, head, closing, tag, note,
                 names=('Capability', 'Service', 'Application', 'Data Domain',
                        'Technology Component', 'Organisation'),
                 defs=('what a body can do', 'how it reaches a citizen or body',
                       'the software that supports it', 'a kind of information, one owner',
                       'the infrastructure underneath', 'owns each of them'),
                 rels=('delivered by', 'supported by', 'uses', 'runs on', 'owns')):
    """The PAERA metamodel as a graph: four entity nodes with labelled arrows, Technology as
    the base everything runs on, Organisation as the owner bracketing them all."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    cap, svc, app, dom, tech, org = names
    nw, nh = 2.9, 1.05
    # column for Organisation on the left, then the 2x2 of entities, Technology underneath
    ox, oy, ow, oh = 0.72, 1.75, 2.0, 4.2
    node(s, ox, oy, ow, oh, org, [defs[5]], fill=PANEL_GREY, head_size=15, size=12.5)
    gx = 3.55
    x1, x2 = gx, gx + 5.4
    y1, y2 = 1.75, 3.55
    node(s, x1, y1, nw, nh, cap, [defs[0]], head_size=15)
    node(s, x2, y1, nw, nh, svc, [defs[1]], head_size=15)
    node(s, x1, y2, nw, nh, app, [defs[2]], head_size=15)
    node(s, x2, y2, nw, nh, dom, [defs[3]], head_size=15)
    # technology base under the 2x2
    ty = 5.05
    node(s, x1, ty, x2 + nw - x1, 0.9, tech, [defs[4]], fill=ITU_BLUE_DARK, ink=WHITE,
         head_ink=WHITE, head_size=15)
    # arrows + relationship labels
    arrow(s, x1 + nw, y1 + nh / 2, x2, y1 + nh / 2)
    label(s, x1 + nw + 0.55, y1 + nh / 2 - 0.2, 1.4, 0.4, rels[0], size=12, bold=True,
          color=ITU_BLUE_DARK, fill=WHITE)
    arrow(s, x1 + nw / 2, y1 + nh, x1 + nw / 2, y2)
    label(s, x1 + nw / 2 + 0.12, y1 + nh + 0.17, 1.5, 0.4, rels[1], size=12, bold=True,
          color=ITU_BLUE_DARK, align=PP_ALIGN.LEFT, fill=WHITE)
    arrow(s, x1 + nw, y2 + nh / 2, x2, y2 + nh / 2)
    label(s, x1 + nw + 0.9, y2 + nh / 2 - 0.2, 0.9, 0.4, rels[2], size=12, bold=True,
          color=ITU_BLUE_DARK, fill=WHITE)
    # everything runs on technology: two short arrows down into the base
    for xx in (x1 + nw / 2, x2 + nw / 2):
        arrow(s, xx, y2 + nh, xx, ty)
    label(s, x1 + nw + 0.3, y2 + nh + 0.05, 2.1, 0.36, 'everything ' + rels[3], size=12, bold=True,
          color=ITU_BLUE_DARK, fill=WHITE)
    # organisation owns each: one arrow into the grid with a label
    arrow(s, ox + ow, oy + oh / 2, x1, oy + oh / 2)
    label(s, ox + ow + 0.02, oy + oh / 2 - 0.42, 0.8, 0.36, rels[4], size=12, bold=True,
          color=ITU_BLUE_DARK, fill=WHITE)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def prompt_anatomy(prs, head, parts, closing, tag, note, panel_heading='THE PROMPT, AS PASTED'):
    """One prompt shown as one block, its four parts bracketed and labelled.
    parts: list of (label, example_line)."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    n = len(parts)
    top, bottom = 1.65, 6.15
    px, pw = 6.55, 6.07
    panel(s, px, top, pw, bottom - top, PANEL_GREY, radius=0.03)
    tb = box(s, px + 0.25, top + 0.12, pw - 0.5, 0.35)
    set_text(tb.text_frame, [[(panel_heading, 10.5, True, GREY, False)]])
    rows_top = top + 0.55
    rh = (bottom - rows_top) / n
    for i, (lab, example) in enumerate(parts):
        y = rows_top + i * rh
        num_chip(s, 0.72, y + 0.08, i + 1)
        tb = box(s, 1.22, y, 4.9, rh - 0.1)
        set_text(tb.text_frame, [[(lab, 16, True, INK, False)]])
        # bracket: a short blue bar on the panel's left edge for this part
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(px + 0.12), Inches(y + 0.1), Inches(0.06), Inches(rh - 0.25))
        solid(r, ITU_BLUE)
        hline(s, 6.2, y + 0.24, 0.3, color=ITU_BLUE, weight_pt=1.5)
        tb = box(s, px + 0.32, y + 0.02, pw - 0.55, rh - 0.1)
        set_text(tb.text_frame, [[(example, 15, False, INK, True)]])
        if i < n - 1:
            hline(s, 0.72, y + rh - 0.05, 5.4)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def gate_flow(prs, head, questions, closing, tag, note, project=('A project proposal', 'before funding'),
              pass_box=('PASS', 'Funded — consume the shared blocks'),
              exc_box=('EXCEPTION', 'Granted in writing, with an expiry date and a reason'),
              log_box=('DECISION LOG', 'Every ruling and every exception, with its reason — into the repository')):
    """The architecture review gate as one mechanism: proposal → the questions → the Board's
    ruling → pass or time-boxed exception → the decision log."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    y0, h = 1.75, 4.3
    mid = y0 + h / 2
    # project box
    node(s, 0.72, mid - 0.7, 1.75, 1.4, project[0], [project[1]], fill=PANEL_GREY, head_size=13,
         anchor=MSO_ANCHOR.MIDDLE)
    # gate panel with numbered questions
    gx, gw = 2.95, 4.55
    p = panel(s, gx, y0, gw, h, LIGHT, radius=0.05)
    tb = box(s, gx + 0.2, y0 + 0.1, gw - 0.4, 0.35)
    set_text(tb.text_frame, [[('THE GATE — THE SAME QUESTIONS, EVERY TIME', 10.5, True, ITU_BLUE_DARK, False)]])
    qy = y0 + 0.55
    qh = (h - 0.65) / len(questions)
    for i, q in enumerate(questions):
        num_chip(s, gx + 0.2, qy + i * qh + 0.04, i + 1, d=0.3)
        tb = box(s, gx + 0.6, qy + i * qh - 0.02, gw - 0.75, qh)
        set_text(tb.text_frame, [[(q, 12.5, False, INK, False)]])
    arrow(s, 2.47, mid, gx, mid)
    # ruling diamond
    dx, dw = 7.9, 1.35
    d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(dx), Inches(mid - 0.75), Inches(dw), Inches(1.5))
    solid(d, ITU_BLUE_DARK)
    tf = d.text_frame
    tf.margin_left = tf.margin_right = 0
    set_text(tf, [[('The Board', 11, True, WHITE, False)], [('rules', 11, True, WHITE, False)]], align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    arrow(s, gx + gw, mid, dx, mid)
    # outcomes
    ox, ow, oh = 9.6, 1.55, 1.45
    node(s, ox, y0 + 0.15, ow, oh, pass_box[0], [pass_box[1]], fill=LIGHT, head_size=12, size=11)
    node(s, ox, y0 + h - oh - 0.15, ow, oh, exc_box[0], [exc_box[1]], fill=PANEL_GREY, head_size=12, size=11)
    arrow(s, dx + dw, mid, ox, y0 + 0.15 + oh / 2)
    arrow(s, dx + dw, mid, ox, y0 + h - 0.15 - oh / 2)
    # decision log
    lx, lw = 11.45, 1.17
    node(s, lx, mid - 1.1, lw, 2.2, log_box[0], [log_box[1]], fill=ITU_BLUE_DARK, ink=WHITE, head_ink=WHITE,
         head_size=11, size=10.5)
    arrow(s, ox + ow, y0 + 0.15 + oh / 2, lx, mid - 0.4)
    arrow(s, ox + ow, y0 + h - 0.15 - oh / 2, lx, mid + 0.4)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def stack_slide(prs, head, top_block, base_block, side_note, closing, tag, note):
    """One block standing on a wider base — 'what Wave 1 really builds'.
    top_block / base_block: (heading, [lines])."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    node(s, 2.9, 1.8, 5.2, 1.75, top_block[0], top_block[1], fill=LIGHT, head_size=15, size=13)
    node(s, 0.72, 3.85, 9.55, 2.1, base_block[0], base_block[1], fill=ITU_BLUE_DARK, ink=WHITE,
         head_ink=WHITE, head_size=15, size=13)
    # bracket note on the right
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.55), Inches(1.8), Inches(0.05), Inches(4.15))
    solid(r, ITU_BLUE)
    tb = box(s, 10.75, 1.8, 1.9, 4.15)
    set_text(tb.text_frame, [[(ln, 12.5, False, GREY, True)] for ln in side_note], space_after=Pt(8))
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def wave_timeline(prs, head, waves, board_label, platform_label, closing, tag, note):
    """Rollout as a wave roadmap: each wave a block whose width and height shrink as the shared
    platforms carry more of the work; a Board bar governs the whole pipeline; a platform bar
    grows underneath. waves: list of (heading, [lines], width_frac, height_frac)."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    x0, w = 0.72, 11.9
    # governance bar across the top
    b = panel(s, x0, 1.7, w, 0.5, ITU_BLUE_DARK, radius=0.3)
    tf = b.text_frame
    set_text(tf, [[(board_label, 12, True, WHITE, False)]], align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # waves
    base_y, max_h = 5.35, 2.85
    gap = 0.18
    total = sum(wf for _, _, wf, _ in waves)
    scale = (w - gap * (len(waves) - 1)) / total
    cx = x0
    for i, (hd, lines, wf, hf) in enumerate(waves):
        ww = wf * scale
        hh = max_h * hf
        node(s, cx, base_y - hh, ww, hh, hd, lines, fill=LIGHT if i else ITU_BLUE, ink=INK if i else WHITE,
             head_ink=ITU_BLUE_DARK if i else WHITE, head_size=13, size=11.5)
        if i:
            arrow(s, cx - gap, base_y - 0.3, cx, base_y - 0.3, color=MIDGREY, weight_pt=1.2)
        cx += ww + gap
    # platform bar: grows with each wave (tapered look via stepped rectangles)
    py = base_y + 0.15
    cx = x0
    for i, (hd, lines, wf, hf) in enumerate(waves):
        ww = wf * scale
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(py + 0.32 - 0.08 * (i + 1)),
                               Inches(ww + (gap if i < len(waves) - 1 else 0)), Inches(0.08 * (i + 1) + 0.12))
        solid(r, ITU_BLUE_DARK if i == 0 else ITU_BLUE)
        cx += ww + gap
    tb = box(s, x0, py + 0.5, w, 0.35)
    set_text(tb.text_frame, [[(platform_label, 11.5, True, ITU_BLUE_DARK, False)]])
    if closing:
        tb = box(s, 0.72, 6.35, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def play_card(prs, head, paste, returns, you, closing, tag, note,
              labels=('YOU PASTE', 'THE AI RETURNS', 'YOU DECIDE')):
    """The spine of an AI play: input → draft → the human step. Three panels, two arrows; the
    human step is the only dark panel so the eye lands on it. paste/returns/you: [lines]."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    y, h = 1.8, 4.2
    pw, gap = 3.55, 0.62
    xs = [0.72, 0.72 + pw + gap, 0.72 + 2 * (pw + gap)]
    fills = (LIGHT, PANEL_GREY, ITU_BLUE_DARK)
    inks = (INK, INK, WHITE)
    heads = (ITU_BLUE_DARK, ITU_BLUE_DARK, WHITE)
    for x, lab, lines, f, ink, hk in zip(xs, labels, (paste, returns, you), fills, inks, heads):
        node(s, x, y, pw, h, lab, lines, fill=f, ink=ink, head_ink=hk, head_size=12, size=14.5)
    for x in xs[1:]:
        arrow(s, x - gap + 0.08, y + h / 2, x - 0.08, y + h / 2, weight_pt=2)
    if closing:
        tb = box(s, 0.72, 6.3, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s
