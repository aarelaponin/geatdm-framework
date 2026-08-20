#!/usr/bin/env python3
"""Split a combined KP module deck into per-video decks, each opening with a
retitled copy of the module cover as its title card. Content slides untouched.

Usage:
    python split_module_deck.py <combined.pptx> <spec.json> [outdir]

spec.json:
{
  "kicker_prefix": "KP1 · Government Enterprise Architecture · Module 1",
  "audience": "CDO · Director-General · sector minister",
  "out_pattern": "KP1_Module1_Video_{code}_v0.1.pptx",
  "cover_title_marker": "Why a PAERA-anchored",   // text that identifies the cover title box
  "cover_subtitle_marker": "Eight standalone videos",
  "videos": [
    {"code": "1.1", "title": "Why your country needs a national EA", "mins": "~4",
     "range": [3, 7],   // 1-indexed inclusive slide range in the combined deck
     "message": "Without a shared plan ..."},
    ...
  ]
}

Numbering rule: the title-card kicker is module-scoped — 'Module N · Video N.x',
never 'Video N.x of <count>' (reads as a claim about the whole KP).
"""
import json
import sys
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x59, 0x59, 0x59)
BLUE_DARK = RGBColor(0x00, 0x6E, 0x96)


def retext(tf, paras):
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    first = True
    for para in paras:
        p = p0 if first else tf.add_paragraph()
        first = False
        for (txt, size, bold, color, italic) in para:
            r = p.add_run()
            r.text = txt
            r.font.name = 'Arial'
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color


def main():
    if len(sys.argv) < 3:
        sys.exit(__doc__)
    src, specpath = sys.argv[1], sys.argv[2]
    outdir = sys.argv[3] if len(sys.argv) > 3 else '.'
    os.makedirs(outdir, exist_ok=True)
    spec = json.load(open(specpath))

    for v in spec['videos']:
        code, name, mins = v['code'], v['title'], v['mins']
        a, b = v['range']
        prs = Presentation(src)
        sldIdLst = prs.slides._sldIdLst
        keep = {1} | set(range(a, b + 1))
        for i, sld in enumerate(list(sldIdLst), start=1):
            if i not in keep:
                prs.part.drop_rel(sld.get(qn('r:id')))
                sldIdLst.remove(sld)

        cover = prs.slides[0]
        tsize = 28 if len(name) < 40 else 24
        for sh in cover.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text
            if spec['cover_title_marker'] in t:
                retext(sh.text_frame, [[('%s — %s' % (code, name), tsize, True, INK, False)]])
            elif spec['cover_subtitle_marker'] in t:
                retext(sh.text_frame, [
                    [('%s · Video %s' % (spec['kicker_prefix'], code), 13, True, BLUE_DARK, False)],
                    [(v['message'], 12, False, GREY, True)],
                ])
            elif t.startswith('Length'):
                retext(sh.text_frame, [
                    [('Length: ', 14, True, INK, False), ('%s mins' % mins, 14, False, INK, False)],
                    [('Target audience: ', 14, True, INK, False), (spec['audience'], 14, False, INK, False)],
                ])
        cover.notes_slide.notes_text_frame.text = (
            'Title card for standalone video %s. Content slides are identical to the corresponding '
            'section of the combined deck %s; voice-over is in each slide’s notes.'
            % (code, os.path.basename(src)))

        out = os.path.join(outdir, spec['out_pattern'].format(code=code))
        prs.save(out)
        print(out, len(list(prs.slides._sldIdLst)), 'slides')


if __name__ == '__main__':
    main()
