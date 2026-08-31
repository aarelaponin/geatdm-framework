#!/usr/bin/env python3
"""Generate the scripts-only companion .md from a built module deck + its split spec.

The deck's speaker notes already are the voice-over, so the companion is derived from
them rather than written a second time — that is what keeps the two in sync. Production
cues live in the same notes and are dropped here: a note paragraph is narration only if
it carries the 'VO:' / 'VO, slide N:' marker.

Usage:
    python scripts_from_deck.py <combined.pptx> <spec.json> <outdir> \\
        --kp KP1 --module 2 --prefix KP1_M2 --version v0.1

Writes <prefix>_Scripts_<version>.md (all videos) and one <prefix>_<code>_Scripts_<version>.md
per video, matching the naming used for Module 1.
"""
import argparse
import json
import os
import re

from pptx import Presentation
from pptx.util import Inches

VO_RE = re.compile(r'^VO(?:,\s*slide\s*\d+)?\s*:\s*', re.I)
# title() places the headline box here; sources_slide and every content slide use it.
TITLE_LEFT, TITLE_TOP = Inches(0.68), Inches(0.30)


def slide_title(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.left == TITLE_LEFT and sh.top == TITLE_TOP:
            return sh.text_frame.text.strip()
    return ''


def is_climax(slide):
    return any(sh.has_text_frame and sh.text_frame.text.strip() == 'IN ONE SENTENCE'
               for sh in slide.shapes)


def narration(slide):
    """The VO paragraphs of a note, marker stripped. Returns (text, dropped_cues)."""
    if not slide.has_notes_slide:
        return '', []
    paras = [p.strip() for p in slide.notes_slide.notes_text_frame.text.split('\n\n') if p.strip()]
    keep = [VO_RE.sub('', p) for p in paras if VO_RE.match(p)]
    return '\n\n'.join(keep), [p for p in paras if not VO_RE.match(p)]


def render(videos, header, intro):
    out = ['# ' + header, '', intro, '']
    for v in videos:
        out += ['---', '', '## %s %s (%s min)' % (v['code'], v['title'], v['mins']), '',
                '> *%s*' % v['message'], '']
        for heading, text in v['slides']:
            out += ['### Slide — ' + heading, '', text, '']
    return '\n'.join(out).rstrip() + '\n'


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('deck')
    ap.add_argument('spec')
    ap.add_argument('outdir')
    ap.add_argument('--kp', required=True)
    ap.add_argument('--module', required=True)
    ap.add_argument('--prefix', required=True)
    ap.add_argument('--version', default='v0.1')
    a = ap.parse_args()

    spec = json.load(open(a.spec))
    prs = Presentation(a.deck)
    slides = list(prs.slides)
    dropped = []

    videos = []
    for v in spec['videos']:
        lo, hi = v['range']
        items = []
        for n in range(lo, hi + 1):
            sl = slides[n - 1]
            head = slide_title(sl)
            if n == lo:
                head = 'Title (%s)' % v['code']
            elif is_climax(sl):
                head = 'In one sentence'
            if head == 'Sources':
                items.append(('Sources', '*(No narration.)*'))
                continue
            text, cues = narration(sl)
            assert text, 'slide %d (%s) has no VO paragraph in its notes' % (n, head)
            dropped += [(n, c) for c in cues]
            items.append((head, text))
        videos.append(dict(v, slides=items))

    codes = '%s – %s' % (videos[0]['code'], videos[-1]['code'])
    intro = ('Spoken narration only, one section per video (%s), slide-by-slide, matching '
             '`%s`. Each video is standalone. Sources slides carry no narration — hold ~5 '
             'seconds; links go in the video description.'
             % (codes, os.path.basename(a.deck)))
    base = '%s Module %s (Topic %s) — Voice-over script' % (a.kp, a.module, a.module)

    os.makedirs(a.outdir, exist_ok=True)
    combined = os.path.join(a.outdir, '%s_Scripts_%s.md' % (a.prefix, a.version))
    open(combined, 'w').write(render(videos, base + 's', intro))
    print(combined)
    for v in videos:
        p = os.path.join(a.outdir, '%s_%s_Scripts_%s.md' % (a.prefix, v['code'], a.version))
        head = '%s — %s %s (%s min)' % (base, v['code'], v['title'], v['mins'])
        open(p, 'w').write(render([v], head, intro))
        print(p)

    print('\nProduction cues dropped (verify none of these is narration):')
    for n, c in dropped:
        print('  slide %-3d %s' % (n, c[:110].replace('\n', ' ')))


if __name__ == '__main__':
    main()
