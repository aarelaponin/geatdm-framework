#!/usr/bin/env python3
# Build the KP1 Module 4 (Topic 4) video deck on the ITU template.
# Content only — every generic helper, branding constant and layout index comes from
# ITU-Giga-KP-Plugin/skills/kp-deck-builder/scripts/deck_lib.py (which also ships the
# template). Conventions and design rules: that skill's SKILL.md.
# Generated .pptx is NEVER hand-edited — fix here, re-render, re-run the split
# (kp-deck-builder/scripts/split_module_deck.py + the split spec next to the decks).
# Override paths with TEMPLATE= and OUT_PATH= env vars.
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                '..', 'ITU-Giga-KP-Plugin', 'skills', 'kp-deck-builder', 'scripts'))
from deck_lib import (
    GREY, INK, ITU_BLUE, ITU_BLUE_DARK, LIGHT, MIDGREY, WHITE,
    LAYOUT_THANKS, LAYOUT_WHITE,
    add_slide, big_slide, block_slide, box, delete_template_slides, edit_agenda,
    edit_cover, footer, notes, open_template, panel, rows_block,
    section_slide, set_text, sources_slide, title, two_panel)

from pptx.enum.text import PP_ALIGN

prs = open_template(os.environ.get('TEMPLATE'))

AUDIENCE = 'Chief architect · senior architect · sector ICT lead'


def section(code, name, message, runtime, note):
    return section_slide(prs, 'KP1 · MODULE 4 · VIDEO %s' % code, code, name, message,
                         runtime + ' · standalone video · voice-over on text slides', note)


# ---------------------------------------------------------------- COVER (edit slide 1)
edit_cover(
    prs,
    title_text='Running the whole method\non one real sector',
    kicker='KP1 · Government Enterprise Architecture · Module 4',
    blurb='Eight standalone videos that run the five-phase lifecycle end to end on Progressa, a '
          'demonstration education sector: discover what exists, rank the gaps, fit the framework, '
          'design the target architecture, sequence the costed roadmap, and stand up the '
          'governance that keeps re-use happening — then the recipe to run it on your own sector.',
    length='~39 mins across 8 videos (4.1 – 4.8)',
    audience=AUDIENCE,
    panel_heading='THE FIVE PHASES, RUN ON ONE SECTOR',
    panel_items=['Discover — an honest picture',
                 'Assess — a ranked gap analysis',
                 'Adapt — framework and sourcing',
                 'Plan — the target, then the roadmap',
                 'Execute & Govern — the living EA'],
    panel_footer='5 phases · 4 sign-offs · 6 deliverables · 1 worked sector',
    note_text='Cover for the combined Module 4 deck. Each section that follows is one standalone '
              '~4–5 minute video. Modules 1 to 3 taught the pieces — why an architecture matters, '
              'how to read and model a government, how to house and govern the practice. Module 4 '
              'puts them in motion: the whole lifecycle run on Progressa, a fictional '
              'demonstration education sector, so the architect watches each phase produce its '
              'deliverable and pass its sign-off.')

# ---------------------------------------------------------------- AGENDA (edit slide 2)
edit_agenda(
    prs,
    header='Module 4 — eight videos',
    items=[
        ('4.1  Meet Progressa and its fragmentation problem', '~5 min'),
        ('4.2  Phase 1, Discover — map what the sector has today', '~5 min'),
        ('4.3  Phase 2, Assess — find the gaps and rank them', '~5 min'),
        ('4.4  Phase 3, Adapt — fit PAERA, decide build, buy or share', '~5 min'),
        ('4.5  Phase 4, Plan — design the target architecture', '~5 min'),
        ('4.6  Phase 4, Plan — sequence the roadmap and cost it', '~5 min'),
        ('4.7  Phase 5, Execute & Govern — stand up the living EA', '~5 min'),
        ('4.8  Run this on your own sector — the transferable recipe', '~4 min'),
    ],
    message_paras=[
        'A method you have only been told about is a method you cannot run. So this module runs '
        'the whole lifecycle, phase by phase, on one sector you can see all of.',
        'Progressa is fictional on purpose — it lets every step be shown in detail. The '
        'institutions change when you take this home. The method does not.',
    ],
    note_text='Navigation slide for the combined deck; the videos ship standalone on YouTube. '
              '4.1 sets the canvas; 4.2–4.7 run the five phases in order, each producing its '
              'deliverable and passing its sign-off; 4.8 strips the method to what transfers. '
              'Note that the Plan phase spans two videos — design the target (4.5), then sequence '
              'the roadmap (4.6) — so the count is five phases, four sign-offs, six deliverables.')

delete_template_slides(prs, keep=2)


# ================================================================ 4.1
T = '4.1 · Meet Progressa — a real sector with a real fragmentation problem'
section('4.1', 'Meet Progressa — a real sector with a real fragmentation problem',
        'Progressa is a demonstration education sector with the fragmentation problem most '
        'countries share — the same learner registered in several places, identity proven on '
        'paper. Meet its institutions and its problem, because the rest of this module runs the '
        'full method on exactly this canvas.',
        '~5 minutes',
        "VO, slide 1: To see the method work, you need a realistic place to run it. Progressa is a "
        "demonstration country with an education sector like many across the continent — real "
        "institutions, a real fragmentation problem, and a minister who wants results. Everything "
        "in this module runs the five-phase lifecycle on Progressa, so you can watch each phase "
        "produce its deliverable, then run the same steps on your own sector.\n\n"
        "On screen: Progressa is a fictional demonstration country. All institutions are "
        "fictional and appear as plain-text labels — no emblems, no logos.")

rows_block(prs, 'Five bodies: a policy unit, a service authority, two registries, a platform',
           [('Ministry of Education, Youth and Sport (MoEYS)',
             'Sets policy and funds schools. The policy unit.'),
            ('National Examination Authority (PNEA)',
             'Runs examinations and certifies results. The service authority.'),
            ('Learner Registry (PLR)',
             'Meant to be the one list of who is a learner.'),
            ('National Identity Authority (PNIA)',
             'Owns the person identity every sector reuses.'),
            ('Digital Government Authority (PDGA)',
             'Runs the shared data-exchange backbone and payments.')],
           'Two registries — and neither is the undisputed owner of the learner.',
           T,
           "VO: Progressa's education sector has five bodies that matter. The Ministry of "
           "Education, Youth and Sport sets policy and funds schools. The National Examination "
           "Authority runs examinations and certifies results. The Learner Registry is meant to be "
           "the one list of who is a learner. The National Identity Authority owns the person "
           "identity every sector reuses. And the Digital Government Authority runs the shared "
           "data-exchange backbone and payments. Five bodies — a policy unit, a service authority, "
           "two registries, and a shared-platform provider.\n\n"
           "Production cue: this is the cast for the whole module. Hold it a beat longer than a "
           "normal list slide — every later video refers back to these five.")

block_slide(prs, 'One learner, three registrations, and the three lists do not agree',
            ['A learner is registered three times — once in the school census, once by the '
             'Examination Authority for exams, once by a social grant programme. None of the three '
             'lists agree.',
             'A parent proves the child\'s identity on paper at every counter, because no system '
             'trusts another\'s. The minister has promised a single learner record that follows '
             'the child from primary school to university — and it cannot be delivered, because '
             'the systems do not fit together.'],
            'Not a Progressa problem. Most education sectors share it.',
            T,
            "VO: Here is the problem the minister feels. A learner is registered three times — "
            "once in the school census, once by the Examination Authority for exams, once by a "
            "social grant programme. None of the three lists agree. A parent proves the child's "
            "identity on paper at every counter, because no system trusts another's. The minister "
            "has promised a single learner record that follows the child from primary school to "
            "university — and it cannot be delivered, because the systems do not fit together. "
            "This is not a Progressa problem. It is the problem most education sectors share.")

block_slide(prs, 'Fictional on purpose — so every step can be shown in full detail',
            ['A fictional canvas lets us show every step in detail without exposing any real '
             'country, and the shapes Progressa has — duplicate registries, paper re-entry, a '
             'stalled flagship — are the shapes you will recognise in your own sector.',
             'So do one thing as the lifecycle runs: keep your own sector in mind, and at each '
             'step ask what the equivalent would be for you.'],
            'The institutions change. The method does not.',
            T,
            "VO: Progressa is fictional, on purpose. It lets us show every step in detail without "
            "exposing any real country, and the shapes it has — duplicate registries, paper "
            "re-entry, a stalled flagship — are the shapes you will recognise in your own sector. "
            "As we run the lifecycle, do one thing: keep your own sector in mind, and at each step "
            "ask what the equivalent would be for you. The institutions change. The method does "
            "not.")

rows_block(prs, 'Five phases, six deliverables — this is what you will watch being built',
           [('A picture of what exists today', 'Discover — the Discovery brief.'),
            ('A ranked gap analysis', 'Assess — the right problems in the right order.'),
            ('A localised framework and a sourcing matrix',
             'Adapt — build, buy, share or sandbox, block by block.'),
            ('A target architecture', 'Plan — the designed future state.'),
            ('A sequenced, costed roadmap', 'Plan — the document the minister takes to cabinet.'),
            ('A living, governed architecture',
             'Execute & Govern — the practice that keeps re-use happening.')],
           'Six deliverables across five phases — Plan produces two of them.',
           T,
           "VO: Running the method on Progressa produces six things across the five phases. A "
           "picture of what exists today. An honest assessment of the gaps, ranked. A framework "
           "fitted to Progressa, with each building block marked build, buy or share. A target "
           "architecture — the designed future state the gaps are closed toward. A sequenced "
           "roadmap the minister can take to cabinet. And a living, governed architecture that "
           "keeps re-use happening after the consultants leave. Each one is built on a real "
           "sector, in detail, where you can see exactly how it is done.",
           top=1.55, bottom=6.35, head_size=17)

big_slide(prs,
          'Progressa is a realistic education sector with the fragmentation problem most '
          'countries share — the canvas for running the whole method, step by step, so you can '
          'reproduce it on your own.',
          T,
          "VO: So this is Progressa. A realistic education sector, five institutions, and the "
          "fragmentation problem most countries live with. It is the canvas. Now the method runs "
          "on it, one phase at a time, and you watch the architecture take shape — so you can do "
          "the same where you work.")

sources_slide(prs, T, [
    'PAERA v1.0 — §2.1 (Problem statement)',
    'PAERA v1.0 — §2.3 (Role of Enterprise Architecture)',
])


# ================================================================ 4.2
T = '4.2 · Phase 1, Discover — map what the sector has today'
section('4.2', 'Phase 1, Discover — map what the sector has today',
        'Discovery produces one deliverable — an honest picture of what the sector has today, with '
        'no recommendations yet — signed off as accurate before any analysis begins. Watch it done '
        'on Progressa.',
        '~5 minutes',
        "VO, slide 1: The first phase is Discover. One question: what exists today? Not what is "
        "wrong with it — that comes later. Just an accurate picture of where the sector is now. On "
        "Progressa, Discovery takes about three to four weeks and produces a single deliverable: "
        "the Discovery brief.")

rows_block(prs, 'Discovery collects five things — and judges none of them',
           [('Strategies in force',
             'The national digital strategy; the education sector plan.'),
            ('Systems that exist',
             'School census, exam management, the social grant system.'),
            ('Registries and who owns them',
             'The learner lists; the identity register.'),
            ('Stakeholders',
             'The ministry, the authorities, the donors funding each piece.'),
            ('The legal framework',
             'The data-protection act; the mandate of each body.')],
           'You gather. You do not judge.',
           T,
           "VO: Discovery collects five things on Progressa. The strategies in force — the "
           "national digital strategy, the education sector plan. The systems that exist — the "
           "school census system, the exam-management system, the social grant system. The "
           "registries and who owns them — the learner lists, the identity register. The "
           "stakeholders — the ministry, the authorities, the donors funding each piece. And the "
           "legal framework — the data-protection act, the mandates of each body. You gather, you "
           "do not judge.")

block_slide(prs, 'Describe without recommending — the judgement belongs to the next phase',
            ['When your architects find that the Examination Authority keeps its own learner list, '
             'the urge is to write "this should be fixed". Resist it.',
             'Discovery records that the list exists, who owns it, how many learners it holds. '
             'Whether it is a problem, how severe, what priority — that belongs to Assess. Mixing '
             'them biases the assessment before you have the full picture.'],
            'Discover first. Judge second.',
            T,
            "VO: The hardest discipline in Discovery is to describe without recommending. When "
            "your architects find that the Examination Authority keeps its own learner list, the "
            "urge is to write 'this should be fixed'. Resist it. Discovery records that the list "
            "exists, who owns it, how many learners it holds. The judgement — is this a problem, "
            "how severe, what priority — belongs to the next phase. Mixing them biases the "
            "assessment before you have the full picture. Discover first. Judge second.")

block_slide(prs, 'The brief records three learner lists as a fact, not yet as a problem',
            ['Progressa\'s Discovery brief lays it out plainly. Five institutions, classified. The '
             'systems each runs. The registries — and here the brief simply records that three '
             'separate learner lists exist, in the census, the Examination Authority, and the '
             'grant programme.',
             'The data flows — mostly paper and manual re-entry between systems. The legal '
             'constraints — what each body is mandated to hold.'],
            'No recommendations. Just the picture, accurate enough to build on.',
            T,
            "VO: Progressa's Discovery brief lays it out plainly. Five institutions, classified. "
            "The systems each runs. The registries — and here the brief simply records the fact "
            "that three separate learner lists exist, in the census, the Examination Authority, "
            "and the grant programme. The data flows — mostly paper and manual re-entry between "
            "systems. The legal constraints — what each body is mandated to hold. No "
            "recommendations. Just the picture, accurate enough to build on.")

block_slide(prs, 'Sign-off one: accurate. Not complete, not perfect — accurate',
            ['Discovery ends with the first of the lifecycle\'s four sign-offs. On Progressa, the '
             'digitalisation officer chairing the EA Board confirms one thing: the picture is '
             'accurate enough to build on.',
             'This gate is cheap and fast, and skipping it is expensive.'],
            'An assessment built on a wrong picture produces wrong priorities.',
            T,
            "VO: Discovery ends with the first of the lifecycle's four sign-offs. The senior "
            "decision-maker — on Progressa, the digitalisation officer chairing the EA Board — "
            "confirms one thing: the picture is accurate enough to build on. Not complete, not "
            "perfect — accurate. This gate is cheap and fast, and skipping it is expensive: an "
            "assessment built on a wrong picture produces wrong priorities. Get the picture signed "
            "off as true, and the next phase stands on solid ground.")

big_slide(prs,
          'Discover produces one honest picture of what the sector has today — described, not '
          'judged — signed off as accurate before any analysis begins.',
          T,
          "VO: That is Phase 1 on Progressa. One question — what exists today. One deliverable — "
          "the Discovery brief. One discipline — describe, do not recommend. One sign-off — the "
          "picture is accurate. Three to four weeks, and the architecture work has a true "
          "foundation.")

sources_slide(prs, T, [
    'PAERA v1.0 — §3.1.3 (Readiness Assessment)',
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
])


# ================================================================ 4.3
T = '4.3 · Phase 2, Assess — find the gaps and rank them'
section('4.3', 'Phase 2, Assess — find the gaps and rank them',
        'Assess turns Progressa\'s picture into a ranked gap analysis — the right problems, in the '
        'right order — signed off as ground truth. It is the assessment method applied to a real '
        'sector.',
        '~5 minutes',
        "VO, slide 1: Phase two is Assess. Now you judge. The question: what is the gap between "
        "where Progressa is and where it needs to be? The deliverable is the current state in four "
        "layers, maturity scorecards, and a gap analysis that ranks the problems. About six to "
        "eight weeks. This is where the picture from Discovery becomes a decision.\n\n"
        "Retrieval prompt — ask before playing on: in your own sector, name the gaps you would "
        "expect a first assessment to find. Hold your list. The four Progressa surfaces are two "
        "slides from now, and almost every first assessment finds the same four.")

rows_block(prs, "Score the capabilities and “things are bad” becomes a map",
           [('Register a learner — LOW', 'Three lists, none of them authoritative.'),
            ('Prove identity — MEDIUM', 'PNIA exists; education does not consume it.'),
            ('Certify a result — HIGH', 'The Examination Authority does this well.'),
            ('Share data across bodies — VERY LOW', 'Almost everything still moves on paper.')],
           'Not everything is broken. Knowing exactly which part is, is the point.',
           T,
           "VO: Start by scoring Progressa's capabilities against PAERA-anchored standards. "
           "Register a learner: low maturity — three lists, none authoritative. Prove identity: "
           "medium — the Identity Authority exists but education does not consume it. Certify a "
           "result: high — the Examination Authority does this well. Share data across bodies: "
           "very low — almost everything is paper. A simple maturity score per capability turns a "
           "vague sense that things are bad into a picture of exactly where the sector is weak and "
           "where it is strong.")

rows_block(prs, 'Four gaps — the ones almost every first assessment finds',
           [('Duplicate registries', 'Three learner lists where there should be one.'),
            ('Identity not consumed',
             'Education re-enters identity on paper instead of using PNIA.'),
            ('No shared data exchange',
             'Every link is a custom point-to-point connection or a paper form.'),
            ('No clear owner of the authoritative learner',
             "Everyone uses “the learner”; nobody is accountable for the copy of record.")],
           'Naming these precisely, on a real sector, is most of the assessment.',
           T,
           "VO: Four gaps surface, the ones almost every first assessment finds. Duplicate "
           "registries — three learner lists where there should be one. Identity not consumed — "
           "education re-enters identity on paper instead of using the Identity Authority. No "
           "shared data exchange — every link is a custom point-to-point connection or a paper "
           "form. And no clear owner — everyone uses 'the learner', nobody is accountable for the "
           "authoritative copy. Naming these precisely, on Progressa, is most of the Assess.\n\n"
           "Production cue: this slide answers the retrieval prompt set on the title slide. Give "
           "the viewer a beat to compare their list against these four before moving on.")

block_slide(prs, 'The ranking is the assessment — the list is only a list',
            ['Score each gap by how much it hurts and how hard it is to close. On Progressa the '
             'duplicate learner registry tops the list: it is the direct cause of the minister\'s '
             'stalled promise, it burdens every parent, and it is closable — make the Learner '
             'Registry authoritative and have the others consume it.',
             'The no-data-exchange gap is high impact but harder, so it sequences later.'],
            'What the next phases act on is the ranking, not the list.',
            T,
            "VO: A list of gaps is not an assessment. The ranking is. Score each gap by how much "
            "it hurts and how hard it is to close. On Progressa, the duplicate learner registry "
            "tops the list — it is the direct cause of the minister's stalled single-learner-record "
            "promise, it burdens every parent, and it is closable by making the Learner Registry "
            "authoritative and having the others consume it. The no-data-exchange gap is high "
            "impact but harder, so it sequences later. The ranking, not the list, is what the next "
            "phases act on.")

block_slide(prs, 'Sign-off two carries an honesty test you can fail quietly',
            ['On Progressa, one of the three learner lists is owned by a powerful programme that '
             'does not want to give it up.',
             'Your job is to name the right problems in the right order — including the political '
             'ones — in language the decision-maker can act on.'],
            'An assessment that softens that gap fails — quietly, a year later.',
            T,
            "VO: Assess ends with the second sign-off: the gap analysis reflects ground truth. "
            "This one has an honesty test. On Progressa, one of the three learner lists is owned "
            "by a powerful programme that does not want to give it up. An assessment that softens "
            "that gap to avoid the fight fails — quietly, a year later, when the single learner "
            "record still does not exist. Your job is to name the right problems in the right "
            "order, including the political ones, in language the decision-maker can act on. Get "
            "that sign-off honestly, and the roadmap has a solid base.")

big_slide(prs,
          'Assess turns Progressa’s picture into a ranked gap analysis — the right problems '
          'in the right order, signed off as ground truth — including the uncomfortable ones.',
          T,
          "VO: That is Phase 2 on Progressa. Score the capabilities. Name the four gaps. Rank by "
          "impact where movement is possible. Sign it off as honest ground truth. Six to eight "
          "weeks, and Progressa knows not just what is wrong, but what to fix first.")

sources_slide(prs, T, [
    'PAERA v1.0 — §5.1 (Capabilities Assessment)',
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
    'PAERA v1.0 — §3.1.3 (Readiness Assessment)',
])


# ================================================================ 4.4
T = '4.4 · Phase 3, Adapt — fit PAERA and decide build, buy or share'
section('4.4', 'Phase 3, Adapt — fit PAERA and decide build, buy or share',
        'Adapt fits PAERA to Progressa — localising principles, setting sector priorities, and '
        'deciding for each building block whether to build, buy, share or sandbox — signed off as '
        'the framework and sourcing approach.',
        '~5 minutes',
        "VO, slide 1: Phase three is Adapt. PAERA is a starting point, not a constraint. Now the "
        "architects shape it to Progressa — the country's own principles, its sector priorities, "
        "and the sourcing decision for each building block. About four to six weeks. The "
        "deliverable is a localised framework with a clear build, buy, share or sandbox call for "
        "each capability.")

block_slide(prs, 'Progressa does not draft principles from scratch — it inherits and localises',
            ['Progressa adopts PAERA\'s principles as its baseline and points each at its own '
             'laws — once-only against the data-protection act, reuse-before-build against the '
             'procurement rules.',
             'It adds one principle its own context needs, about offline access for rural schools.'],
            'Every principle gets a written implication the Board can rule with.',
            T,
            "VO: First, localise the principles. Progressa adopts PAERA's principles as its "
            "baseline and points each at its own laws — once-only against the data-protection act, "
            "reuse-before-build against the procurement rules. It adds one principle its context "
            "needs, about offline access for rural schools, and gives every principle a written "
            "implication so the Board can use it to settle an argument. Progressa does not draft "
            "principles from scratch. It inherits the thinking and localises the wording.")

# The sourcing matrix — the heart of Adapt. A verdict chip per building block, so the pattern
# (share and buy outnumber build) is visible before a word of the reason column is read.
s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Four kinds of decision, made on purpose, one building block at a time')
matrix = [('Learner Registry (PLR)', 'BUILD', ITU_BLUE_DARK,
           'The authoritative core. Nobody else can own it.'),
          ('Identity', 'SHARE', ITU_BLUE,
           'Consume PNIA’s platform. Do not mint a second identity.'),
          ('Exam management', 'BUY', GREY,
           'A solved problem with a mature market product.'),
          ('Data exchange', 'SHARE', ITU_BLUE,
           'PDGA’s backbone — paid for once, for the whole government.'),
          ('New analytics', 'SANDBOX', MIDGREY,
           'Nobody is sure yet. Try it before committing.')]
y = 1.68
for name, verdict, chip_fill, reason in matrix:
    panel(s, 0.72, y, 11.9, 0.82, LIGHT)
    tb = box(s, 0.95, y + 0.16, 3.5, 0.5)
    set_text(tb.text_frame, [[(name, 17, True, INK, False)]])
    panel(s, 4.62, y + 0.2, 1.65, 0.42, chip_fill, radius=0.35)
    tb = box(s, 4.62, y + 0.24, 1.65, 0.36)
    set_text(tb.text_frame, [[(verdict, 12, True, WHITE, False)]], align=PP_ALIGN.CENTER)
    tb = box(s, 6.55, y + 0.17, 5.9, 0.55)
    set_text(tb.text_frame, [[(reason, 15.5, False, INK, False)]])
    y += 0.92
tb = box(s, 0.72, 6.42, 11.9, 0.6)
set_text(tb.text_frame, [[('Four kinds of decision — each written down, each with its reason.',
                           15.5, True, ITU_BLUE_DARK, False)]])
footer(s, T)
notes(s, "VO: Then the heart of Adapt: for each building block, a deliberate sourcing call. "
         "Progressa's Learner Registry — build it properly, because it is the authoritative core "
         "nobody else can own. Identity — do not build; share the national Identity Authority's "
         "platform. Exam management — buy a mature market product, it is a solved problem. Data "
         "exchange — share the Digital Government Authority's backbone. A new analytics capability "
         "nobody is sure about — sandbox it first, before committing. Four kinds of decision, made "
         "on purpose, written down with a reason each.\n\n"
         "On screen: no vendor or product names anywhere on this slide — the call is about the "
         "kind of sourcing, not the supplier.")

block_slide(prs, 'Reuse before build is where the country-level saving is actually chosen',
            ['Notice the pattern in those calls. The default is reuse and share; build is the '
             'exception that has to justify itself. Every share — identity, data exchange — is a '
             'capability Progressa pays for once instead of in every programme.',
             'Reuse before buy, buy before build, applied building block by building block, is how '
             'Adapt turns a principle into real sourcing decisions.'],
            'A procurement rule cannot make these calls. An architect looking across the sector can.',
            T,
            "VO: Notice the pattern in those calls. The default is reuse and share; build is the "
            "exception that must justify itself. Every share — identity, data exchange — is a "
            "capability Progressa pays for once instead of in every programme. This is the phase "
            "where the whole-of-government saving is actually chosen. A procurement rule could not "
            "make these calls; only an architect looking across the whole sector, deciding "
            "deliberately, can. Reuse before buy, buy before build — applied building block by "
            "building block — is how Adapt turns the re-use principle into real sourcing decisions.")

block_slide(prs, 'Sign-off three commits the shape of everything that follows',
            ['On Progressa the Board signs the localised principle set and the '
             'build-buy-share-sandbox matrix. This is a consequential gate.',
             'Approve it, and the roadmap has a clear set of decisions to sequence.'],
            'Skip it and every project re-litigates build-versus-reuse on its own.',
            T,
            "VO: Adapt ends with the third sign-off: the localised framework and the sourcing "
            "approach are approved. On Progressa, the Board signs the principle set and the "
            "build-buy-share-sandbox matrix. This is a consequential gate — it commits the shape "
            "of everything that follows. Approve it, and the roadmap has a clear set of decisions "
            "to sequence. Skip the deliberation, and every project re-litigates build-versus-reuse "
            "on its own, which is exactly the fragmentation the EA exists to prevent.")

big_slide(prs,
          'Adapt fits PAERA to Progressa — localised principles and a deliberate build, buy, share '
          'or sandbox call for every building block — with reuse as the default that build must '
          'justify against.',
          T,
          "VO: That is Phase 3 on Progressa. Localise the principles. Make a deliberate sourcing "
          "call for each building block, with reuse as the default. Sign off the framework and the "
          "matrix. Four to six weeks, and Progressa has decided not just what it wants, but how it "
          "will get each piece — and where it will pay once instead of many times.")

sources_slide(prs, T, [
    'PAERA v1.0 — §5.6 (Sourcing Strategy)',
    'PAERA v1.0 — §5.2 (Principles)',
])


# ================================================================ 4.5
T = '4.5 · Phase 4, Plan — design the target architecture'
section('4.5', 'Phase 4, Plan — design the target architecture',
        'Between the gaps and the roadmap sits the deliverable the whole EA is for — the target '
        'architecture: a future-state picture of the capabilities, the data owners, the shared '
        'platforms and the integration map your government should have. Design it by applying your '
        'principles to the gaps, and the roadmap finally has somewhere to go.',
        '~5 minutes',
        "VO, slide 1: Assess told you where Progressa is and where the gaps are. The roadmap will "
        "tell you how to get somewhere. But somewhere where? Between the two sits the deliverable "
        "the whole architecture exists to produce: the target architecture — the picture of the "
        "future state you are building toward. Skip it, and you sequence a roadmap to a "
        "destination nobody drew. This is where you draw it — the first half of the Plan phase, "
        "before you sequence anything.")

rows_block(prs, 'The target is the current-state picture, designed instead of observed',
           [('Target capabilities',
             'Each owned by exactly one body. The duplicates resolved.'),
            ('Target data domains',
             'One authoritative owner for the learner, the person, the result.'),
            ('Target shared platforms',
             'The identity platform and the data-exchange backbone, consumed by all.'),
            ('Target technology standards',
             'What everything that follows is built to.')],
           'The same four layers as the AS-IS — drawn forward, not recorded.',
           T,
           "VO: The target architecture is the future-state picture, drawn in the same four layers "
           "you used for the current state — but designed, not observed. The target capability "
           "map: each capability owned by exactly one body, the duplicates resolved. The target "
           "data domains: one authoritative owner for the learner, the person, the result — the "
           "three lists collapsed to one. The target shared platforms: the identity platform and "
           "the data-exchange backbone every sector will consume. And the target technology "
           "standards. It is the current-state picture's mirror — what the sector should look like "
           "once the gaps are closed.")

block_slide(prs, 'You do not invent the target — you apply your principles to your gaps',
            ['Once-only says the three learner lists become one, so the target has a single '
             'authoritative Learner Registry. Reuse-before-build says education consumes the '
             'national identity rather than minting its own, so the target shows the Examination '
             'Authority consuming the Identity Authority.',
             'One-owner-per-domain resolves who owns what. Apply the principles to the current '
             'state and the target architecture almost draws itself.'],
            'The principles are not decoration. They are the design rules.',
            T,
            "VO: You do not invent the target from imagination. You design it by applying the "
            "principles you adopted in Adapt to the gaps you found in Assess. Once-only says the "
            "three learner lists become one — so the target has a single authoritative Learner "
            "Registry. Reuse-before-build says education consumes the national identity rather "
            "than minting its own — so the target shows the Examination Authority consuming the "
            "Identity Authority. One-owner-per-domain resolves who owns what. The principles are "
            "not decoration; they are the design rules. Apply them to the current state, and the "
            "target architecture almost draws itself.")

block_slide(prs, 'The integration map is the one artefact the current state never had',
            ['It names which bodies exchange which data, over what mechanism, in what priority '
             'order. On Progressa: the Examination Authority reads the learner from the Learner '
             'Registry and the identity from the Identity Authority, both over the Digital '
             'Government Authority\'s backbone, replacing the paper re-entry.',
             'It is the target\'s connective tissue — the picture of a government whose systems '
             'finally talk to each other.'],
            'A first cut here. The full design is the interoperability knowledge product.',
            T,
            "VO: The target adds one artefact the current state never had: the integration map. It "
            "names which bodies need to exchange which data, over what mechanism, in what priority "
            "order. On Progressa: the Examination Authority reads the learner from the Learner "
            "Registry and the identity from the Identity Authority, both over the Digital "
            "Government Authority's data-exchange backbone, replacing the paper re-entry. The "
            "integration map is the target's connective tissue — the picture of a government whose "
            "systems finally talk to each other. Here you draw it as a first cut; designing those "
            "exchanges in full detail is the interoperability work that the next knowledge product "
            "covers.")

block_slide(prs, 'Every target element traces back to a gap and forward to a sourcing call',
            ['The target must be reachable, not a fantasy. Each element should trace back to a gap '
             'it closes, and forward to a sourcing decision — build, buy or share — that obtains '
             'it.',
             'So the target architecture and the sourcing decisions are designed together, and '
             'signed off together, as the agreed destination the roadmap then sequences.'],
            'A target capability with no way to obtain it is a wish, not an architecture.',
            T,
            "VO: One discipline keeps the target useful. It must be reachable, not a fantasy. "
            "Every element of the target should trace back to a gap it closes and forward to a "
            "sourcing decision — build, buy or share — that obtains it. A target capability with "
            "no way to build or acquire it is a wish, not an architecture. So the target "
            "architecture and the sourcing decisions are designed together: the target says what "
            "the future state is; the sourcing says how each piece is obtained. Signed off "
            "together, they are the agreed destination the roadmap then sequences.")

big_slide(prs,
          'The target architecture is the designed future state — target capabilities, one owner '
          'per data domain, shared platforms and the integration map — reached by applying your '
          'principles to the gaps.',
          T,
          "VO: So before you sequence anything, design the target. The future-state picture in "
          "four layers, plus the integration map, reached by applying your principles to the gaps "
          "you found, every element traceable to a sourcing decision. That target is what the "
          "whole architecture exists to produce — and what the roadmap then turns into a sequence "
          "your minister can fund.")

sources_slide(prs, T, [
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
    'PAERA v1.0 — §3.4.3 (Interoperability)',
    'PAERA v1.0 — Annex 2 (Metamodel)',
])


# ================================================================ 4.6
T = '4.6 · Phase 4, Plan — sequence the roadmap and cost it'
section('4.6', 'Phase 4, Plan — sequence the roadmap and cost it',
        'Plan turns Progressa\'s decisions into a sequenced, costed roadmap in waves — the '
        'deliverable the minister takes to cabinet — signed off with budget committed.',
        '~5 minutes',
        "VO, slide 1: Phase four is Plan. Now the decisions become a sequence. The question: how "
        "does Progressa get from today to the target, in what order, at what cost? The deliverable "
        "is a roadmap in waves, with investment estimates — the single document the minister takes "
        "to cabinet. About six to eight weeks.")

# The module's centrepiece: the four-wave roadmap. Bands rather than rows, because a wave is a
# block of time with a thing in it — the shape carries the argument that each wave delivers.
s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Four waves, and each one delivers something the minister can show')
waves = [('WAVE 1', 'Inception',
          'Stand up the repository and the Board; make the Learner Registry authoritative.'),
         ('WAVE 2', 'High-priority use case',
          'The single learner record — PNEA and the grant programme consume PLR and PNIA.'),
         ('WAVE 3', 'The shared backbone',
          'Build out PDGA’s data-exchange backbone.'),
         ('WAVE 4', 'Mass scale',
          'Every education service running on the backbone.')]
y = 1.68
for label, name, detail in waves:
    panel(s, 0.72, y, 11.9, 1.05, LIGHT)
    tb = box(s, 0.98, y + 0.16, 1.5, 0.35)
    set_text(tb.text_frame, [[(label, 12, True, ITU_BLUE_DARK, False)]])
    tb = box(s, 0.98, y + 0.5, 4.2, 0.45)
    set_text(tb.text_frame, [[(name, 19, True, INK, False)]])
    tb = box(s, 5.5, y + 0.22, 6.9, 0.7)
    set_text(tb.text_frame, [[(detail, 15.5, False, INK, False)]])
    y += 1.15
tb = box(s, 0.72, 6.35, 11.9, 0.6)
set_text(tb.text_frame, [[('Each wave delivers something visible — not just groundwork.',
                           15.5, True, ITU_BLUE_DARK, False)]])
footer(s, T)
notes(s, "VO: Progressa sequences the work into waves, the way PAERA's recommended roadmap "
         "suggests. Wave one is inception — stand up the repository and the Board, and make the "
         "Learner Registry the authoritative learner. Wave two is the high-priority use case — the "
         "single learner record, by having the Examination Authority and the grant programme "
         "consume the Learner Registry and the Identity Authority instead of their own lists. Wave "
         "three builds out the shared data-exchange backbone. Wave four takes it to mass scale — "
         "every education service running on the backbone. Each wave delivers something the "
         "minister can show, not just groundwork.\n\n"
         "Production cue: this is the module's centrepiece slide. Hold it. The wave names follow "
         "PAERA §5.7.2–§5.7.5 (Inception, High-priority Use Case, Initial Transformation, "
         "Mass-scale Transformation).")

block_slide(prs, 'Cost each wave honestly — directional estimates, not quotations',
            ['Wave one is mostly people: the initial engagement and the small permanent EA team. '
             'Wave two adds the integration work connecting the authorities to the registries. '
             'Wave three is the backbone investment.',
             'Directional is enough for cabinet to commit a five-year envelope, rather than an '
             'annual line that vanishes.'],
            'The promised single learner record lands in Wave 2 — inside the first year.',
            T,
            "VO: Each wave gets an honest cost. Wave one is mostly people — the initial engagement "
            "and the small permanent EA team. Wave two adds the integration work to connect the "
            "authorities to the registries. Wave three is the backbone investment. The estimates "
            "are directional, not quotations — enough for cabinet to commit a five-year envelope, "
            "not an annual line that vanishes. And the sequencing is deliberate: the minister sees "
            "the single learner record — the thing they promised — land in wave two, inside the "
            "first year, not in year five.")

two_panel(prs, 'The roadmap is the one page both sides point at',
          ('THE MINISTER READS',
           ['Outcomes and dates.',
            'The single learner record by next year.',
            'A five-year envelope to take to cabinet.']),
          ('THE ARCHITECT READS',
           ['A sequence with dependencies.',
            'Wave 2 cannot start until PLR is authoritative in Wave 1.',
            'Where the integration work actually sits.']),
          'Same page, two readings, no translation lost.',
          T,
          "VO: The roadmap does a second job beyond planning. It is the one object the business "
          "side and the IT side both point at. The minister reads it as outcomes and dates — the "
          "single learner record by next year. The architects read it as a sequence with "
          "dependencies — you cannot do wave two before the Learner Registry is authoritative in "
          "wave one. Same page, two readings, no translation lost. The roadmap is where the policy "
          "goal and the technical sequence finally line up in a single picture the whole "
          "leadership shares.")

block_slide(prs, 'Sign-off four is the most consequential gate of the four',
            ['The decision-maker and the EA Board approve the roadmap and commit the budget — '
             'ideally a multi-year envelope, not an annual line. This is the deliverable the '
             'minister takes to cabinet.',
             'With it approved, Progressa moves from planning to building.'],
            'Four sign-offs done, roughly six months from the first day of Discovery.',
            T,
            "VO: Plan ends with the fourth sign-off, the most consequential of the four. The "
            "decision-maker and the EA Board approve the roadmap and commit the budget — ideally a "
            "multi-year envelope, not an annual line. This is the deliverable the minister takes "
            "to cabinet. With it approved, Progressa moves from planning to building. Four "
            "sign-offs done, in roughly six months from the first day of Discovery. Now the work "
            "shifts to making it live.")

big_slide(prs,
          'Plan turns Progressa’s decisions into a wave-sequenced, costed roadmap — the '
          'cabinet-ready document business and IT both read — with the minister’s flagship '
          'landing inside the first year.',
          T,
          "VO: That is Phase 4 on Progressa. Sequence into waves, each delivering something "
          "visible. Cost each wave honestly. Put the minister's promise in an early wave. Sign off "
          "the roadmap and commit the budget. Six to eight weeks, and Progressa has a plan its "
          "whole leadership shares — and can fund.")

sources_slide(prs, T, [
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
    'PAERA v1.0 — §5.7 (Recommended Roadmap — wave sequence §5.7.2–§5.7.5)',
])


# ================================================================ 4.7
T = '4.7 · Phase 5, Execute & Govern — stand up the living EA'
section('4.7', 'Phase 5, Execute & Govern — stand up the living EA',
        'Execute & Govern makes Progressa\'s architecture live — the repository, the Board, and '
        'the review gate that catches a project trying to build its own learner list and tells it '
        'to consume the registry. This is where re-use actually happens.',
        '~5 minutes',
        "VO, slide 1: Phase five is Execute and Govern — the phase that never ends. The approved "
        "roadmap becomes a project pipeline, and the small permanent EA team turns the "
        "architecture from a one-time delivery into a living practice. On Progressa, three things "
        "make it real: the repository, the Board, and the review gate. Watch the gate do its work, "
        "because this is where everything the earlier phases decided either holds or quietly "
        "unravels.\n\n"
        "Retrieval prompt — ask before playing on: a funded project in a hurry wants to build its "
        "own small learner list, because integrating with the registry would be slower for them. "
        "What should the gate do, and what happens to the country if it does nothing? The answer "
        "is two slides away.")

block_slide(prs, 'The repository goes live and a small permanent team keeps it true',
            ['The repository holds Progressa\'s layers, its sourcing decisions, its principle set '
             'and its decision log. One named architect keeps it current, updated as each wave '
             'delivers.',
             'The permanent team — two to four architects reporting to the digitalisation officer '
             '— is Progressa\'s standing muscle for cross-cutting decisions.'],
            'The team exists whether or not any single project is running.',
            T,
            "VO: First, the repository goes live — the single store of Progressa's architecture: "
            "the layers, the sourcing decisions, the principle set, the decision log. One named "
            "architect keeps it current, updated as each wave delivers. And the permanent team — "
            "two to four architects reporting to the digitalisation officer — exists whether or "
            "not any single project is running. It is Progressa's standing muscle for cross-cutting "
            "decisions.")

block_slide(prs, 'A funded project in a hurry plans to build its own learner list',
            ['A new scholarship programme arrives at the gate. Its plan: build its own small '
             'learner list, because integrating with the Learner Registry would mean learning it '
             'and waiting on its team.',
             'The gate asks its few questions. Does a shared building block already exist for '
             'this? Yes — the Learner Registry, made authoritative back in Wave 1.'],
            'For the project, building is faster. Every project makes that same rational choice.',
            T,
            "VO: Now the gate does its work. A new scholarship programme arrives, funded and in a "
            "hurry. Its plan: build its own small learner list, because integrating with the "
            "Learner Registry would mean learning it and waiting on its team. For the project, "
            "building is faster — the same rational choice every project makes. The gate asks its "
            "few questions. Does a shared building block already exist for this? Yes — the Learner "
            "Registry, made authoritative back in wave one.")

# The module's emotional peak — the only full-colour punch block in the deck.
block_slide(prs, 'The Board rules: consume the Registry. Do not build a fourth list',
            ['Left alone, the scholarship programme would have built Progressa\'s fourth learner '
             'list, and the fragmentation the assessment found would have grown instead of shrunk. '
             'The gate caught it while it was still a line in a project plan — cheap to change.',
             'A procurement rule could not have done this. Only the whole-of-government view, plus '
             'a Board that can say no, can. And the decision goes in the log, with its reason.'],
            'That moment, repeated for every project, is how re-use actually happens.',
            T,
            "VO: The Board rules: consume the Learner Registry; do not build a fourth list. This "
            "single decision is the whole point of the architecture. Left alone, the scholarship "
            "programme would have built Progressa's fourth learner list, and the fragmentation the "
            "assessment found would have grown instead of shrinking. The gate caught it while it "
            "was still a line in a project plan — cheap to change. A procurement rule could not "
            "have done this; only the whole-of-government view, plus a Board that can say no, can. "
            "That moment, repeated for every project, is how re-use actually happens — not because "
            "a strategy required it, but because the gate enforced it. And the decision goes in the "
            "log, with its reason, so the next architect knows why.\n\n"
            "Production cue: this is the pivotal slide of the whole module and its one full-colour "
            "block. Hold it a beat longer. It also answers the retrieval prompt set on the title "
            "slide.",
            punch_fill=ITU_BLUE, punch_ink=WHITE)

block_slide(prs, 'The Board is the standing rhythm no single project can provide',
            ['Four times a year the business side — the ministry, the authorities — and the IT '
             'side — the architects — sit over the same repository and decide together.',
             'That standing forum is what lets Progressa keep redesigning how it serves learners, '
             'instead of freezing after the first delivery.'],
            'The repository gives the picture, the metamodel the words, the Board the rhythm.',
            T,
            "VO: The Board does a second job. Four times a year, the business side — the ministry, "
            "the authorities — and the IT side — the architects — sit over the same repository and "
            "decide together. The repository gives them the shared picture, the metamodel the "
            "shared words, the Board the shared rhythm. That standing forum is what lets Progressa "
            "keep redesigning how it serves learners, instead of freezing after the first "
            "delivery. It is the part of an EA that no single project can provide.")

rows_block(prs, 'Six months to a roadmap, then forever',
           [('Weeks 1 to 26', 'Discovery through to an approved roadmap — the four sign-offs.'),
            ('Then, every quarter',
             'Board reviews, the repository kept current, the team protected from the urgent '
             'project of the week.'),
            ('That is what “months, not years” means',
             'Months to a roadmap; then a permanent practice.')],
           'And what keeps one learner record from drifting back into three.',
           T,
           "VO: Notice the rhythm of the whole thing. The first four phases took about six months, "
           "from Discovery to an approved roadmap. Phase five runs forever — quarterly Board "
           "reviews, the repository kept current, the team protected from being pulled onto the "
           "urgent project of the week. Months to a roadmap; then a permanent practice. That is "
           "what 'months, not years' actually means on Progressa — and what keeps the single "
           "learner record, once built, from drifting back into three.",
           numbered=False)

big_slide(prs,
          'Execute & Govern makes the architecture live — a repository kept true, a Board that '
          'enforces re-use at the gate, and the standing rhythm where business and IT keep '
          'deciding together.',
          T,
          "VO: That is Phase 5 on Progressa. The repository live and current. The Board binding. "
          "The gate catching the fourth learner list before it is built. The business and IT sides "
          "deciding together, quarterly, forever. This is where the architecture stops being a "
          "document and becomes how Progressa works.")

sources_slide(prs, T, [
    'PAERA v1.0 — §4.2.1 (Management)',
    'PAERA v1.0 — §4.2.2 (Architecture)',
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
    'GovStack Building Block specifications — govstack.global',
])


# ================================================================ 4.8
T = '4.8 · Run this on your own sector — the transferable recipe'
section('4.8', 'Run this on your own sector — the transferable recipe',
        'The five phases, four sign-offs and six deliverables you watched on Progressa are the '
        'recipe — change the institutions and the data domains, and the same method runs on any '
        'sector you are handed.',
        '~4 minutes',
        "VO, slide 1: You have watched the whole method run on Progressa. The last step is to make "
        "it yours. Because the value of a worked example is not Progressa — it is that the same "
        "five phases, four sign-offs and six deliverables run on any public-sector domain you are "
        "handed. Here is the recipe, stripped to what transfers.")

rows_block(prs, 'The recipe: five phases, six deliverables, four sign-offs',
           [('Discover', 'An honest picture of what exists today.'),
            ('Assess', 'A ranked gap analysis.'),
            ('Adapt', 'A localised framework and a build-buy-share matrix.'),
            ('Plan, first half', 'The target architecture — the designed future state.'),
            ('Plan, second half', 'The wave-sequenced, costed roadmap that reaches it.'),
            ('Execute & Govern', 'A living, governed architecture.')],
           'A sign-off after each of the first four phases. That order is the method.',
           T,
           "VO: The recipe is five phases and six deliverables, with a sign-off after each of the "
           "first four phases. Discover produces an honest picture. Assess produces a ranked gap "
           "analysis. Adapt produces a localised framework and a build-buy-share matrix. Plan "
           "produces two: the target architecture — the designed future state — and then the "
           "wave-sequenced, costed roadmap that reaches it. Execute and Govern produces a living, "
           "governed architecture. Discover before you judge; judge before you adapt; adapt before "
           "you design the target; design the target before you sequence the roadmap; govern "
           "always. That order is the method.",
           top=1.55, bottom=6.35, head_size=17)

two_panel(prs, 'Swap the contents. Keep the method',
          ('WHAT CHANGES',
           ['The institutions — your ministry, your authorities, your registries.',
            'The data domains — a patient, a farmer, a taxpayer instead of a learner.',
            'The specific gaps you find.']),
          ('WHAT DOES NOT',
           ['The five phases.',
            'The four sign-offs.',
            'The six deliverables.',
            'Reuse before build.',
            'The binding Board.']),
          'The contents are yours. The method is the same one you just watched.',
          T,
          "VO: Moving to your sector, be clear about what changes and what does not. What changes: "
          "the institutions — your ministry, your authorities, your registries instead of "
          "Progressa's. The data domains — a patient, a farmer, a taxpayer instead of a learner. "
          "The specific gaps you find. What does not change: the five phases, the four sign-offs, "
          "the six deliverables, the reuse-before-build default, and the binding Board. You swap "
          "the contents; you keep the method.")

block_slide(prs, 'The second sector is cheaper than the first, and the third cheaper still',
            ['The same five phases run on health, where the duplicated thing is a patient record; '
             'on agriculture, where it is a farmer registry; on social protection, where it is a '
             'beneficiary record.',
             'The shapes are the same — duplicate registries, paper re-entry, a stalled flagship — '
             'and so are the steps to fix them.'],
            'The team, the framework and the governance are already there.',
            T,
            "VO: And the method travels beyond education. The same five phases run on health, "
            "where the duplicated thing is a patient record; on agriculture, where it is a farmer "
            "registry; on social protection, where it is a beneficiary record. The shapes are the "
            "same — duplicate registries, paper re-entry, a stalled flagship — and so are the "
            "steps to fix them. Once your country has run the lifecycle once and built the EA "
            "muscle, the second sector is cheaper than the first, and the third cheaper still, "
            "because the team, the framework and the governance are already there.\n\n"
            "On screen: these sectors are named as places the method applies, not as worked "
            "examples — education is the worked sector throughout.")

rows_block(prs, 'Three things to carry from Progressa into your own sector',
           [('Start small',
             'Pick the one high-priority use case and make it land in an early wave.'),
            ('Sign off honestly',
             'Including the politically uncomfortable gap — the flattering assessment fails later.'),
            ('Protect the team',
             'The practice dies when the architects are pulled onto the urgent project of the '
             'week.')],
           'Get those three right and the method does the rest.',
           T,
           "VO: Three things to carry from Progressa when you start on your own sector. Start "
           "small — pick the one high-priority use case, the equivalent of the single learner "
           "record, and make it land in an early wave. Sign off honestly — including the "
           "politically uncomfortable gap, because the assessment that flatters fails a year "
           "later. And protect the team — the most common way the practice dies is the architects "
           "being pulled onto the urgent project of the week. Get those three right and the method "
           "does the rest.")

big_slide(prs,
          'Five phases, four sign-offs, six deliverables — change the institutions and the data '
          'domains, and the method you watched on Progressa runs on any sector you are handed.',
          T,
          "VO: So that is the method, end to end, on a real sector — and the recipe to run it on "
          "yours. Five phases. Four sign-offs. Six deliverables, the target architecture and the "
          "roadmap among them. Reuse before build. A Board that can say no. Change the "
          "institutions and the domains, keep the method, and you can take any sector in your "
          "country from a fragmented start to a living, governed architecture.")

sources_slide(prs, T, [
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
    'PAERA v1.0 — §5.7 (Recommended Roadmap)',
])


# ================================================================ Thank you
s = add_slide(prs, LAYOUT_THANKS)
notes(s, 'Closing slide for the combined deck. Individual videos end on their sources slide instead.')

# Self-check: the split spec's slide ranges depend on this count, and a helper that
# silently stops drawing shows up first as a slide with no voice-over.
assert len(prs.slides._sldIdLst) == 60, 'slide count changed — update decks/split_spec.json'
assert all(sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip() for sl in prs.slides), \
    'every slide carries its voice-over in the notes'

OUT = os.environ.get('OUT_PATH') or os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    'videos', 'module_4', 'en', 'decks', 'KP1_M4_Deck_v0.1.pptx')
prs.save(OUT)
print('slides:', len(prs.slides._sldIdLst))
print('saved', OUT)
