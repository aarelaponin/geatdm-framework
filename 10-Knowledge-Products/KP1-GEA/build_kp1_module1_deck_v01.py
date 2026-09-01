#!/usr/bin/env python3
# Build the KP1 Module 1 (Topic 1) video deck on the ITU template.
# Content only — every generic helper, branding constant and layout index comes from
# ITU-Giga-KP-Plugin/skills/kp-deck-builder/scripts/deck_lib.py (which also ships the
# template). Conventions and design rules: that skill's SKILL.md.
# Generated .pptx is NEVER hand-edited — fix here, re-render, re-run the split
# (kp-deck-builder/scripts/split_module_deck.py + the split spec next to the decks).
# Override paths with TEMPLATE= and OUT_PATH= env vars.
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                '..', 'ITU-Giga-KP-Plugin', 'skills', 'kp-deck-builder', 'scripts'))
from deck_lib import (
    GREY, INK, ITU_BLUE, ITU_BLUE_DARK, LIGHT, MIDGREY, PANEL_GREY, WHITE,
    LAYOUT_THANKS, LAYOUT_WHITE,
    add_slide, big_slide, box, delete_template_slides, footer, hline, mini_strip, notes,
    open_template, panel, panel_text, rows_slide, section_slide, set_text, solid,
    sources_slide, title)

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = open_template(os.environ.get('TEMPLATE'))


# ---------------------------------------------------------------- COVER (edit slide 1)
cover = prs.slides[0]
for sh in list(cover.shapes):
    if sh.shape_type == 6 and sh.name == 'Group 19':
        sh._element.getparent().remove(sh._element)
    if sh.name == 'Picture Placeholder 6':
        sh._element.getparent().remove(sh._element)

tb_map = [sh for sh in cover.shapes if sh.has_text_frame]
for sh in tb_map:
    t = sh.text_frame.text
    if 'Title of the Video' in t:
        sh.text_frame.clear()
        set_text(sh.text_frame, [[('Why a PAERA-anchored\nEnterprise Architecture', 30, True, INK, False)]])
    elif 'Subtitle' in t:
        sh.text_frame.clear()
        set_text(sh.text_frame, [
            [('KP1 · Government Enterprise Architecture · Module 1', 16, True, ITU_BLUE_DARK, False)],
            [('Eight standalone videos: why your country needs a national EA, what it is, why PAERA, the lifecycle, the four asks — and four countries that did the work.', 14, False, GREY, False)],
        ], space_after=Pt(6))
    elif 'Lenght' in t or 'Length' in t:
        sh.text_frame.clear()
        set_text(sh.text_frame, [
            [('Length: ', 14, True, INK, False), ('~34 mins across 8 videos (1.1 – 1.8)', 14, False, INK, False)],
            [('Target audience: ', 14, True, INK, False), ('CDO · Director-General · sector minister', 14, False, INK, False)],
        ], space_after=Pt(4))

# right panel with lifecycle motif (replaces [add image] block)
rp = slide_panel = cover.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.68), 0, Inches(5.653), Inches(7.5))
solid(rp, ITU_BLUE)
phases = ['Discover', 'Assess', 'Adapt', 'Plan', 'Execute & Govern']
py = 1.05
tbh = box(cover, 8.15, 0.5, 4.7, 0.5)
set_text(tbh.text_frame, [[('THE LIFECYCLE THIS MODULE TEACHES', 11, True, LIGHT, False)]])
for i, ph in enumerate(phases):
    bxs = cover.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.15), Inches(py), Inches(4.2), Inches(0.72))
    bxs.adjustments[0] = 0.14
    solid(bxs, WHITE if i < 4 else ITU_BLUE_DARK)
    tf = bxs.text_frame
    tf.margin_left = Inches(0.18)
    set_text(tf, [[(ph, 15, True, INK if i < 4 else WHITE, False)]])
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if i < 4:
        d = cover.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(8.02), Inches(py + 0.72 + 0.065), Inches(0.14), Inches(0.14))
        solid(d, LIGHT)
    py += 1.02
tbf = box(cover, 8.15, 6.35, 4.7, 0.8)
set_text(tbf.text_frame, [[('4 sign-offs · 6 months to a roadmap · then ongoing', 12, False, WHITE, False)]])
notes(cover, "Cover for the combined Topic 1 deck. Each subtopic section that follows is one standalone ~4-5 minute video. The right panel previews the module's centrepiece: the five-phase, four-sign-off lifecycle taught in video 1.6.")

# ---------------------------------------------------------------- AGENDA (edit slide 2)
agenda = prs.slides[1]
for sh in agenda.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text
    if t.strip() == 'Agenda':
        sh.text_frame.clear()
        set_text(sh.text_frame, [[('Module 1 — eight videos', 14, True, GREY, False)]])
    elif 'Subtopic 1' in t:
        sh.text_frame.clear()
        items = [
            ('1.1  Why your country needs a national EA', '~4 min'),
            ('1.2  What an EA actually is', '~3 min'),
            ("1.3  Why projects can't do this themselves", '~5 min'),
            ('1.4  Why an EA matters more now', '~5 min'),
            ('1.5  Why PAERA-anchored', '~4 min'),
            ('1.6  The lifecycle on one page', '~5 min'),
            ('1.7  What you will need from your minister', '~5 min'),
            ('1.8  Four signposts — three African, one international', '~5 min'),
        ]
        paras = []
        for name, rt in items:
            paras.append([(name, 15, True, INK, False), ('   ' + rt, 12, False, GREY, False)])
        set_text(sh.text_frame, paras, space_after=Pt(9), align=PP_ALIGN.LEFT)
    elif 'core message' in t:
        sh.text_frame.clear()
        set_text(sh.text_frame, [
            [('Without a shared plan, every new programme rebuilds what others have already built. ', 19, False, INK, True)],
            [('This module makes the case for a PAERA-anchored national EA — and names the four commitments your minister must make.', 19, False, INK, True)],
        ], space_after=Pt(10))
notes(agenda, 'Navigation slide for the combined deck; the videos themselves ship standalone on YouTube. Videos 1.1-1.4 make the case for an EA, 1.5-1.6 make the case for PAERA and show the method, 1.7-1.8 turn conviction into commitments and evidence.')

# ---------------------------------------------------------------- delete template slides 3..38
delete_template_slides(prs, keep=2)


# ---------------------------------------------------------------- section slide (module-scoped kicker)
def section(code, name, message, runtime, note):
    return section_slide(prs, 'KP1 · MODULE 1 · VIDEO %s' % code, code, name, message,
                         runtime + ' · standalone video · voice-over on text slides', note)


# ================================================================ 1.1
T = '1.1 · Why your country needs a national EA'
section('1.1', 'Why your country needs a national EA',
        'Without a shared plan for your government’s digital systems, every new programme rebuilds what others have already built. The country pays. The citizen pays. Your minister cannot deliver what they promised.',
        '~4 minutes',
        "VO, slide 1: In your ministry, you have probably seen this pattern. One programme builds a system to register citizens. Another programme builds another system to register the same citizens — for a different service. A third programme builds a third. Each takes years. Each is funded separately — often by a different donor: the World Bank, the African Development Bank, a bilateral partner. None of them work together. And your citizen still fills the same form, five times, in five different counters.\n\nYou cannot fix this inside any one programme. Each programme is doing exactly what it was funded to do.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Four signs your government has no shared plan')
rows_slide(s, [
    ('More than one register of the same people',
     'The school list has one. The national ID register has another. The social register has a third. None of them agree.'),
    ('Every system-to-system link built from scratch',
     'Last year your team connected the tax office to the business register. This year, the same work again — for the health ministry.'),
    ('Vendor lock-in on systems built years ago',
     'One vendor built it ten years ago. To change anything, you call that vendor and pay what they ask. The vendor knows this, and prices accordingly.'),
    ('Ministries that do not connect across each other',
     'The citizen who gave her ID number to the health ministry gives it again, on paper, when she enrols her child in school. “Ask once” exists on paper only.'),
])
footer(s, T)
notes(s, "Rows reveal one at a time, cumulative.\n\nRETRIEVAL MOMENT (before revealing row 1): ask the viewer — how many separate registers of the same citizens does your government keep? Hold two seconds, then reveal.\n\nVO: There are four signs that this pattern is happening in your government. Sign one — your ministry has more than one register of the same people... Sign two — every time two systems need to share data, you build a new connection from zero... Sign three — a vendor built one of your systems ten years ago; nobody else knows how it works... Sign four — each ministry has its own digital systems, and they do not connect across ministries. The pledge to ask each citizen for information only once exists on paper. In practice, it is impossible.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'The country pays in four directions at once')
cells = [
    ('Money', 'The same systems built many times over.'),
    ('Time', 'Every new programme waits for cross-system work nobody planned for.'),
    ('Citizens', 'The same form, filled in five places.'),
    ('Your minister', 'Flagship cross-ministry programmes cannot land — the systems will not talk.'),
]
gx, gy, gw, gh, gap = 0.72, 1.65, 5.9, 2.45, 0.25
for i, (h, b) in enumerate(cells):
    x = gx + (i % 2) * (gw + gap)
    y = gy + (i // 2) * (gh + gap)
    panel(s, x, y, gw, gh, LIGHT if i < 3 else ITU_BLUE)
    panel_text(s, x, y, gw, gh, [
        [(h, 20, True, INK if i < 3 else WHITE, False)],
        [(b, 15, False, GREY if i < 3 else WHITE, False)],
    ])
footer(s, T)
notes(s, "VO: The cost of this pattern runs in four directions at once. You pay more, because you build the same things many times. The country moves more slowly, because every new programme waits for cross-system work that nobody planned for. Your citizens carry the burden, because they fill the same form in five places. And your minister cannot deliver a flagship cross-ministry programme, because the systems will not talk to each other.\n\nThe minister cell is the one to land hardest with a Strategist audience — it is highlighted deliberately.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'One root cause: there is no shared plan')
tb = box(s, 0.72, 2.0, 11.8, 2.6)
set_text(tb.text_frame, [
    [('No one has written down what your government’s digital systems are, who owns them, and how they should fit together.', 22, False, INK, False)],
], space_after=Pt(10))
p = panel(s, 0.72, 4.4, 11.8, 1.6, LIGHT)
panel_text(s, 0.72, 4.45, 11.8, 1.6, [
    [('A national Enterprise Architecture is that plan.', 26, True, ITU_BLUE_DARK, False)],
])
footer(s, T)
notes(s, "VO: All four signs come from one root cause. There is no shared plan. No one has written down what your government's digital systems are, who owns them, and how they should fit together. That shared plan is what a national Enterprise Architecture provides. The rest of this knowledge product shows you how to commission one, what it will deliver, and what you will need from your minister to make it work.\n\nThis is the most important slide of video 1.1 — hold it a beat longer.")

sources_slide(prs, T, [
    'PAERA v1.0 — §2.1 Problem statement (paera.govstack.global)',
    'PAERA v1.0 — §5.2 Principle #5, Once-Only',
    'EU European Interoperability Framework — Once-Only principle',
])

# ================================================================ 1.2
T = '1.2 · What an EA actually is'
section('1.2', 'What an EA actually is',
        'An EA is the picture everyone agrees describes your government — minister, ministry CIO, donor, vendor. With it, you can lead the conversation. Without it, others lead it for you.',
        '~3 minutes',
        "VO, slide 1: An Enterprise Architecture is a set of documents and diagrams. Together, they describe how your government works. What services it delivers, and to whom. What data it holds, and who owns it. What software supports those services. What infrastructure runs underneath.\n\nAn EA is not software. It is not a vendor product. It is not a tool you buy. It is the agreed picture.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'An EA is the agreed picture — not software you buy')
colw, colh, cy = 5.9, 4.7, 1.75
panel(s, 0.72, cy, colw, colh, LIGHT)
panel_text(s, 0.72, cy, colw, colh, [
    [('An EA is', 20, True, ITU_BLUE_DARK, False)],
    [('the agreed picture of your government', 17, False, INK, False)],
    [('a set of documents and diagrams', 17, False, INK, False)],
    [('the same description everyone uses — minister, CIO, donor, vendor', 17, False, INK, False)],
])
panel(s, 6.87, cy, colw, colh, PANEL_GREY)
panel_text(s, 6.87, cy, colw, colh, [
    [('An EA is not', 20, True, GREY, False)],
    [('software', 17, False, INK, False)],
    [('a vendor product', 17, False, INK, False)],
    [('a tool you procure', 17, False, INK, False)],
])
footer(s, T)
notes(s, "VO: An EA is not software. It is not a vendor product. It is not a tool you buy. It is the agreed picture.\n\nWhy does it matter that everyone has the same picture? Because every important conversation in your ministry breaks down on this point. This slide is the Strategist's first line of defence against vendor pitches — when a vendor says 'our product gives you an EA', this is the slide that says no, it does not.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Minister, donor, vendor — all use the same picture')
rows_slide(s, [
    ('Your minister — to brief cabinet',
     'They cannot describe what the country’s digital spend is buying without an agreed picture of what the systems are.'),
    ('The donor — before funding the next programme',
     'They want to see how their investment fits with the others.'),
    ('The vendor — to match what they propose', 'What they build must fit what is already there.'),
    ('You — to keep all three aligned',
     'So the donor funds what the country needs, the vendor builds what fits, and the minister tells a coherent story.'),
], numbered=False)
footer(s, T)
notes(s, "VO: Your minister uses the picture to brief cabinet... The donor uses the picture before they fund the next programme... The vendor uses the picture when they propose a system... And you, the middle manager, use the picture to keep all three aligned.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'The picture has four parts')
layers = [
    ('Services', 'What your government does, and for whom. The layer your minister talks about — the layer the public sees.'),
    ('Data', 'The information your government holds, who owns it, where the authoritative copy lives.'),
    ('Applications', 'The software that uses the data to deliver the services. What gets bought, built, integrated, replaced.'),
    ('Technology', 'The infrastructure underneath — networks, hosting, identity, security.'),
]
ly, lh = 1.6, 1.22
for i, (h, b) in enumerate(layers):
    y = ly + i * (lh + 0.08)
    fill = LIGHT if i != 1 else ITU_BLUE
    panel(s, 0.72, y, 9.6, lh, fill)
    panel_text(s, 0.72, y, 9.6, lh, [
        [(h, 18, True, INK if i != 1 else WHITE, False), ('   ' + b, 14, False, GREY if i != 1 else WHITE, False)],
    ], m=0.25)
side = box(s, 10.55, 2.55, 2.3, 2.2)
set_text(side.text_frame, [
    [('Applications come and go. Technology cycles every decade.', 13, False, GREY, True)],
    [('Data outlasts them all.', 15, True, ITU_BLUE_DARK, False)],
], space_after=Pt(8))
footer(s, T)
notes(s, "Rows build top to bottom, cumulative. Order is deliberate — BDAT with data before applications, reflecting data-first architecture.\n\nVO: Every EA looks at your government in four parts. Services — what your government does, and for whom... Data — the information your government holds. Data is the longest-lived part of any government. Applications come and go. Technology cycles every decade. Data outlasts them all. Applications — the software that uses that data... Technology — the basics that must be running for anything else to work.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Written down, a list of projects becomes a system')
tb = box(s, 0.72, 1.9, 11.8, 4.5)
set_text(tb.text_frame, [
    [('When you write down all four parts, your digital landscape stops being a list of unrelated projects. It starts looking like a system —', 22, False, INK, False)],
    [('a system you can change deliberately,', 22, False, INK, False)],
    [('a system you can plan investments against,', 22, False, INK, False)],
    [('a system you can talk about across ministries, with donors and vendors, without re-explaining the basics every time.', 22, False, INK, False)],
], space_after=Pt(14))
footer(s, T)
notes(s, "VO: When you write down all four parts, your government's digital landscape stops being a list of unrelated projects. It starts looking like a system. A system you can change deliberately. A system you can plan investments against. A system you can talk about across ministries, with donors and with vendors, without re-explaining the basics every time.")

big_slide(prs, 'An EA is the agreed picture of your government — and the tool that lets you lead the conversation about what comes next.',
          T,
          "VO: That is what an Enterprise Architecture is. Not software. Not a tool. The agreed picture of your government — and the tool that lets you, instead of the vendor or the donor, lead the conversation about what comes next.\n\nScreenshot-ready summary; Strategists reuse this line in their own briefings.")

sources_slide(prs, T, [
    'PAERA v1.0 — §2.3 Role of Enterprise Architecture',
    'TOGAF — BDAT layering reference',
])

# ================================================================ 1.3
T = "1.3 · Why projects can't do this themselves"
section('1.3', "Why projects can't do this themselves",
        'Procurement rules can require interoperability. They cannot deliver it. Only planning at the level of the whole government, supported by reference architectures, can.',
        '~5 minutes',
        "VO, slide 1: You may be thinking: my country already requires this. Every new digital project must specify open APIs. Every new contract must require interoperability. The national digital strategy is signed by cabinet. So why does the citizen still fill the same form five times? Why do new programmes still build their own version of identity and payments? The answer is uncomfortable. Procurement rules can require behaviour. They cannot make that behaviour the cheapest choice for the project doing the work.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Each project is rational — and builds its own')
rows_slide(s, [
    ('The question every project asks', '“How do we ship on time and on budget?”'),
    ('What re-use costs the project', 'Learn another ministry’s system. Negotiate with their team. Accept their delays.'),
    ('So the team builds its own — and that is not a failure of discipline',
     'Building your own version is faster and simpler. The project is doing exactly what it was funded to do.'),
], numbered=False, head_size=20, sub_size=16)
footer(s, T)
notes(s, "VO: Inside any new programme, the team is rational. They have a contract, a budget, a deadline. Re-using another ministry's identity system means learning that system, negotiating with that ministry's team, accepting their delays. Building your own version is faster. So the team builds their own. That is not a failure of discipline. That is the project doing exactly what it was funded to do.\n\nTone: the math, not the morality — no blame on project teams.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'The math changes only at whole-of-government level')
chain = [
    ('Ministry 1', 'pays to build the identity system'),
    ('Ministry 2', 'does not — it consumes the first'),
    ('Ministry 3', 'does not — it consumes the same'),
    ('Ministry 4', 'does not — it consumes the same'),
]
cw, ch2, cx0, cy0, cgap = 2.85, 1.7, 0.72, 1.8, 0.15
for i, (h, b) in enumerate(chain):
    x = cx0 + i * (cw + cgap)
    fill = ITU_BLUE if i == 0 else LIGHT
    panel(s, x, cy0, cw, ch2, fill)
    panel_text(s, x, cy0, cw, ch2, [
        [(h, 16, True, WHITE if i == 0 else INK, False)],
        [(b, 13, False, WHITE if i == 0 else GREY, False)],
    ], m=0.18)
tb = box(s, 0.72, 4.0, 11.8, 1.1)
set_text(tb.text_frame, [
    [('Over five years, the country saves a meaningful share of its sectoral digital spend.', 20, True, INK, False)],
])
tb = box(s, 0.72, 5.15, 11.8, 1.6)
set_text(tb.text_frame, [
    [('But this view does not exist inside any single project. It exists only at the level of your country’s whole digital portfolio — the view an EA gives you, and the view your minister needs to make funding decisions.', 16, False, GREY, False)],
])
footer(s, T)
notes(s, "VO: Re-use becomes rational only when you can see across the whole government. From that view, the math changes. The first ministry pays to build the identity system. The second ministry does not — it consumes the first one. The third does not. The fourth does not. Over five years, the country saves a meaningful share of its sectoral digital spend.\n\nBut this view does not exist inside any single project. It exists only at the level of your country's whole digital portfolio. That is the view an EA gives you, and it is the view your minister needs to make funding decisions that look different from project-level cost choices.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Two more things projects cannot deliver')
panel(s, 0.72, 1.7, 5.9, 4.6, LIGHT)
panel_text(s, 0.72, 1.7, 5.9, 4.6, [
    [('Sustainability', 20, True, ITU_BLUE_DARK, False)],
    [('Twelve years from now, the original vendor is gone. The open-source library has forked. The technology has moved on. The system is still in service — and nobody fully understands it.', 16, False, INK, False)],
    [('Projects are not incentivised to plan for that moment. They are incentivised to ship.', 16, True, INK, False)],
])
panel(s, 6.87, 1.7, 5.9, 4.6, LIGHT)
panel_text(s, 6.87, 1.7, 5.9, 4.6, [
    [('Complexity reduction', 20, True, ITU_BLUE_DARK, False)],
    [('A project says yes to most feature requests, because requests come from people the project must please. Five years later, the system is too complex to maintain or change.', 16, False, INK, False)],
    [('Saying no requires authority the project does not have.', 16, True, INK, False)],
])
footer(s, T)
notes(s, "VO: There are two more things projects do not deliver. The first is sustainability. A project ships on time and moves on. Twelve years later, the original vendor is gone... The second is complexity reduction. A project says yes to most feature requests... Each yes has accumulated. Saying no requires authority the project does not have.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, "A reference architecture is other countries' learning, written down")
rows_slide(s, [
    ('Other countries have built these systems for decades', ''),
    ('They learned what survives technology change — and what cost them dearly', ''),
    ('A reference architecture is that learning, written down. PAERA is one such reference.',
     'Adopting it means starting with their lessons, instead of paying for the same lessons yourself over the next twenty years.'),
], numbered=False, head_size=21, sub_size=16)
footer(s, T)
notes(s, "VO: Reference architectures answer the questions projects do not. Other countries have built these systems for decades. They learned what survives technology change and what does not. They learned which architectural decisions still hold up after twenty years and which decisions cost them dearly. A reference architecture is that learning, written down. PAERA is one such reference. Adopting it means starting with their lessons, instead of paying for the same lessons yourself over the next twenty years.")

big_slide(prs, 'Re-use, sustainability, complexity reduction — none of these come from projects. All of them come from planning.',
          T,
          "VO: So when you make the case for an EA, the case is this. The four-part picture is the artefact. Planning is the function. Re-use, sustainability and complexity reduction — three things every digital ministry says it wants — none of these come from projects alone. All of them come from the planning view that an EA gives, and from the reference architectures that EA practice connects you to. That is why projects cannot fix this themselves. That is why an EA is work your minister must commission, separately, deliberately, with sustained funding.",
          sub='The four-part picture is the artefact. Planning is the function. That is why your minister must commission an EA separately, deliberately, with sustained funding.')

sources_slide(prs, T, [
    'PAERA v1.0 — §1.3 GovStack Vision',
    'PAERA v1.0 — §3.3 Digital Infrastructure principles',
    'PAERA v1.0 — §5.2 Principles',
])

# ================================================================ 1.4
T = '1.4 · Why an EA matters more now'
section('1.4', 'Why an EA matters more now',
        'For thirty years, digital work meant putting paper online. That era is ending. The work now is to redesign how your ministry serves citizens — and that work needs business and IT in the same room, using the same words.',
        '~5 minutes',
        "VO, slide 1: For a long time, digital transformation in government meant one thing. Take a paper process and put it online. The application form becomes a web form. The queue becomes an online appointment. The certificate becomes a PDF. The ministry still does the same work, in the same order, with the same roles. Only the medium changes.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Putting paper online is no longer the main work')
panel(s, 0.72, 1.7, 5.9, 4.6, PANEL_GREY)
panel_text(s, 0.72, 1.7, 5.9, 4.6, [
    [('The old work', 19, True, GREY, False)],
    [('Take a paper process and put it online', 16, False, INK, False)],
    [('The ministry’s operating model stays the same', 16, False, INK, False)],
    [('Business decides what; IT delivers how', 16, False, INK, False)],
])
panel(s, 6.87, 1.7, 5.9, 4.6, LIGHT)
panel_text(s, 6.87, 1.7, 5.9, 4.6, [
    [('The new work', 19, True, ITU_BLUE_DARK, False)],
    [('Redesign how citizens are served', 16, False, INK, False)],
    [('The ministry’s operating model changes', 16, False, INK, False)],
    [('Business and IT decide together', 16, True, INK, False)],
])
footer(s, T)
notes(s, "VO: That work is still important. But it is no longer the most important work your ministry is being asked to do. The countries — and the ministries — that are delivering real citizen results today are not just digitising forms. They are redesigning how citizens are served.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'The new work redesigns how citizens are served')
rows_slide(s, [
    ('One farmer registry',
     'Shared by the agriculture ministry, the cooperative bank, the input subsidy programme and the climate-resilience programme — the same farmer, recognised the same way, by all of them.'),
    ('One national identity',
     'A citizen proves who they are at any service, without paper.'),
    ('One learner record',
     'It follows the child from primary school to university.'),
], numbered=False, head_size=21, sub_size=16)
tb = box(s, 0.72, 6.15, 11.8, 0.7)
set_text(tb.text_frame, [[('None of these is “put it online.” Each is a redesign of how the ministry works.', 17, True, ITU_BLUE_DARK, False)]])
footer(s, T)
notes(s, "VO: Look at what the new work actually looks like. One farmer registry, used by the agriculture ministry, the cooperative bank, the input subsidy programme and the climate-resilience programme — the same farmer, recognised the same way, by all of them. One national identity that lets a citizen prove who they are at any service, without paper. One learner record that follows the child from primary school to university. None of these is 'put it online.' Each of them is a redesign of how the ministry works.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'The new work needs business and IT to decide together')
tb = box(s, 0.72, 2.0, 11.8, 4.3)
set_text(tb.text_frame, [
    [('When you redesign how your ministry works, two groups must decide together.', 22, False, INK, False)],
    [('The business side: ', 20, True, INK, False), ('your minister, your director-general, your head of policy.', 20, False, INK, False)],
    [('The IT side: ', 20, True, INK, False), ('your architects, your engineers.', 20, False, INK, False)],
    [('In the old work, business decided what and IT delivered how. In the new work, they decide together — about the same questions, with the same seriousness.', 20, False, GREY, False)],
], space_after=Pt(16))
footer(s, T)
notes(s, "VO: But the new work comes with a new problem. When you redesign how the ministry works, two groups must decide together. The business side: your minister, your director-general, your head of policy. The IT side: your architects, your engineers. In the old work, they did not need to talk much. Business decided what; IT delivered how. In the new work, they decide together. About the same questions. With the same level of seriousness.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'They sit in the same meeting — and miss each other')
panel(s, 0.72, 1.8, 5.9, 2.4, LIGHT)
panel_text(s, 0.72, 1.8, 5.9, 2.4, [
    [('Your minister talks about', 15, False, GREY, False)],
    [('citizen services and policy goals', 20, True, INK, False)],
])
panel(s, 6.87, 1.8, 5.9, 2.4, LIGHT)
panel_text(s, 6.87, 1.8, 5.9, 2.4, [
    [('Your chief architect talks about', 15, False, GREY, False)],
    [('systems, APIs and data', 20, True, INK, False)],
])
tb = box(s, 0.72, 4.7, 11.8, 1.8)
set_text(tb.text_frame, [
    [('The decision does not get made well — or it does not get made at all.', 21, True, INK, False)],
    [('And the ministry stays in the old work, even though the new work is what is needed.', 17, False, GREY, False)],
], space_after=Pt(10))
footer(s, T)
notes(s, "VO: But they do not share a language. Your minister talks about citizen services and policy goals. Your chief architect talks about systems, APIs and data. They sit in the same meeting and miss each other. The decision does not get made well — or it does not get made at all. And the ministry stays in the old work, even though the new work is what is needed.\n\nThis is the Strategist's recognition moment — they have lived this meeting.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'An EA gives both sides one language')
rows_slide(s, [
    ('The agreed picture', 'Something both sides can point at.'),
    ('Plain words for the basic terms', 'Service, capability, data domain — meaning the same thing to both sides.'),
    ('A regular forum — the EA Board', 'Where business and IT sit together and decide.'),
], head_size=21, sub_size=16)
tb = box(s, 0.72, 6.2, 11.8, 0.7)
set_text(tb.text_frame, [[('With these in place, the redesign conversation finally happens. Without them, it does not.', 16, False, GREY, True)]])
footer(s, T)
notes(s, "VO: This is the second job an Enterprise Architecture does. It gives both sides a shared language for the new work. The agreed picture — something both sides can point at. Plain words for the basic terms — service, capability, data domain — that mean the same thing to both sides. And a regular forum — the EA Board — where they sit together and decide. With these in place, the conversation about redesigning how the ministry works finally happens. Without them, it does not.")

big_slide(prs, 'In the era of digitising paper, an EA was useful. In the era of redesigning how the ministry works, an EA is necessary.',
          T,
          "VO: When you make the case to your minister, this is the second half of the case. Planning is one half. Shared language is the other. In the era of digitising paper, an EA was useful. In the era of redesigning how the ministry works, an EA is necessary — because without it, the business side and the IT side cannot have the conversation the redesign requires.")

sources_slide(prs, T, [
    'PAERA v1.0 — §2.3 Role of Enterprise Architecture',
    'PAERA v1.0 — §4.5 Digital Co-creation',
])

# ================================================================ 1.5
T = '1.5 · Why PAERA-anchored'
section('1.5', 'Why PAERA-anchored',
        'PAERA gives your team five years of head start. Adopt it, and the architecture work begins on day one. Do not adopt it, and your first year is spent inventing what others have already published.',
        '~4 minutes',
        "VO, slide 1: An Enterprise Architecture is the agreed picture of your government. The next question is: which framework do you use to draw that picture? You have two paths.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Two paths — one costs you a year before work begins')
b1 = panel(s, 0.72, 2.1, 11.9, 1.5, PANEL_GREY)
panel_text(s, 0.72, 2.1, 11.9, 1.5, [
    [('Design your own framework', 18, True, INK, False)],
    [('12 to 18 months of consultants and groundwork before any architect draws a single picture of any ministry', 15, False, GREY, False)],
], m=0.28)
b2 = panel(s, 0.72, 4.1, 3.6, 1.5, ITU_BLUE)
panel_text(s, 0.72, 4.1, 3.6, 1.5, [
    [('Anchor on PAERA', 18, True, WHITE, False)],
    [('work begins on day one', 15, False, WHITE, False)],
], m=0.28)
tb = box(s, 4.6, 4.35, 8.0, 1.3)
set_text(tb.text_frame, [
    [('The Public Administration Ecosystem Reference Architecture — published 2024 under GovStack, built for the public-sector use case, with the groundwork already done.', 15, False, GREY, False)],
])
footer(s, T)
notes(s, "Bar length is the visual: the long grey bar is the time you spend before architecture work starts; the short blue bar is PAERA. Do not label the top bar pejoratively — the lengths do the work.\n\nVO: The first path: hire consultants to design an EA framework specifically for your country. A new way of organising ministries. A new set of terms. A new set of principles. A new method. This takes twelve to eighteen months, before any architect has drawn a single picture of any ministry. You have paid for groundwork, not for architecture.\n\nThe second path: anchor on a framework that already exists... PAERA is that framework — the Public Administration Ecosystem Reference Architecture, published in 2024 under GovStack.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Five things you do not pay to invent')
rows_slide(s, [
    ('A way to organise ministries and authorities',
     'Every body is a policy unit, a regulatory agency or a service-delivery authority — the classification is published, not argued over for three months.'),
    ('Plain words for the basic terms', 'Capability. Service. Application. Data domain. Defined on day one.'),
    ('A library of common building blocks',
     'Identity, payments, information sharing, registries — published specifications with working examples.'),
    ('Ten architectural principles',
     'From rule of law to user-centred design — already debated across many countries. You extend them; you do not draft them.'),
    ('A way to run the work',
     'The lifecycle — Discover, Assess, Adapt, Plan, Execute & Govern — with roles, decisions and sign-offs per phase. Next video.'),
], top=1.5, bottom=6.95, head_size=17.5, sub_size=13.5)
footer(s, T)
notes(s, "Rows reveal one at a time. This is the substantive payload of video 1.5.\n\nVO: Five things PAERA gives your team on day one. Five things you do not pay to invent. First, a way to organise ministries and authorities... Second, plain words for the basic terms... Third, a library of common building blocks... Fourth, a set of architectural principles — ten of them... Fifth, a way to run the work. We walk the lifecycle in the next video.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'One framework — every sector reuses the investment')
sectors = ['Education', 'Health', 'Social protection', 'Agriculture']
swd, sx0, sy0 = 2.75, 0.9, 2.2
for i, sec in enumerate(sectors):
    x = sx0 + i * (swd + 0.15)
    pnl = panel(s, x, sy0, swd, 1.0, LIGHT)
    tf = pnl.text_frame
    set_text(tf, [[(sec, 16, True, INK, False)]], align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
base = panel(s, 0.9, 3.5, 11.45, 1.0, ITU_BLUE)
tf = base.text_frame
set_text(tf, [[('PAERA — the same framework underneath', 18, True, WHITE, False)]], align=PP_ALIGN.CENTER)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tb = box(s, 0.9, 5.0, 11.5, 1.4)
set_text(tb.text_frame, [
    [('Once your country builds the EA muscle for one sector, every next sector reuses the same investment.', 18, False, INK, False)],
    [('The second sector is cheaper than the first. The third is cheaper still.', 18, True, ITU_BLUE_DARK, False)],
], space_after=Pt(8))
footer(s, T)
notes(s, "VO: PAERA works across sectors. The same framework applies to education, to health, to social protection, to agriculture. Once your country builds the EA muscle for one sector, every next sector reuses the same investment. The framework does not get re-bought, re-trained, re-customised. The second sector is cheaper than the first. The third is cheaper still.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'PAERA connects your team to a working network')
sats = [
    ('Building Block specifications', 0.9, 1.75),
    ('GovMarket — compliance-checked implementations', 7.6, 1.75),
    ('Certification programme', 0.9, 5.0),
    ('Sandbox for prototyping', 7.6, 5.0),
    ('Shared knowledge base', 4.4, 5.9),
]
# connectors first, so panels sit on top of the line ends
for (txt, x, y) in sats:
    w = 4.6
    cn = s.shapes.add_connector(1, Inches(x + w / 2), Inches(y + 0.475), Inches(6.67), Inches(3.7))
    cn.line.color.rgb = MIDGREY
    cn.line.width = Pt(1.2)
for (txt, x, y) in sats:
    w = 4.6
    pnl = panel(s, x, y, w, 0.95, LIGHT)
    tf = pnl.text_frame
    tf.margin_left = Inches(0.15)
    set_text(tf, [[(txt, 14.5, True, INK, False)]], align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
hub = panel(s, 5.27, 3.15, 2.8, 1.1, ITU_BLUE)
tf = hub.text_frame
set_text(tf, [[('PAERA', 22, True, WHITE, False)]], align=PP_ALIGN.CENTER)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
footer(s, T)
notes(s, "Keep this factual, not promotional — these are real assets a country gets access to, not buzzwords.\n\nVO: PAERA is not a standalone document. It connects to a working network — building block specifications, a marketplace of vendor-built implementations that have been checked for compliance, a certification programme, a sandbox you can test in, a shared knowledge base. Adopting PAERA means joining a network of countries and partners actively building these resources together.")

big_slide(prs, 'PAERA is not a vendor choice. It is the choice to start with the work other countries have already done.',
          T,
          "VO: When you make the case for PAERA to your minister, this is the case. It is not a vendor choice. It is the choice to start with the work other countries have already done — so your team can spend their time on what is specific to your country. That is the head start. That is what your minister is being asked to commit to.",
          sub='So your team spends its time on what is specific to your country.')

sources_slide(prs, T, [
    'PAERA v1.0 — paera.govstack.global (§1.2 Motivation; §1.3 GovStack Vision)',
    'GovStack — govstack.global',
    'GovMarket — marketplace of compliance-checked implementations',
])

# ================================================================ 1.6
T = '1.6 · The lifecycle on one page'
section('1.6', 'The lifecycle on one page',
        'Six months from start to a roadmap your minister can take to cabinet. Then ongoing governance. Five phases. Four sign-offs. One continuous practice.',
        '~5 minutes',
        "VO, slide 1: If you take one picture away from this knowledge product, take this one. The EA lifecycle on a single page. The picture that goes on the wall of the EA Board room. The picture your minister puts in every cabinet briefing. The picture your team points at when they explain where the work is.\n\nRETRIEVAL MOMENT: before showing the next slide, ask the viewer to guess — from a standing start, how long until a roadmap your minister can take to cabinet? Answer comes on the timeline: six months.")

PHASES = [
    ('Discover', 'What exists today?', '3–4 wks'),
    ('Assess', 'What is the gap?', '6–8 wks'),
    ('Adapt', 'What fits our country?', '4–6 wks'),
    ('Plan', 'How do we get there?', '6–8 wks'),
    ('Execute\n& Govern', 'How do we sustain this?', 'Ongoing'),
]

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Five phases. Four sign-offs. Six months to a roadmap.')
fx0, fy0, fw, fh, fgap = 0.6, 2.15, 2.28, 1.55, 0.19
for i, (name, q, dur) in enumerate(PHASES):
    x = fx0 + i * (fw + fgap)
    fill = LIGHT if i < 4 else ITU_BLUE
    pnl = panel(s, x, fy0, fw, fh, fill)
    tf = pnl.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    set_text(tf, [
        [(name.replace('\n', ' '), 16, True, INK if i < 4 else WHITE, False)],
        [(q, 12, False, GREY if i < 4 else WHITE, True)],
    ], align=PP_ALIGN.CENTER)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tb = box(s, x, fy0 + fh + 0.12, fw, 0.4)
    set_text(tb.text_frame, [[(dur, 13, True, ITU_BLUE_DARK if i < 4 else GREY, False)]], align=PP_ALIGN.CENTER)
    if i < 4:
        d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(fx0 + (i + 1) * (fw + fgap) - fgap / 2 - 0.1), Inches(fy0 + fh / 2 - 0.1), Inches(0.2), Inches(0.2))
        solid(d, ITU_BLUE_DARK)
tb = box(s, 0.6, 4.55, 12.1, 0.5)
set_text(tb.text_frame, [[('◆ = sign-off by the senior decision-maker. Each phase answers one question and produces one deliverable.', 13.5, False, GREY, False)]])
br = panel(s, 0.6, 5.25, 9.65, 0.55, PANEL_GREY)
tf = br.text_frame
set_text(tf, [[('Phases 1–4: about six months, four sign-offs', 13, True, GREY, False)]], align=PP_ALIGN.CENTER)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
br2 = panel(s, 10.45, 5.25, 2.28, 0.55, ITU_BLUE)
tf = br2.text_frame
set_text(tf, [[('permanent practice', 13, True, WHITE, False)]], align=PP_ALIGN.CENTER)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tb = box(s, 0.6, 6.1, 12.1, 0.6)
set_text(tb.text_frame, [[('The centrepiece of Module 1 — designed to be screenshot and put on the EA Board room wall.', 13, False, GREY, True)]])
footer(s, T)
notes(s, "VO: Five phases. Each one answers a single question. Each produces a single deliverable. Each ends with a sign-off by the senior decision-maker. The first four phases together take about six months. The fifth phase is ongoing — your country's permanent way of working.")


def phase_slide(idx, deliverable, signoff, dur_label, vo):
    name, q, dur = PHASES[idx]
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, 'Phase %d — %s' % (idx + 1, name.replace('\n', ' ')))
    mini_strip(s, len(PHASES), idx)
    fields = [
        ('Question', q),
        ('Deliverable', deliverable),
        ('Sign-off' if idx < 4 else 'Cadence', signoff),
        ('Duration', dur_label),
    ]
    y = 1.65
    heights = [1.0, 1.55, 1.45, 0.85]
    for k, (lab, val) in enumerate(fields):
        tb = box(s, 0.72, y, 2.3, 0.6)
        set_text(tb.text_frame, [[(lab.upper(), 13, True, ITU_BLUE_DARK, False)]])
        tb = box(s, 3.1, y - 0.02, 9.4, heights[k])
        set_text(tb.text_frame, [[(val, 18, k in (1, 2), INK, False)]])
        y += heights[k] + 0.18
    footer(s, T)
    notes(s, vo)
    return s


phase_slide(0, 'Discovery brief — the current picture: strategies in force, systems, sector plans, stakeholders, legal framework. No recommendations yet.',
            'The senior decision-maker confirms the picture is accurate enough to build on.', '3 to 4 weeks',
            "VO: Phase one. Discover. The question: what exists today? Your architects map the current digital landscape — the strategies in force, the systems that exist, the sector plans, the stakeholders, the legal framework. The deliverable is a Discovery brief: the picture of where your country is, with no recommendations yet. The sign-off: the senior decision-maker confirms the picture is accurate enough to build on. Three to four weeks.")
phase_slide(1, 'Assessment report — current-state picture in the four parts of an EA, maturity scorecards, gap analysis against PAERA-anchored standards.',
            'The gap analysis reflects ground truth — it names the right problems in the right priority.', '6 to 8 weeks',
            "VO: Phase two. Assess. The question: what is the gap? Architects compare what they found in Discovery against PAERA-anchored standards. They write the current-state picture in the four parts of an EA. They produce maturity scorecards. They write the gap analysis. The sign-off: the gap analysis reflects ground truth. Six to eight weeks.")
phase_slide(2, 'Localised framework with sourcing decisions — your principles, your sector priorities; per building block: build, buy, share, or sandbox first.',
            'The framework and the build / buy / share approach.', '4 to 6 weeks',
            "VO: Phase three. Adapt. The question: what fits our country? PAERA is a starting point, not a constraint. Architects work with sector CIOs and your EA Board to shape the framework to your country. For each building block, the question is: do we build it, do we buy it from the marketplace, do we share another country's, or do we test in a sandbox first? Four to six weeks.")
phase_slide(3, 'Roadmap and investment plan — target state, sequenced work, investment estimates. The deliverable your minister can take to cabinet.',
            'Roadmap approved, budget committed — by the senior decision-maker and the EA Board.', '6 to 8 weeks',
            "VO: Phase four. Plan. The question: how do we get there? Architects describe the target state and sequence the work into a roadmap with investment estimates. The senior decision-maker and the EA Board approve the roadmap and commit budget. This is the deliverable your minister can take to cabinet. Six to eight weeks.")
s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Your minister reviews four times — not every diagram')
hline(s, 0.9, 2.75, 11.5, MIDGREY, weight_pt=2.5)
gates = [
    ('The picture is accurate', 'end of Discover'),
    ('The gap reflects ground truth', 'end of Assess'),
    ('Framework + sourcing approach', 'end of Adapt'),
    ('Roadmap approved, budget committed', 'end of Plan'),
]
for i, (g, when) in enumerate(gates):
    x = 1.5 + i * 2.95
    d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(2.62), Inches(0.28), Inches(0.28))
    solid(d, ITU_BLUE_DARK)
    tb = box(s, x - 1.15, 3.15, 2.6, 1.3)
    set_text(tb.text_frame, [
        [(g, 14, True, INK, False)],
        [(when, 12, False, GREY, False)],
    ], align=PP_ALIGN.CENTER, space_after=Pt(3))
tb = box(s, 0.9, 4.9, 11.5, 1.6)
set_text(tb.text_frame, [
    [('Between sign-offs, the architects work.', 18, False, INK, False)],
    [('The minister’s job in between is to remove obstacles — the political ones, mostly. The technical ones are why the architects were hired.', 18, False, GREY, False)],
], space_after=Pt(8))
footer(s, T)
notes(s, "VO: Notice the rhythm. Four sign-offs in six months. Your minister does not review every diagram. They review at four moments, each tied to a defined deliverable. Between sign-offs, the architects work. The minister's job in between is to remove obstacles — the political ones, mostly. The technical ones are why the architects were hired.")

phase_slide(4, 'A living, governed EA — repository plus practice. The 2–4 permanent architects keep it alive; the EA Board reviews projects against it.',
            'Quarterly Board reviews — indefinitely. New projects reviewed against the architecture; cross-ministry decisions approved; domain boundaries enforced.', 'Ongoing — your country’s permanent way of working',
            "VO: Phase five. Execute and Govern. The question: how do we sustain this? The approved roadmap becomes a project pipeline. The small permanent EA team — two to four senior architects — turns the EA from a one-time delivery into a living repository. The EA Board reviews new projects against the architecture, approves cross-ministry decisions, enforces boundaries between domains. Quarterly reviews. Indefinitely.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'You can skip a phase — and pay for it several times over')
tb = box(s, 0.9, 1.9, 11.5, 3.4)
set_text(tb.text_frame, [
    [('Discover before you measure.', 24, True, INK, False)],
    [('Measure before you adapt.', 24, True, INK, False)],
    [('Adapt before you plan.', 24, True, INK, False)],
    [('Plan before you execute.', 24, True, INK, False)],
    [('Govern always.', 24, True, ITU_BLUE_DARK, False)],
], space_after=Pt(10))
tb = box(s, 0.9, 5.6, 11.5, 0.9)
set_text(tb.text_frame, [[('Skipping usually costs several times what it appeared to save.', 17, False, GREY, True)]])
footer(s, T)
notes(s, "VO: The phases depend on each other, in order. Discover before you measure. Measure before you adapt. Adapt before you plan. Plan before you execute. Govern always. You can skip a phase, but you will pay for it later — usually several times what skipping appeared to save.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Week 1 to week 26 — then forever')
hline(s, 0.9, 3.3, 11.5, MIDGREY, weight_pt=2.5)
miles = [
    ('Week 1', 'Discovery begins', 1.3),
    ('Week 26', 'Approved roadmap — to cabinet', 6.0),
    ('Ongoing', 'Governed EA — permanent practice', 10.6),
]
for lab, desc, x in miles:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(3.12), Inches(0.4), Inches(0.4))
    solid(c, ITU_BLUE)
    tb = box(s, x - 1.0, 2.2, 2.5, 0.6)
    set_text(tb.text_frame, [[(lab, 18, True, INK, False)]], align=PP_ALIGN.CENTER)
    tb = box(s, x - 1.25, 3.8, 3.0, 1.0)
    set_text(tb.text_frame, [[(desc, 14, False, GREY, False)]], align=PP_ALIGN.CENTER)
tb = box(s, 0.9, 5.3, 11.5, 1.3)
set_text(tb.text_frame, [
    [('“Months, not years” is not a slogan. It is the consequence of five phases, each with a clear question, a clear deliverable, and a clear sign-off.', 17, False, INK, False)],
])
footer(s, T)
notes(s, "VO: Six months from Discovery to an approved roadmap. Then your country in permanent EA-governed mode. That is what 'months not years' actually means. It is not a slogan. It is the consequence of sequencing five phases — each with a clear question, a clear deliverable, a clear sign-off — and committing to the practice that runs after.\n\n(This answers the prediction question posed at the start of the video.)")

sources_slide(prs, T, [
    'PAERA v1.0 — §3.1.3 Readiness Assessment',
    'PAERA v1.0 — §5.4 Organisational Assessment & Roadmap',
])

# ================================================================ 1.7
T = '1.7 · What you will need from your minister'
section('1.7', 'What you will need from your minister',
        'Four asks. A small permanent EA team. An EA Board with real authority. About two per cent of digital budget, sustained for five years. And one promise — that the team will not be pulled onto the urgent project of the week.',
        '~5 minutes',
        "VO, slide 1: Suppose you have made the case. Your minister is convinced an EA is the right work. Now the harder part: agreeing the four specific things the minister must commit to. Each one is necessary. Without any one of them, the EA programme will struggle to deliver what it could.")


def ask_slide(n, heading, rows, vo, accent_last=False):
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, heading)
    y = 1.8
    for i, (lab, val) in enumerate(rows):
        tb = box(s, 0.72, y, 2.5, 0.6)
        set_text(tb.text_frame, [[(lab.upper(), 13, True, ITU_BLUE_DARK, False)]])
        tb = box(s, 3.3, y - 0.02, 9.2, 1.15)
        bold = accent_last and i == len(rows) - 1
        set_text(tb.text_frame, [[(val, 17, bold, INK, False)]])
        y += 1.18
    footer(s, T)
    notes(s, vo)
    return s


ask_slide(1, 'Ask 1 — A small permanent EA team', [
    ('Size', '2 to 4 senior architects — permanent, reporting to your CDO or equivalent'),
    ('Typically', 'A chief architect, domain architects, a methodology lead. Smaller countries: 2 people, each carrying several domains.'),
    ('Not', 'Consultants. Not a temporary unit.'),
    ('The point', 'The institutional home of architecture — it exists whether or not any single programme is running.'),
], "VO: First ask. A small permanent EA team. Two to four senior architects. Permanent. Reporting to your CDO or its equivalent. Not project consultants who arrive and leave. Not a temporary unit. The institutional home of architecture work in your country. Tell your minister: this team will exist whether or not any single programme is running. That is the point. It is your country's permanent muscle for cross-cutting digital decisions.", accent_last=True)

ask_slide(2, 'Ask 2 — An EA Board with binding authority', [
    ('Chair', 'Your CDO — or the minister directly'),
    ('Members', 'Sector ministry CIOs, owners of the major state registers, optionally an external advisor'),
    ('Cadence', 'Quarterly main meetings; ad-hoc for urgent decisions'),
    ('Mandate', 'BINDING — not advisory. Without binding authority, the EA becomes documentation nobody reads.'),
], "VO: Second ask. An EA Board with real authority. Chaired by your CDO, or by the minister directly. Members: sector ministry CIOs, owners of the major state registers, and where useful an external advisor. Cadence: quarterly meetings, with ad-hoc sessions for urgent decisions. The hardest part: the mandate must be binding, not advisory. The Board reviews new digital projects against the architecture, approves cross-ministry integrations, and enforces boundaries between architectural domains. With it, the EA becomes the place every digital decision passes through.", accent_last=True)

ask_slide(3, 'Ask 3 — About 2% of digital budget, for five years', [
    ('Initial engagement', 'The first four phases — about 10 to 15 senior person-months over six months'),
    ('Permanent practice', 'Your 2 to 4 architects, ongoing'),
    ('Governance overhead', 'Board time, occasional external review'),
    ('The framing', 'A five-year envelope, not an annual line that disappears when priorities shift. Every other digital programme runs more efficiently when this 2% is in place.'),
], "VO: Third ask. A sustained budget envelope. Three parts. The initial six-month engagement that runs the first four phases — about ten to fifteen senior person-months. The permanent practice — your two to four architects, ongoing. The governance overhead — Board time and occasional external review. As a share of your digital-government budget, the whole thing is typically about two per cent. What you want is a five-year envelope, not an annual budget line that disappears when ministerial priorities shift. Tell the minister: this is the leverage decision.", accent_last=True)

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Ask 4 — One promise: the team stays protected')
p = panel(s, 0.72, 1.8, 11.9, 2.6, ITU_BLUE)
panel_text(s, 0.72, 1.9, 11.9, 2.6, [
    [('The EA team will not be pulled onto the urgent project of the week.', 24, True, WHITE, False)],
    [('Not by you. Not by the minister. Not by anyone in cabinet.', 18, False, WHITE, False)],
], m=0.35)
tb = box(s, 0.72, 4.8, 11.9, 1.7)
set_text(tb.text_frame, [
    [('This is the most commonly broken promise — and the one that quietly kills EA programmes in their second year.', 18, True, INK, False)],
    [('Get it in writing if you can. Recommit it whenever the minister changes.', 16, False, GREY, False)],
], space_after=Pt(8))
footer(s, T)
notes(s, "VO: Fourth ask. One promise. The EA team will not be pulled onto the urgent project of the week. Not by you. Not by the minister. Not by anyone in cabinet. Get this in writing if you can. The most common way EA programmes quietly die is in their second year, when the team is moved onto a flagship delivery and the architecture work stops. The promise to protect the team must be explicit, must be visible, and must be recommitted whenever the minister changes.\n\nThis is the emotional peak of video 1.7 — the one slide in the module with a full-colour block.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Be honest about the time horizon')
rows_slide(s, [
    ('6 months', 'to an approved roadmap'),
    ('18 to 24 months', 'to a fully operating practice'),
    ('5+ years', 'to mature governance'),
], numbered=False, head_size=24, sub_size=17, top=1.6, bottom=5.3)
tb = box(s, 0.72, 5.5, 11.9, 1.3)
set_text(tb.text_frame, [
    [('The minister who launches this work will not be the one who completes it.', 19, True, INK, False)],
    [('A minister who hears that as a problem is sponsoring a deliverable, not an EA. A minister who hears it as a feature is the right minister.', 15, False, GREY, False)],
], space_after=Pt(6))
footer(s, T)
notes(s, "VO: And one note on time horizon to put to your minister honestly. Six months to an approved roadmap. Eighteen to twenty-four months to a fully operating practice. Five years to mature governance. The minister who launches this work will not be the one who completes it. That is uncomfortable, but it is the truth — and a minister who hears it as a problem is sponsoring a deliverable, not an EA. A minister who hears it as a feature is the right minister to commission this work.")

big_slide(prs, 'A team. A Board. A budget envelope. A promise. Put them on one page — and ask for all four together.',
          T,
          "VO: Four asks. A team. A Board. A budget envelope. A promise. Put them on a single page. Bring them to the meeting. Ask for all four together — not in pieces. With all four committed, you have what you need. Without any one of them, the work will struggle.")

sources_slide(prs, T, [
    'PAERA v1.0 — §4.2.1 Management',
    'PAERA v1.0 — §4.2.2 Architecture',
    'PAERA v1.0 — §5.4 Organisational Assessment & Roadmap',
])

# ================================================================ 1.8
T = '1.8 · Four signposts'
section('1.8', 'Four signposts — three African, one international',
        'Rwanda, Kenya and South Africa show the pattern at African scale and in different governance shapes. Estonia is the international polestar. The pattern travels. Your country can apply it too.',
        '~5 minutes',
        "VO, slide 1: Four countries did the work. Three of them are in Africa. They differ in size, in resources, in governance type — and the EA pattern shows in all three. The fourth is Estonia, the most-cited international example, useful as a reference but with a very different starting context. Looking at the four together shows what the pattern looks like in practice.")


def country_slide(heading, rows, vo):
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, heading)
    rows_slide(s, rows, numbered=False, head_size=18, sub_size=15)
    footer(s, T)
    notes(s, vo)
    return s


country_slide('Rwanda — small country, strong centre: the lifecycle can compress', [
    ('Ministry of ICT and Innovation; Irembo as the citizen-services platform',
     'From a small base in the mid-2010s to one of the most ambitious digital-government programmes on the continent.'),
    ('A single national ID linked across services', ''),
    ('A small institutional muscle, coordinated centrally, disciplined', ''),
    ('Lesson: with political will, the lifecycle can be compressed',
     'And the gains in the second and third sectors come faster than the first.'),
], "VO: Rwanda. The Ministry of ICT and Innovation, with Irembo as the unifying citizen-services platform. Starting from a small base in the mid-2010s, Rwanda built one of the most ambitious digital-government programmes on the continent. A single citizen-service platform. A national ID linked across services. Strong central coordination. The institutional muscle is small but disciplined. The lesson for other countries: in a smaller country with political will, the lifecycle can be compressed, and the gains in the second and third sectors come faster than the first.")

country_slide('Kenya — the obstacles are concrete and documented: plan for them', [
    ('Huduma Centres — physical one-stop shops for government services', ''),
    ('Huduma Namba — the unifying identity programme underneath', ''),
    ('Mixed results, openly debated in public',
     'Obstacles met in courts, in parliament, and in implementation.'),
    ('Lesson: the lesson is not “Kenya solved this”',
     'It is that Kenya tried, the obstacles are documented — legal, political, technical — and the debate tells you what to plan for.'),
], "VO: Kenya. The Huduma Centres and the Huduma Namba experience. Kenya took a different path — physical one-stop centres where citizens could access many government services in one place, with the Huduma Namba programme attempting to provide a unifying digital identity to underpin them. The results are mixed and openly debated in Kenya's public arena. The lesson is not 'Kenya solved this' — it is that Kenya tried, encountered concrete obstacles in courts, in parliament, and in implementation, and the debate is documented and useful. For countries thinking about similar moves, Kenya's experience tells you what to plan for.\n\nHonesty about the contested experience is a feature of this example, not a defect.")

country_slide('South Africa — where you cannot impose, you federate', [
    ('SITA — the State Information Technology Agency as coordinating body', ''),
    ('Strong provincial governments and constitutionally autonomous statutory bodies',
     'A single architecture cannot be imposed from the top.'),
    ('Shared standards, common procurement, a maintained reference architecture',
     'Agencies adopt it rather than have it imposed.'),
    ('Lesson: federation — coordination without coercion',
     'The realistic pattern wherever sub-national governments or autonomous bodies carry real authority.'),
], "VO: South Africa. State Information Technology Agency, known as SITA, and the federated digital-government model. South Africa is a federal democracy with strong provincial governments and constitutionally autonomous statutory bodies. You cannot impose a single architecture from the top. SITA acts as a coordinating body — shared standards, common procurement frameworks, a maintained reference architecture that agencies adopt rather than have imposed. The lesson: in any country where sub-national governments or autonomous statutory bodies carry real authority, the federated model — coordination without coercion — is the realistic pattern. Many African countries with strong provincial or county governments will recognise this shape.")

country_slide('Estonia — the polestar, not the template', [
    ('RIA — the Information System Authority; X-Road as the data-exchange backbone', ''),
    ('Distributed registries — population, business, land — owned by their accountable agencies', ''),
    ('Anchoring principle: Once-Only',
     'The state never asks a citizen for the same information twice.'),
    ('Started late 1990s; most public services online today',
     'A small unitary state with very different starting conditions — use it as a polestar for what mature digital government looks like, not a template to copy.'),
], "VO: Estonia. The Information System Authority, known as RIA. The most-cited international example of mature digital government. Starting in the late 1990s, Estonia built X-Road as the data-exchange backbone, distributed state registries owned by their accountable agencies, and the Once-Only principle — that the state never asks a citizen for the same information twice. Today, almost every public service in Estonia runs online. Estonia is a small unitary state with very different starting conditions from most African countries. Use it as a polestar — what fully mature digital government can look like — not as a template to copy directly. The pattern is the same. The path is your country's.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Four very different countries — the same four elements')
tblshape = s.shapes.add_table(5, 5, Inches(0.72), Inches(1.7), Inches(11.9), Inches(4.4))
tbl = tblshape.table
tbl.columns[0].width = Inches(5.1)
for c in range(1, 5):
    tbl.columns[c].width = Inches(1.7)
headers = ['', 'Rwanda', 'Kenya', 'South Africa', 'Estonia']
elements = [
    'Small central team with real authority',
    'Published framework other agencies adopt',
    'Binding governance, not advisory',
    'Multi-year horizon, results visible in months',
]
for c, h in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.text = h
    for pgh in cell.text_frame.paragraphs:
        pgh.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        for r in pgh.runs:
            r.font.name = 'Arial'; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = INK
for i, el in enumerate(elements):
    cell = tbl.cell(i + 1, 0)
    cell.text = el
    for pgh in cell.text_frame.paragraphs:
        for r in pgh.runs:
            r.font.name = 'Arial'; r.font.size = Pt(15); r.font.color.rgb = INK
    for c in range(1, 5):
        cell = tbl.cell(i + 1, c)
        cell.text = '✓'
        for pgh in cell.text_frame.paragraphs:
            pgh.alignment = PP_ALIGN.CENTER
            for r in pgh.runs:
                r.font.name = 'Arial'; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = RGBColor(0x00, 0x9C, 0xD6)
tb = box(s, 0.72, 6.3, 11.9, 0.6)
set_text(tb.text_frame, [[('These are not outliers. This is what committing to the lifecycle, the team, the governance and the time horizon looks like in practice.', 14.5, False, GREY, False)]])
footer(s, T)
notes(s, "VO: Four very different countries. The same architectural elements show in all four. A small central team with real authority. A published framework that other agencies adopt rather than fight. A governance mechanism that is binding, not advisory. A time horizon measured in years for full maturity, with intermediate results visible inside months. These are not outliers. They are what committing to the lifecycle and the team and the governance and the time horizon looks like in practice.")

big_slide(prs, 'The pattern travels — across small countries and large, unitary states and federations. Your country can apply it too.',
          T,
          "VO: The pattern travels. Across small countries and large. Across unitary states and federations. Across countries with strong central authority and countries where authority is distributed. Your country can apply it too — once your minister commits to the four asks from the previous video. The rest of this knowledge product shows you how to do the work, using a fictional country called Progressa so every step is visible in detail.")

sources_slide(prs, T, [
    'Rwanda — Irembo (irembo.gov.rw); Ministry of ICT and Innovation',
    'Kenya — Huduma Kenya (huduma.go.ke); Huduma Namba',
    'South Africa — SITA (sita.co.za)',
    'Estonia — e-Estonia.com; RIA (ria.ee)',
])

# ================================================================ Thank you
s = add_slide(prs, LAYOUT_THANKS)
notes(s, 'Closing slide for the combined deck. Individual videos end on their sources slide instead.')

# Self-check: the split spec's slide ranges depend on this count, and a helper that
# silently stops drawing shows up first as a slide with no voice-over.
assert len(prs.slides._sldIdLst) == 64, 'slide count changed — update decks/split_spec.json'
assert all(sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip() for sl in prs.slides), \
    'every slide carries its voice-over in the notes'

OUT = os.environ.get('OUT_PATH') or os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    'videos', 'module_1', 'en', 'decks', 'KP1_M1_Deck_v0.1.pptx')
prs.save(OUT)
print('slides:', len(prs.slides._sldIdLst))
print('saved', OUT)
