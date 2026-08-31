#!/usr/bin/env python3
# Build the KP1 Module 3 (Topic 3) video deck on the ITU template.
# Content only — every generic helper, branding constant and layout index comes from
# ITU-Giga-KP-Plugin/skills/kp-deck-builder/scripts/deck_lib.py (which also ships the
# template). Conventions and design rules: that skill's SKILL.md.
# Generated .pptx is NEVER hand-edited — fix here, re-render, re-run the split
# (kp-deck-builder/scripts/split_module_deck.py + the split spec next to the decks).
# Override paths with TEMPLATE= and OUT_PATH= env vars.
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                '..', 'ITU-Giga-KP-Plugin', 'skills', 'kp-deck-builder', 'scripts'))
from deck_lib import (
    GREY, INK, ITU_BLUE, ITU_BLUE_DARK, LIGHT, PANEL_GREY, WHITE,
    LAYOUT_THANKS, LAYOUT_WHITE,
    add_slide, big_slide, block_slide, box, delete_template_slides, edit_agenda,
    edit_cover, footer, two_panel,
    mini_strip, notes, open_template, panel, panel_text, rows_slide,
    section_slide, set_text, sources_slide, title)

from pptx.util import Pt

prs = open_template(os.environ.get('TEMPLATE'))

AUDIENCE = 'Chief architect · senior architect · sector ICT lead'


# ---------------------------------------------------------------- module-local composites
def rows_block(head, rows, closing, tag, note, numbered=True, top=1.6, bottom=6.3, head_size=19):
    """Stacked rows with a closing line under them. Module 3 is a module of short lists —
    the repository's contents, the graduation signs, the update triggers, the gate questions —
    so this shape recurs; the closing line is where each list lands."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    rows_slide(s, rows, top=top, bottom=bottom, numbered=numbered, head_size=head_size)
    if closing:
        tb = box(s, 0.72, bottom + 0.12, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def fade_slide(head, idx, fade, counter, tag, note):
    """The four fade-modes of 3.7: same How-it-fades / The-counter bands every time, so the
    pattern itself teaches. Only the progress strip moves."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    mini_strip(s, 4, idx, y=1.32)   # under the headline: at 0.42 it paints over a two-line title
    y = 1.85
    for label, text, fill, ink in (('HOW IT FADES', fade, LIGHT, INK),
                                   ('THE COUNTER', counter, ITU_BLUE, WHITE)):
        panel(s, 0.72, y, 11.9, 2.1, fill)
        tb = box(s, 0.95, y + 0.2, 2.5, 0.4)
        set_text(tb.text_frame, [[(label, 11, True, ITU_BLUE_DARK if ink is INK else WHITE, False)]])
        tb = box(s, 3.5, y + 0.17, 8.85, 1.75)
        set_text(tb.text_frame, [[(text, 18, False, ink, False)]])
        y += 2.45
    footer(s, tag)
    notes(s, note)
    return s


def section(code, name, message, runtime, note):
    return section_slide(prs, 'KP1 · MODULE 3 · VIDEO %s' % code, code, name, message,
                         runtime + ' · standalone video · voice-over on text slides', note)


# ---------------------------------------------------------------- COVER (edit slide 1)
edit_cover(
    prs,
    title_text='Housing, governing and\nsustaining the architecture',
    kicker='KP1 · Government Enterprise Architecture · Module 3',
    blurb='Seven standalone videos for the architect who runs the practice: the repository the '
          'architecture lives in, tooling chosen without lock-in, the update discipline, an EA '
          'Board that can say no, the project review gate, the few metrics that matter, and the '
          'four ways a practice fades.',
    length='~32 mins across 7 videos (3.1 – 3.7)',
    audience=AUDIENCE,
    panel_heading='THE PRACTICE THIS MODULE STANDS UP',
    panel_items=['Repository — where the architecture lives',
                 'Tooling — chosen without lock-in',
                 'Governance — a Board and a gate',
                 'Sustainment — alive past year two'],
    panel_footer='5 gate questions · 4 health metrics · 4 ways a practice fades',
    note_text='Cover for the combined Module 3 deck. Each section that follows is one standalone '
              '~4–5 minute video. Module 2 taught the architect to read, model and assess a '
              'government; Module 3 is the practitioner rendering of Phase 5, Execute and Govern — '
              'the permanent practice that keeps the architecture alive after the assessment is done.')

# ---------------------------------------------------------------- AGENDA (edit slide 2)
edit_agenda(
    prs,
    header='Module 3 — seven videos',
    items=[
        ('3.1  Set up the one place your architecture lives', '~5 min'),
        ('3.2  Choose EA tooling without locking yourself in', '~4 min'),
        ('3.3  Keep the repository true — the update discipline', '~4 min'),
        ('3.4  Stand up an EA Board that can actually say no', '~5 min'),
        ('3.5  Review projects against the architecture', '~5 min'),
        ('3.6  Show the EA is working — the few metrics that matter', '~4 min'),
        ('3.7  Keep the practice alive past year two', '~5 min'),
    ],
    message_paras=[
        'An architecture that is not housed, kept current and given teeth stops being an '
        'architecture within a year. It becomes a document about one.',
        'This module stands up the permanent practice: one place the architecture lives, a Board '
        'with the authority to say no, a gate every project passes through, and the moves that '
        'keep all of it alive past year two.',
    ],
    note_text='Navigation slide for the combined deck; the videos ship standalone on YouTube. '
              '3.1–3.3 build the home of the architecture and the discipline that keeps it true; '
              '3.4–3.5 give it authority through the Board and the review gate; 3.6–3.7 are about '
              'proving it works and keeping it alive.')

delete_template_slides(prs, keep=2)


# ================================================================ 3.1
T = '3.1 · Set up the one place your architecture lives'
section('3.1', 'Set up the one place your architecture lives',
        'An EA repository is the single agreed place the architecture lives — the four layers, the '
        'entities and the decisions — so that one picture of your government exists instead of many '
        'private copies. Set it up first; everything else governs what goes into it.',
        '~5 minutes',
        "VO, slide 1: You have run an assessment. You have the four-layer picture of a sector. Now "
        "the question is: where does it live? If the answer is a slide deck on your laptop, the "
        "architecture will be out of date within a month and disagreed with within two. An "
        "Enterprise Architecture needs a home — one place where the current picture lives, that "
        "everyone works from. That home is the repository.")

rows_block('A repository is a discipline, not a product',
           [('The four layers', 'The capabilities, data domains, applications and technology you mapped.'),
            ('The relationships', 'How they connect, expressed with the shared entities.'),
            ('The decisions', 'What the EA Board ruled — and why it ruled that way.')],
           'Not a building. Not a particular software product. One agreed store — one version everyone trusts.',
           T,
           "VO: A repository is not a building, and it is not a particular software product. It is "
           "the single, agreed store of your architecture. It holds three things. The four layers — "
           "the capabilities, data domains, applications and technology you mapped. The "
           "relationships — how they connect, using the shared entities. And the decisions — what "
           "the EA Board ruled, and why. One place. One version everyone trusts.")

block_slide(prs, 'A second copy does not add a backup — it destroys the single source of truth',
            ['The moment there are two copies, they drift. Then they disagree.',
             'One ministry works from last year’s picture; another from a newer one. People '
             'stop arguing about the architecture and start arguing about whose copy is right.'],
            'One place — or effectively none.',
            T,
            "VO: Why insist on one place? Because the moment there are two copies, they drift, and "
            "then they disagree. One ministry works from last year's picture; another from a newer "
            "one. People stop arguing about the architecture and start arguing about whose copy is "
            "right. A second copy does not add a backup — it destroys the single source of truth. "
            "One place, or effectively none.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'What it holds, concretely — a sector anyone can look up')
cells = [('DATA DOMAINS, WITH OWNERS',
          ['Person → the Identity Authority.',
           'Learner → the Learner Registry.',
           'Examination result → the Examination Authority.']),
         ('THE CAPABILITY MAP AND PORTFOLIO',
          ['What the sector can do, and the system that supports each capability.']),
         ('THE SHARED TECHNOLOGY',
          ['The platforms, standards and infrastructure all of it runs on.']),
         ('THE DECISION LOG',
          ['The Board ruled that the Examination Authority must consume the Learner Registry '
           'rather than keep its own list — with the reason.'])]
cw, ch = 5.9, 2.2
for i, (head, lines) in enumerate(cells):
    x = 0.72 + (i % 2) * (cw + 0.23)
    y = 1.68 + (i // 2) * (ch + 0.22)
    panel(s, x, y, cw, ch, LIGHT)
    panel_text(s, x, y, cw, ch,
               [[(head, 14, True, ITU_BLUE_DARK, False)]]
               + [[(ln, 15.5, False, INK, False)] for ln in lines])
tb = box(s, 0.72, 6.4, 11.9, 0.5)
set_text(tb.text_frame, [[('Anyone can look up not just what the architecture is, but why it is '
                           'that way.', 15.5, True, ITU_BLUE_DARK, False)]])
footer(s, T)
notes(s, "VO: Concretely, take Progressa's education sector. The repository holds the data domains "
         "and their owners — Person owned by the Identity Authority, Learner by the Learner "
         "Registry, examination results by the Examination Authority. It holds the capability map, "
         "the application portfolio, the shared technology. And it holds the decisions: when the "
         "Board ruled that the Examination Authority must consume the Learner Registry rather than "
         "keep its own list, that ruling lives in the repository, with its reason. Anyone can look "
         "up not just what the architecture is, but why it is that way.\n\n"
         "On screen: Progressa is a fictional demonstration country; institution names are "
         "plain-text labels, no emblems.")

block_slide(prs, 'A spreadsheet everyone uses beats a platform nobody does',
            ['You do not need to buy anything to start. A well-structured spreadsheet, or a shared '
             'wiki, is a real repository — if it holds the four layers, the relationships and the '
             'decisions, and if everyone uses it as the one source.',
             'Many countries run a perfectly good early EA on a spreadsheet and a document store. '
             'When the spreadsheet starts to hurt, you graduate — and the choice of tool is its '
             'own deliberate decision.'],
            'Structure matters more than software.',
            T,
            "VO: You do not need to buy anything to start. A well-structured spreadsheet, or a "
            "shared wiki, is a real repository if it holds the four layers, the relationships and "
            "the decisions, and if everyone uses it as the one source. Structure matters more than "
            "software. Many countries run a perfectly good early EA on a spreadsheet and a document "
            "store. The discipline of one place is what counts — not the price of the tool. When a "
            "spreadsheet starts to hurt, you graduate to a dedicated tool; the choice of tool is "
            "its own deliberate decision.")

two_panel(prs, 'The repository is the object both sides point at',
          ('THE MINISTER, BRIEFING CABINET',
           ['Which services the sector delivers.',
            'Where the money and the risk sit.',
            'What the Board decided, and why.']),
          ('THE ARCHITECT, DESIGNING A SYSTEM',
           ['Which capability the system serves.',
            'Which data domain it may use, and who owns it.',
            'What the Board decided, and why.']),
          'Same picture — because there is only one.',
          T,
          "VO: One more reason the repository matters. It is the object both sides point at. When "
          "your minister briefs cabinet and when your architect designs a system, they are looking "
          "at the same picture — because there is only one. The repository is what makes the "
          "architecture a shared language between the business side and the IT side, instead of two "
          "private ones. Without it, each side keeps its own version, and the conversation breaks "
          "down.")

big_slide(prs,
          'An EA repository is the single agreed place your architecture lives — the layers, the '
          'relationships and the decisions — so there is one picture of your government, not many.',
          T,
          "VO: So before tooling, before governance, set up the one place. The single store of the "
          "layers, the relationships and the decisions. Get that right, and everything else in this "
          "module is about protecting what goes into it. Get it wrong — let copies multiply — and "
          "no amount of governance will save the architecture.")

sources_slide(prs, T, [
    'PAERA v1.0 — §4.2.2 (Architecture — the EA function)',
    'PAERA v1.0 — §3.1 (Governance & Policy)',
    'PAERA v1.0 — Annex 2 (Metamodel)',
])


# ================================================================ 3.2
T = '3.2 · Choose EA tooling without locking yourself in'
section('3.2', 'Choose EA tooling without locking yourself in',
        'Choose your EA tooling the way you would choose any system — reuse before buy, buy before '
        'build, keep your data in open formats you control, and never let the EA tool itself become '
        'the vendor trap it is meant to help you avoid.',
        '~4 minutes',
        "VO, slide 1: At some point a spreadsheet stops coping. You have hundreds of entities, "
        "several sectors, relationships you cannot see in rows and columns. It is time for a "
        "dedicated EA tool. And here is the danger: in buying a tool to help your government avoid "
        "vendor lock-in, you can lock yourself into the tool. Here is how to choose without falling "
        "into that trap.")

rows_block('Graduate when the spreadsheet hurts in three specific ways',
           [('You cannot see the relationships any more', 'They are scattered across tabs.'),
            ('Several people edit, and their changes collide', ''),
            ('A simple cross-layer question costs an afternoon',
             'Which systems touch the Person domain?')],
           'Then graduate — not before. A heavy tool bought early is an empty, expensive database.',
           T,
           "VO: First, know when to graduate. The signs: you cannot see the relationships any more, "
           "because they are scattered across tabs. Several people edit and their changes collide. "
           "You cannot answer a simple cross-layer question — which systems touch the Person domain "
           "— without an afternoon of manual work. When the spreadsheet hurts in these specific "
           "ways, graduate. Not before. Buying a heavy tool too early just gives you an empty, "
           "expensive database.")

block_slide(prs, 'The rule you enforce on every ministry applies to your own tool too',
            ['Reuse before buy, buy before build. Is there already a shared tooling platform your '
             'government runs that you can use? Is there an open-source EA tool that fits? Only '
             'then a bought product.'],
            'Building your own EA tool is almost never right.',
            T,
            "VO: Then apply the same sourcing rule you apply to every system — reuse before buy, "
            "buy before build. Is there already a shared tooling platform your government runs that "
            "you can use? Is there an open-source EA tool that fits? Only then a bought product. "
            "Building your own EA tool is almost never right — it is the most bespoke choice for a "
            "problem many others have already solved. The rule you enforce on every ministry's "
            "project applies to your own tool too.")

block_slide(prs, 'One question decides whether you own your architecture, or the tool does',
            ['Before you adopt any tool, ask it: can I export everything, in a format I can read '
             'without this tool, whenever I want?',
             'If the answer is no, the tool owns your architecture, not you.'],
            'The content is the asset; the tool is only a viewer.',
            T,
            "VO: The single most important rule: keep your architecture content in an open format "
            "you control. Before you adopt any tool, ask one question — can I export everything, in "
            "a format I can read without this tool, whenever I want? If the answer is no, the tool "
            "owns your architecture, not you. The content — the layers, the relationships, the "
            "decisions — is the asset. The tool is just a viewer. You must be able to change the "
            "viewer without losing the asset.\n\nThe operative beat of the video. Hold the export "
            "question a beat longer than feels comfortable.")

block_slide(prs, 'Bend the tool to your metamodel, never the metamodel to the tool',
            ['Every EA tool ships with its own built-in way of organising things — its own '
             'metamodel. The mistake is to let the tool’s model quietly replace the one you '
             'adopted.',
             'You bend the tool to fit your entities and your principles — not the other way around.'],
            'If a tool cannot represent your metamodel, question the tool.',
            T,
            "VO: One more trap. Every EA tool ships with its own built-in way of organising things "
            "— its own metamodel. The mistake is to let the tool's model quietly replace the one you "
            "adopted. You bend the tool to fit your entities and your principles — not the other "
            "way around. If a tool cannot represent your metamodel, that is a reason to question "
            "the tool, not to change your metamodel. The architecture is the content and the model. "
            "The tool serves them.")

big_slide(prs,
          'Choose EA tooling like any system — reuse before buy, keep your data in open formats you '
          'control, and bend the tool to your metamodel, never the reverse.',
          T,
          "VO: So choose your tooling the way you would choose any government system. Graduate when "
          "the spreadsheet truly hurts. Reuse before buy, buy before build. Keep your content in "
          "open formats you control, so you can always leave. And bend the tool to your metamodel. "
          "Do that, and the tool helps you. Skip it, and you have bought yourself the very lock-in "
          "an EA exists to prevent.")

sources_slide(prs, T, [
    'PAERA v1.0 — §5.6 (Sourcing Strategy — build / buy / share / sandbox)',
    'PAERA v1.0 — §3.3 (Digital Infrastructure principles — technology neutrality)',
])


# ================================================================ 3.3
T = '3.3 · Keep the repository true — the update discipline'
section('3.3', 'Keep the repository true — the update discipline',
        'A repository is only worth what it is current. Decide who owns it, what event triggers an '
        'update, and how a change is checked — so the architecture tracks reality instead of slowly '
        'becoming a confident work of fiction.',
        '~4 minutes',
        "VO, slide 1: A repository has one enemy, and it is not technical. It is staleness. An "
        "architecture that is six months behind reality is worse than none — because people trust "
        "it, and it lies to them. The work of keeping the repository true is unglamorous and "
        "constant, and it is what separates a living EA from a binder on a shelf.")

block_slide(prs, 'A missing architecture makes people ask questions; a wrong one answers them',
            ['If the repository says the Examination Authority consumes the Learner Registry, and '
             'in reality it quietly built its own copy last quarter, then every decision made on '
             'the repository is built on a falsehood.'],
            'Currency is not a nice-to-have. It is the whole value.',
            T,
            "VO: Start from why this matters. A missing architecture makes people ask questions. A "
            "wrong architecture answers them — wrongly. If the repository says the Examination "
            "Authority consumes the Learner Registry, and in reality it quietly built its own copy "
            "last quarter, then every decision made on the repository is built on a falsehood. "
            "Stale is worse than absent. Currency is not a nice-to-have. It is the whole value.")

block_slide(prs, 'When everyone owns currency, no one does',
            ['One person — usually the chief architect, or a named custodian on the EA team — is '
             'accountable for the repository being current. Not a committee. Not everyone.',
             'That person does not do all the updating. They are answerable for whether the picture '
             'is true.'],
            'One name, one accountability.',
            T,
            "VO: First, name one owner. One person — usually the chief architect or a named "
            "custodian on the EA team — is accountable for the repository being current. Not a "
            "committee. Not everyone. When everyone owns currency, no one does. One name, one "
            "accountability. That person does not do all the updating, but they are answerable for "
            "whether the picture is true.")

rows_block('Tie updates to events, because a yearly review is eleven months stale',
           [('A system goes live', 'Update the application layer.'),
            ('A system is retired', 'Remove it.'),
            ('A new data domain appears', 'Add it, with its owner.'),
            ('The Board makes a decision', 'Record it, with its reason.')],
           'Events drive updates, not the calendar.',
           T,
           "VO: Second, define what triggers an update. Not review it every year — that guarantees "
           "it is eleven months stale. Tie updates to events. A new system goes live: update the "
           "application layer. A system is retired: remove it. A new data domain appears: add it "
           "with its owner. The Board makes a decision: record it. A ministry reorganises: revise "
           "the capabilities. When a real-world change happens, the repository changes with it. "
           "Events drive updates, not the calendar.")

rows_block('A light gate keeps the repository conformant as it grows',
           [('Does the update use the shared entities correctly?', ''),
            ('Does every data domain still have exactly one owner?', ''),
            ('Is every decision logged with its reason?', '')],
           'The check is small. Skipping it fills the repository with orphans and private language.',
           T,
           "VO: Third, check each change. Not a heavy process — a light gate. When something is "
           "added, confirm it uses the shared entities correctly, that every data domain still has "
           "exactly one owner, and that any decision is logged with its reason. This keeps the "
           "repository conformant as it grows. The check is small. Skipping it lets the repository "
           "fill with private-language entries and orphaned items until it is as messy as the "
           "reality it was meant to clarify.")

block_slide(prs, 'The cheapest moment to update the repository is when a project comes to the Board',
            ['The project tells you what it will build, what it will consume, and what data it '
             'touches. That is exactly the information the repository needs.',
             'So tie the two together: a project that passes the gate leaves its architecture '
             'change in the repository as it goes.'],
            'The governance process and the update discipline are one motion, done once.',
            T,
            "VO: The cheapest moment to update the repository is when a project comes to the Board "
            "for review. The project tells you what it will build, what it will consume, what data "
            "it touches. That is exactly the information the repository needs. So tie the two "
            "together: a project that passes the gate leaves its architecture change in the "
            "repository as it goes. The governance process and the update discipline are the same "
            "motion, done once.")

big_slide(prs,
          'A repository is only worth what it is current — one owner, event-driven updates, a light '
          'conformance check — so the architecture tracks reality instead of becoming fiction.',
          T,
          "VO: So keep the repository true. One named owner accountable for currency. Updates "
          "triggered by real events, not a yearly review. A light check that each change stays "
          "conformant. Do this, and the repository stays a trustworthy picture of your government. "
          "Neglect it, and within a year you have a confident, detailed, widely-trusted work of "
          "fiction.")

sources_slide(prs, T, [
    'PAERA v1.0 — §4.2.2 (Architecture)',
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
])


# ================================================================ 3.4
T = '3.4 · Stand up an EA Board that can actually say no'
section('3.4', 'Stand up an EA Board that can actually say no',
        'An EA Board with binding authority — the right chair, the right members, a regular cadence, '
        'and a mandate that lets it say no — is what turns the architecture from a document into the '
        'place every digital decision passes through.',
        '~5 minutes',
        "VO, slide 1: The repository holds the architecture. The Board is what gives it authority. "
        "Without a governance board, the architecture is a document people can ignore. With one — a "
        "real one, that can say no — the architecture becomes the place every significant digital "
        "decision passes through. The difference between those two outcomes is whether the Board "
        "has binding authority. Everything else about the Board is detail; that is the point.")

block_slide(prs, 'An advisory board produces minutes; a binding board produces decisions',
            ['An advisory board looks at a project, offers an opinion, and the project does '
             'whatever it was going to do. That is governance theatre.',
             'A binding board’s rulings are a condition of proceeding: a project told to '
             'consume the shared identity platform does so, or it does not get funded.'],
            'Without binding authority you have a discussion group, not a governance board.',
            T,
            "VO: Start with the one thing that matters. The Board must be binding, not advisory. An "
            "advisory board looks at a project, offers an opinion, and the project does whatever it "
            "was going to do. That is governance theatre — it produces minutes, not decisions. A "
            "binding board's rulings are a condition of proceeding: a project told to consume the "
            "shared identity platform does so, or it does not get funded. If you cannot get binding "
            "authority, you do not yet have a governance board. You have a discussion group.")

block_slide(prs, 'The chair decides whether the no sticks',
            ['Your Chief Digitalisation Officer or its equivalent — or, in a smaller government, '
             'the minister directly. Senior enough that when the Board declines a powerful '
             'ministry’s pet project, the decision holds.',
             'Note who does not chair it: you. You prepare the Board, you advise it, you bring the '
             'analysis.'],
            'The architect informs; the chair decides.',
            T,
            "VO: Who chairs it decides whether no sticks. The chair is your Chief Digitalisation "
            "Officer or its equivalent — or, in a smaller government, the minister directly. Senior "
            "enough that when the Board declines a powerful ministry's pet project, the decision "
            "holds. Note who does not chair it: you, the architect. You prepare the Board, you "
            "advise it, you bring the analysis. The authority to decide sits with someone who can "
            "carry it politically. The architect informs; the chair decides.")

rows_block('Membership follows the affected systems, not convenience',
           [('The sector ministry CIOs', 'Their systems carry the decisions.'),
            ('The owners of the major state registries', 'The identity authority, the population register.'),
            ('The data-protection regulator', 'Wherever data-sharing is in play.'),
            ('One external advisor', 'Where an outside perspective is useful.')],
           'A Board of the right twelve is better than a Board of the convenient forty.',
           T,
           "VO: Who sits on it: the people whose systems the decisions affect. The sector ministry "
           "CIOs. The owners of the major state registries — the identity authority, the population "
           "register. Where data-sharing is in play, the data-protection regulator. And, where "
           "useful, one external advisor for perspective. Keep it small enough to decide and broad "
           "enough that the decisions are owned by the people who must live with them. A Board of "
           "the right twelve is better than a Board of the convenient forty.")

two_panel(prs, 'Govern on a cadence that does not become the bottleneck',
          ('QUARTERLY — THE MAIN MEETING',
           ['Review the architecture as a whole.',
            'What changed, what is drifting, what to prioritise next.']),
          ('AD HOC — THE FAST PATH',
           ['Urgent project decisions.',
            'So a programme on a deadline is not stuck for three months waiting for the next '
            'quarterly.']),
          'If the Board is slow, projects route around it — and you are back to no governance.',
          T,
          "VO: Set a cadence that governs without becoming the bottleneck. Quarterly main meetings "
          "to review the architecture as a whole — what changed, what is drifting, what to "
          "prioritise. And a fast ad-hoc path for urgent project decisions, so a programme on a "
          "deadline is not stuck for three months waiting for the next quarterly. If the Board is "
          "slow, projects will route around it, and you are back to no governance. Binding "
          "authority plus responsiveness is what keeps projects coming through the front door.")

rows_block('The Board is the standing rhythm no single project can provide',
           [('The repository gives the shared picture', ''),
            ('The metamodel gives the shared words', ''),
            ('The Board gives the shared rhythm',
             'A standing forum where the business side and the IT side decide together.')],
           'Four times a year, and whenever it is urgent, both sides sit down and decide together.',
           T,
           "VO: There is a deeper reason the Board matters. Redesigning how government serves "
           "citizens needs the business side and the IT side to decide together — and they need a "
           "regular place to do it. The Board is that place. The repository gives them the shared "
           "picture; the metamodel gives them shared words; the Board gives them the shared rhythm "
           "— a standing forum where, four times a year and whenever it is urgent, business and IT "
           "sit down and decide together. That standing rhythm is part of what an EA provides that "
           "no single project ever can.")

big_slide(prs,
          'An EA Board with binding authority, the right chair and members, and a cadence that '
          'doesn’t block delivery is what turns the architecture into the place every digital '
          'decision passes through.',
          T,
          "VO: So stand up a Board that can actually say no. Binding, not advisory. Chaired by "
          "someone senior enough that the no sticks. Made of the people whose systems are affected. "
          "Meeting on a rhythm that governs without blocking. Get that in place and the architecture "
          "has teeth. Leave it advisory, and you have a beautiful repository that everyone is free "
          "to ignore.")

sources_slide(prs, T, [
    'PAERA v1.0 — §4.2.1 (Management)',
    'PAERA v1.0 — §3.1 (Governance & Policy)',
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
])


# ================================================================ 3.5
T = '3.5 · Review projects against the architecture'
section('3.5', 'Review projects against the architecture',
        'The architecture review gate — a short, consistent set of questions every project passes '
        'through before funding — is what turns principles and re-use from good intentions into the '
        'actual path of least resistance.',
        '~5 minutes',
        "VO, slide 1: A Board with authority needs something to do with it. That something is the "
        "architecture review gate: the point where every significant new project passes through a "
        "short, consistent set of questions before it gets funded. This is where the architecture "
        "does its real work — not on a wall, but at the moment a project would otherwise quietly "
        "build its own version of something the country already has.")

rows_block('Five questions, asked of every project, the same way every time',
           [('Does a shared building block already exist for this?', ''),
            ('What data domains do you touch — and do you consume the owner’s copy?', ''),
            ('Do you meet the architecture principles — security by design, once-only, the rest?', ''),
            ('Is your sourcing choice deliberate — build, buy, share, sandbox?', ''),
            ('Can you export your data in an open format?', '')],
           'Consistency is what makes the gate fair and predictable — so projects prepare for it '
           'rather than resent it.',
           T,
           "VO: The gate is a few questions, asked of every project, the same way every time. Does "
           "a shared building block already exist for what you want to build? What data domains "
           "will you touch, and do you consume the owning body's copy rather than make your own? Do "
           "you meet the architecture principles — security by design, once-only, the rest? What is "
           "your sourcing choice, and is it deliberate? Can you export your data in an open format? "
           "Five questions. Asked consistently. That consistency is what makes the gate fair and "
           "predictable, so projects prepare for it rather than resent it.\n\n"
           "Centrepiece of 3.5 — reveal one question at a time, then hold the full list on screen.",
           top=1.55, bottom=6.2, head_size=18)

block_slide(prs, 'Building is rational for the project and ruinous for the country',
            ['Left alone, a project builds its own identity function, its own learner list, its own '
             'payment integration — because, inside the project, building is faster than reusing.',
             'The gate is the one place where the whole-of-government view meets that decision, '
             'while it can still be changed. The Board says: consume the shared one.'],
            'That single moment, repeated across every project, is how re-use actually happens.',
            T,
            "VO: Here is why this gate matters more than any document. Left alone, a project builds "
            "its own identity function, its own learner list, its own payment integration — because, "
            "inside the project, building is faster than reusing. That choice is rational for the "
            "project and ruinous for the country. The gate is the one place where the "
            "whole-of-government view meets that decision, while it can still be changed. When a "
            "project proposes to build what already exists as a shared block, the Board, at the "
            "gate, says: consume the shared one. That single moment, repeated across every project, "
            "is how re-use actually happens. Not because a strategy required it — because the gate "
            "enforced it.\n\nThe emotional peak of the module. This is the pivotal slide of 3.5 — "
            "let the closing line land before moving on.",
            punch_fill=ITU_BLUE, punch_ink=WHITE)

block_slide(prs, 'An exception without an expiry quietly becomes the permanent normal',
            ['Sometimes the project is right — the shared identity platform genuinely cannot do '
             'what it needs yet. So the gate is not a wall. It grants exceptions.',
             'But every exception is written down, with its reason, and given a sunset — a date '
             'when it is revisited. Ten open-ended exceptions are the fragmentation you were trying '
             'to prevent.'],
            'Grant exceptions honestly, record them, make them temporary by default.',
            T,
            "VO: Sometimes the project is right. The shared identity platform genuinely cannot do "
            "what this project needs yet. So the gate is not a wall — it grants exceptions. But "
            "every exception is written down, with its reason, and given a sunset — a date when it "
            "is revisited. An exception without an expiry quietly becomes the permanent normal, and "
            "ten of them become the fragmentation you were trying to prevent. Grant exceptions "
            "honestly, record them, and make them temporary by default.")

block_slide(prs, 'The decision log is the architecture’s memory',
            ['Every gate decision and every exception goes into the repository with its reason.',
             'Two years on, when someone asks why the Examination Authority is not allowed its own '
             'learner list, the answer is in the log — not lost with the architect who has moved on.'],
            'It is also what keeps the gate consistent: the precedent is written down.',
            T,
            "VO: Every decision the gate makes goes into the repository's decision log — what was "
            "decided, and why. This is the architecture's memory. Two years on, when someone asks "
            "why the Examination Authority is not allowed its own learner list, the answer is in "
            "the log, with the reasoning, not lost with the architect who has moved on. The decision "
            "log is also what makes the gate consistent: this project is treated the way the last "
            "similar one was, because the precedent is written down.")

big_slide(prs,
          'The review gate — a few consistent questions every project answers before funding, with '
          'logged decisions and time-boxed exceptions — is what turns re-use from a wish into the '
          'default.',
          T,
          "VO: So build the gate. A few questions, asked of every project the same way. A clear "
          "ruling — consume the shared block, or a written, time-boxed exception. Every decision "
          "logged with its reason. This is where principles stop being words and become the path "
          "projects actually take. The gate is the engine of the whole EA. Everything else — the "
          "repository, the Board, the principles — exists so that this moment, repeated, goes the "
          "right way.")

sources_slide(prs, T, [
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
    'PAERA v1.0 — §5.2 (Principles)',
    'PAERA v1.0 — §5.6 (Sourcing Strategy)',
    'GovStack Building Block specifications — for the shared building blocks referenced',
])


# ================================================================ 3.6
T = '3.6 · Show the EA is working — the few metrics that matter'
section('3.6', 'Show the EA is working — the few metrics that matter',
        'A handful of honest metrics — coverage, re-use rate, open exceptions, decisions made — show '
        'the minister and the team that the EA is working, and tell you where it isn’t, without '
        'drowning anyone in vanity numbers.',
        '~4 minutes',
        "VO, slide 1: Sooner or later your minister asks the fair question: is this EA work actually "
        "doing anything? You need an answer that is honest, short, and true. Not a fifty-page "
        "report. A handful of metrics that show whether the architecture is working — and, just as "
        "usefully, where it is not. Pick them carefully, because the wrong metrics make a stalled "
        "EA look healthy.\n\nRETRIEVAL MOMENT: before the four-metrics slide, ask the viewer which "
        "single number they would put in front of the budget authority. The answer — the re-use "
        "rate — is delivered two slides on.")

rows_block('Four numbers carry most of the signal',
           [('Coverage', 'How much of your government is actually in the repository — and current.'),
            ('Re-use rate', 'Of the projects that came through the gate, how many consumed a shared '
                            'block instead of building their own.'),
            ('Open exceptions', 'How many waivers are outstanding, and how old.'),
            ('Decisions', 'How many gate decisions the Board made — and how fast.')],
           'A growing pile of aged exceptions is fragmentation returning; a slow gate is one '
           'projects route around.',
           T,
           "VO: Four metrics carry most of the signal. Coverage — how much of your government is "
           "actually in the repository, and current, not how many pages exist. Re-use rate — of the "
           "projects that came through the gate, how many consumed a shared building block instead "
           "of building their own. This is the one that shows the EA is paying for itself. Open "
           "exceptions — how many waivers are outstanding, and how old; a growing pile of aged "
           "exceptions is fragmentation returning. And decisions — how many gate decisions the "
           "Board made, and how fast, because a slow gate is one projects route around.")

block_slide(prs, 'The re-use rate is the one to put in front of the budget authority',
            ['Every project that consumed the shared identity platform instead of building its own '
             'is a system the country paid for once instead of many times.',
             'Tracked over a year, a rising re-use rate is the closest thing you have to proof that '
             'the EA is returning its cost — the abstract argument turned into a number that goes up.'],
            'That is the metric that protects your funding.',
            T,
            "VO: Of the four, the re-use rate is the one to put in front of the budget authority. "
            "Every project that consumed the shared identity platform instead of building its own "
            "is a system the country paid for once instead of many times. Tracked over a year, a "
            "rising re-use rate is the closest thing you have to proof that the EA is returning its "
            "cost. It turns the abstract argument — planning enables re-use — into a number that "
            "goes up. That is the metric that protects your funding.\n\nThis answers the retrieval "
            "prompt set on the section slide.")

two_panel(prs, 'If a metric would still rise while the architecture stopped mattering, it is wrong',
          ('MEASURES EFFECT', ['Re-use rate.', 'Coverage that is current.',
                               'Gate decisions made, and how fast.']),
          ('MEASURES ACTIVITY', ['Diagrams drawn.', 'Pages written.', 'Meetings held.']),
          'A programme can produce a thousand pages and change nothing.',
          T,
          "VO: Avoid the vanity metrics. The number of diagrams drawn, pages written, meetings held "
          "— these measure activity, not effect. A programme can produce a thousand pages and "
          "change nothing. Worse, busy-looking metrics let a stalled EA hide. If a metric would "
          "still go up while the architecture quietly stopped mattering, it is the wrong metric. "
          "Measure what the EA changes — re-use, coverage, decisions — not how busy the team looks.")

block_slide(prs, 'A scorecard that is always green is believed once',
            ['One page, every quarter, to the Board and the minister — including what is not '
             'working: the coverage gap in a sector nobody will let you near, the aging exceptions '
             'a powerful ministry will not close.',
             'The temptation is to show only green. Resist it.'],
            'A scorecard honestly amber where it should be is what keeps the EA funded across a '
            'change of minister.',
            T,
            "VO: Report the four on a single page, every quarter, to the Board and the minister. "
            "Include what is not working — the coverage gap in a sector nobody will let you near, "
            "the aging exceptions a powerful ministry will not close. The temptation is to show only "
            "green. Resist it. A scorecard that is honestly amber where it should be is what builds "
            "the trust that keeps the EA funded across a change of minister. A scorecard that is "
            "always green is one nobody believes the second time.")

big_slide(prs,
          'A one-page quarterly scorecard — coverage, re-use rate, open exceptions, decisions — '
          'shows honestly whether the EA is working, and the re-use rate is what proves it pays.',
          T,
          "VO: So measure a few things, honestly. Coverage. Re-use rate. Open exceptions. Decisions "
          "made and how fast. Put them on one page, every quarter, amber where they should be "
          "amber. This is how you show the minister the EA is working, how you catch it when it is "
          "not, and how you protect the funding that keeps the practice alive.")

sources_slide(prs, T, [
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
    'PAERA v1.0 — §5.7 (Recommended Roadmap — intermediate-results pattern)',
])


# ================================================================ 3.7
T = '3.7 · Keep the practice alive past year two'
section('3.7', 'Keep the practice alive past year two',
        'EA programmes rarely fail technically; they fade — the team gets pulled away, the repository '
        'goes stale, the Board drifts to advisory, the sponsor changes. Naming these four fade-modes '
        'and the move that counters each is how you keep the practice alive.',
        '~5 minutes',
        "VO, slide 1: Most EA programmes do not fail in a dramatic way. They fade. The first six "
        "months go well — there is energy, a roadmap, a Board. Then, somewhere in the second year, "
        "the practice quietly stops mattering, and one day someone notices the repository is a year "
        "out of date and the Board has not met in two quarters. The fade is predictable. It comes "
        "in four forms. Name them, and you can counter each.\n\nRETRIEVAL MOMENT: before the first "
        "fade slide, ask the viewer to name the ways they have seen a programme fade. The four "
        "answers are delivered on the next four slides, one each.")

fade_slide('Fade one — your architects are the obvious people to second to the crisis', 0,
           'A flagship programme runs into trouble and needs skilled people. Your architects go '
           '“just for a few months”, and the architecture work stops. This is the most '
           'common way EA programmes die.',
           'A written protection promise — an explicit commitment that the EA team will not be '
           'pulled onto the urgent project of the week — recommitted whenever the minister or the '
           'sponsor changes.',
           T,
           "VO: The first fade: the team gets pulled away. A flagship programme runs into trouble "
           "and needs skilled people. Your architects are the obvious choice. They are seconded just "
           "for a few months, and the architecture work stops. This is the most common way EA "
           "programmes die. The counter is a protection promise — an explicit, written commitment "
           "that the EA team will not be pulled onto the urgent project of the week — and it must be "
           "recommitted whenever the minister or the sponsor changes, because the new one never "
           "feels bound by the old one's promise.")

fade_slide('Fade two — an unused repository is a dead one', 1,
           'Updates slip during a busy quarter. The picture drifts from reality. Someone relies on '
           'it, gets burned, and word spreads that the repository cannot be trusted. Once trust '
           'goes, people stop using it.',
           'The update discipline: one named owner accountable for currency, updates tied to real '
           'events and to the review gate — so the repository stays true as a by-product of work '
           'that happens anyway.',
           T,
           "VO: The second fade: the repository goes stale. Updates slip during a busy quarter. The "
           "picture drifts from reality. Someone relies on it, gets burned, and word spreads that "
           "the repository cannot be trusted. Once trust goes, people stop using it, and an unused "
           "repository is a dead one. The counter is the update discipline — one named owner "
           "accountable for currency, updates tied to real events and to the review gate, so the "
           "repository stays true as a by-product of work that happens anyway.")

fade_slide('Fade three — a Board overruled in silence is a Board being dismantled', 2,
           'Under delivery pressure, a powerful project is let through despite the Board’s '
           'objection. Then another. The no becomes a suggestion, projects learn they can route '
           'around it, and within a year it is a discussion group again.',
           'Protect the binding authority deliberately — and track and report every time the Board '
           'is overruled. Sunlight on overrides is what keeps the authority real.',
           T,
           "VO: The third fade: the Board drifts back to advisory. Under delivery pressure, a "
           "powerful project is let through despite the Board's objection. Then another. The Board's "
           "no quietly becomes a suggestion, projects learn they can route around it, and within a "
           "year it is a discussion group again. The counter is to protect the binding authority "
           "deliberately — and to track and report every time the Board is overruled. Sunlight on "
           "overrides is what keeps the authority real; a Board overruled in silence is a Board "
           "being dismantled.")

fade_slide('Fade four — an EA that depends on one champion dies with that champion', 3,
           'The minister or the digitalisation officer who championed the EA moves on. The successor '
           'has their own priorities and no attachment to this one.',
           'Institutionalise it: a legal mandate for the Board, not just a memo. A budget line, not '
           'an annual favour. And the one-page scorecard, so a new sponsor can see in five minutes '
           'that the EA is returning its cost.',
           T,
           "VO: The fourth fade: the sponsor changes. The minister or the digitalisation officer "
           "who championed the EA moves on. The successor has their own priorities and no attachment "
           "to this one. The counter is to make the EA outlive its sponsor — institutionalise it. A "
           "legal mandate for the Board, not just a memo. A budget line, not an annual favour. And "
           "the one-page scorecard, so a new sponsor can see in five minutes that the EA is "
           "returning its cost. An EA that depends on one champion dies with that champion's tenure. "
           "One that is institutionalised survives the handover.")

block_slide(prs, 'You are building something designed to outlast you',
            ['A mature EA practice takes about five years. The architect who stands it up is rarely '
             'the one who sees it mature. That is not a reason for discouragement — it is the job.',
             'A repository that stays true without you. A Board that holds without you. A mandate '
             'that survives the next election.'],
            'Build it to survive you, and you have done the work.',
            T,
            "VO: One honest thing to carry. A mature EA practice takes about five years. The "
            "architect who stands it up is rarely the one who sees it mature. That is not a reason "
            "for discouragement — it is the job. You are building something designed to outlast "
            "you: a repository that stays true without you, a Board that holds without you, a "
            "mandate that survives the next election. Build it to survive you, and you have done "
            "the work. Build it to depend on you, and it fades the moment you leave.")

big_slide(prs,
          'EA programmes fade in four ways — team pulled, repository stale, Board gone advisory, '
          'sponsor changed — and you counter each by institutionalising the practice so it survives '
          'you.',
          T,
          "VO: So watch for the four fades. The team pulled away — counter with a protected, "
          "recommitted promise. The repository stale — counter with the update discipline. The "
          "Board gone advisory — counter by protecting and reporting its authority. The sponsor "
          "changed — counter by institutionalising the mandate, the budget and the scorecard. Keep "
          "the practice alive past year two, and the architecture becomes how your government "
          "works. Let it fade, and it becomes another binder on a shelf that once cost a great "
          "deal.")

sources_slide(prs, T, [
    'PAERA v1.0 — §4.2.1 (Management)',
    'PAERA v1.0 — §4.2.2 (Architecture)',
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
])


# ================================================================ Thank you
s = add_slide(prs, LAYOUT_THANKS)
notes(s, 'Closing slide for the combined deck. Individual videos end on their sources slide instead.')

# Self-check: the split spec's slide ranges depend on this count, and a helper that
# silently stops drawing shows up first as a slide with no voice-over.
assert len(prs.slides._sldIdLst) == 56, 'slide count changed — update decks/split_spec.json'
assert all(sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip() for sl in prs.slides), \
    'every slide carries its voice-over in the notes'

OUT = os.environ.get('OUT_PATH') or os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    'videos', 'module_3', 'en', 'decks', 'KP1_M3_Deck_v0.1.pptx')
prs.save(OUT)
print('slides:', len(prs.slides._sldIdLst))
print('saved', OUT)
