#!/usr/bin/env python3
# Build the KP1 Module 2 (Topic 2) video deck on the ITU template.
# Content only — every generic helper, branding constant and layout index comes from
# ITU-Giga-KP-Plugin/skills/kp-deck-builder/scripts/deck_lib.py (which also ships the
# template). Conventions and design rules: that skill's SKILL.md.
# Generated .pptx is NEVER hand-edited — fix here, re-render, re-run the split
# (kp-deck-builder/scripts/split_module_deck.py + the split spec next to the decks).
# Override paths with TEMPLATE= and OUT_PATH= env vars.
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                '..', 'ITU-Giga-KP-Plugin', 'skills', 'kp-deck-builder', 'scripts'))
from deck_lib import (
    GREY, INK, ITU_BLUE, ITU_BLUE_DARK, LIGHT, MIDGREY, PANEL_GREY, SEPARATOR, WHITE,
    LAYOUT_THANKS, LAYOUT_WHITE,
    add_slide, big_slide, box, delete_template_slides, edit_agenda, edit_cover, footer,
    hline, mini_strip, notes, num_chip, open_template, panel, panel_text, rows_slide,
    section_slide, set_text, solid, sources_slide, title)

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = open_template(os.environ.get('TEMPLATE'))

AUDIENCE = 'Chief architect · senior architect · sector ICT lead'


# ---------------------------------------------------------------- module-local composites
def qdm_slide(head, layer_idx, question, deliverable, mistake, tag, note):
    """The four-layer slides of 2.1: same Question / Deliverable / Mistake bands every time.
    The repetition is the teaching device (bundle §3.1 slide spec), so the layout is fixed
    and only the progress strip moves."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    mini_strip(s, 4, layer_idx, y=1.32)   # under the headline: at 0.42 it paints over a two-line title
    bands = [('QUESTION', question, LIGHT, INK),
             ('DELIVERABLE', deliverable, LIGHT, INK),
             ('MISTAKE', mistake, ITU_BLUE, WHITE)]
    y = 1.75
    for label, text, fill, ink in bands:
        panel(s, 0.72, y, 11.9, 1.55, fill)
        tb = box(s, 0.95, y + 0.17, 2.5, 0.4)
        set_text(tb.text_frame, [[(label, 11, True, ITU_BLUE_DARK if ink is INK else WHITE, False)]])
        tb = box(s, 3.5, y + 0.14, 8.85, 1.25)
        set_text(tb.text_frame, [[(text, 18, False, ink, False)]])
        y += 1.72
    footer(s, tag)
    notes(s, note)
    return s


def trace_stack(head, levels, closing, tag, note, label_w=2.5):
    """Vertical trace from a citizen service down to the technology it runs on — the
    centrepiece visual of 2.1 and 2.5. Text boxes and diamonds only, per the ITU guide."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    n = len(levels)
    h, gap = 0.72, 0.24
    y = 1.62
    for i, (label, text) in enumerate(levels):
        last = i == n - 1
        panel(s, 1.5, y, 10.4, h, ITU_BLUE_DARK if last else LIGHT)
        tb = box(s, 1.72, y + 0.10, label_w, h - 0.2)
        set_text(tb.text_frame, [[(label, 13, True, WHITE if last else ITU_BLUE_DARK, False)]])
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tb = box(s, 1.72 + label_w, y + 0.10, 10.0 - label_w, h - 0.2)
        set_text(tb.text_frame, [[(text, 15, False, WHITE if last else INK, False)]])
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        if not last:
            d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(6.6), Inches(y + h + 0.055),
                                   Inches(0.17), Inches(0.17))
            solid(d, MIDGREY)
        y += h + gap
    tb = box(s, 0.72, 6.42, 11.9, 0.5)
    set_text(tb.text_frame, [[(closing, 15, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def block_slide(head, lead, punch, tag, note, punch_fill=LIGHT, punch_ink=ITU_BLUE_DARK,
                lead_size=20):
    """A paragraph of argument, then the line it lands on. Used where the bundle specifies
    a single text block."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    tb = box(s, 0.72, 1.75, 11.9, 2.9)
    set_text(tb.text_frame, [[(p, lead_size, False, INK, False)] for p in lead], space_after=Pt(12))
    panel(s, 0.72, 4.85, 11.9, 1.45, punch_fill)
    panel_text(s, 0.72, 4.9, 11.9, 1.45, [[(punch, 23, True, punch_ink, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def two_panel(head, left, right, closing, tag, note, right_fill=PANEL_GREY):
    """Comparison columns: (heading, [lines]) each side."""
    s = add_slide(prs, LAYOUT_WHITE)
    title(s, head)
    cw, ch, cy = 5.9, 4.35, 1.7
    for i, ((ph, lines), fill) in enumerate(zip((left, right), (LIGHT, right_fill))):
        x = 0.72 + i * (cw + 0.23)
        panel(s, x, cy, cw, ch, fill)
        paras = [[(ph, 19, True, ITU_BLUE_DARK, False)]]
        paras += [[(ln, 16, False, INK, False)] for ln in lines]
        panel_text(s, x, cy, cw, ch, paras)
    if closing:
        tb = box(s, 0.72, 6.25, 11.9, 0.6)
        set_text(tb.text_frame, [[(closing, 15.5, True, ITU_BLUE_DARK, False)]])
    footer(s, tag)
    notes(s, note)
    return s


def section(code, name, message, runtime, note):
    return section_slide(prs, 'KP1 · MODULE 2 · VIDEO %s' % code, code, name, message,
                         runtime + ' · standalone video · voice-over on text slides', note)


# ---------------------------------------------------------------- COVER (edit slide 1)
edit_cover(
    prs,
    title_text='Reading, modelling and\nassessing a government',
    kicker='KP1 · Government Enterprise Architecture · Module 2',
    blurb='Seven standalone videos for the architect who does the work: the four layers, the '
          'shared metamodel, ten principles to adopt, the taxonomy of public bodies, a worked '
          'sector walkthrough, and how to run a Phase 2 Assess.',
    length='~32 mins across 7 videos (2.1 – 2.7)',
    audience=AUDIENCE,
    panel_heading='THE FOUR LAYERS THIS MODULE TEACHES',
    panel_items=['Business — what the body does',
                 'Data — what it holds, who owns it',
                 'Application — what supports the work',
                 'Technology — what it all runs on'],
    panel_footer='6 shared entities · 10 principles to adopt · 2 traps to catch at Assess',
    note_text='Cover for the combined Module 2 deck. Each section that follows is one standalone '
              '~4–5 minute video. Module 1 addressed the Strategist who commissions a national EA; '
              'Module 2 addresses the architect who does the work, and deliberately introduces the '
              'vocabulary Module 1 set aside — metamodel, taxonomy, capability, data domain.')

# ---------------------------------------------------------------- AGENDA (edit slide 2)
edit_agenda(
    prs,
    header='Module 2 — seven videos',
    items=[
        ('2.1  Read any government in four layers', '~5 min'),
        ('2.2  The shared vocabulary that makes re-use possible', '~5 min'),
        ("2.3  Adopt your principles, don't draft them", '~4 min'),
        ('2.4  Classify any public body before you model it', '~4 min'),
        ('2.5  BDAT on a real ministry — the Progressa walkthrough', '~5 min'),
        ('2.6  Run a Phase 2 Assess', '~5 min'),
        ('2.7  The two traps to catch at Assess', '~4 min'),
    ],
    message_paras=[
        'An architecture only connects to other architectures if everyone reads a government the '
        'same way and draws it with the same entities.',
        'This module teaches that shared reading — then puts it to work on one sector and on the '
        'assessment your roadmap will be built from.',
    ],
    note_text='Navigation slide for the combined deck; the videos ship standalone on YouTube. '
              '2.1–2.4 build the shared apparatus (layers, metamodel, principles, taxonomy), '
              '2.5 applies all four to one sector, 2.6–2.7 turn the apparatus into an assessment '
              'and the traps it must catch.')

delete_template_slides(prs, keep=2)


# ================================================================ 2.1
T = '2.1 · Read any government in four layers'
section('2.1', 'Read any government in four layers',
        'Every government, in any sector, can be read in four layers — Business, Data, Application, '
        'Technology. Learn the question each layer answers, the deliverable it produces, and the '
        'mistake first-time architects make, and you can decompose any ministry put in front of you.',
        '~5 minutes',
        "VO, slide 1: An Enterprise Architecture describes a government in four layers. Business, "
        "Data, Application, Technology. As the architect, you do not just name them — you work "
        "inside them every day. So for each layer you need three things: the question it answers, "
        "the deliverable you produce, and the mistake that catches first-time architects.\n\n"
        "RETRIEVAL MOMENT: before the Business slide, ask the viewer to name the mistake they think "
        "catches first-time architects on each layer. The four answers are delivered on the next "
        "four slides, one per layer.")

qdm_slide('The Business layer describes what a body does — not how it is arranged', 0,
          'What does this body do, for whom, and how well?',
          'A capability map and a service catalogue.',
          'Copying the org chart instead of describing capabilities.',
          T,
          "VO: Start with the Business layer. The question it answers: what does this body do, for "
          "whom, and how well? Not how it is organised inside — what it actually does. The "
          "deliverable is a capability map and a service catalogue. A capability is something the "
          "body can do — register a learner, run an examination, certify a teacher. A service is how "
          "a citizen or another body receives that capability. The mistake first-time architects "
          "make is to copy the organisation chart and call it the Business layer. The org chart "
          "tells you who reports to whom. It does not tell you what the body does. Two bodies with "
          "the same chart can do completely different work. Describe the capabilities, not the boxes.")

qdm_slide('The Data layer outlasts every application above it', 1,
          'What information does the body hold, who owns it, where is the authoritative copy?',
          'A data-domain catalogue with one owner per domain.',
          'Listing databases instead of data domains.',
          T,
          "VO: Second, the Data layer. The question: what information does the body hold, who owns "
          "each kind, and where does the authoritative copy live? The deliverable is a data-domain "
          "catalogue. A data domain is a kind of information the government agrees on — a person, a "
          "learner, a school, a payment. For each domain you name the one body that owns it and "
          "holds the authoritative copy. Data is the longest-lived layer in any government. "
          "Applications come and go. Technology cycles every decade. The data outlasts them all. The "
          "mistake first-time architects make is to list the databases they can see, one per system. "
          "But five systems may each hold their own version of the same learner. The Data layer is "
          "not a list of databases. It is the agreed set of domains, each with one owner, sitting "
          "above all the databases.\n\nHold the one-owner-per-domain line a beat longer — it is the "
          "payload of this slide and the thing 2.5 makes concrete.")

qdm_slide('An application portfolio tells you why a system exists; an inventory does not', 2,
          'What software supports the work, and which capability does each one serve?',
          'An application portfolio mapped to capabilities and data domains.',
          'A software inventory with no link to the business.',
          T,
          "VO: Third, the Application layer. The question: what software supports the work, and "
          "which capability does each application serve? The deliverable is an application portfolio "
          "— every system, mapped to the capability it supports and the data domain it touches. The "
          "mistake here is producing a software inventory that lists names and versions but never "
          "connects them to what the body does. An inventory tells you that a system exists. A "
          "portfolio tells you why it exists, what would break if it failed, and whether two systems "
          "are quietly doing the same job. The link up to the Business layer and across to the Data "
          "layer is the whole value. Without it, you have an IT asset list, not an architecture.")

qdm_slide('On the Technology layer, the mistake runs the other way — too deep, too early', 3,
          'What does everything run on?',
          'A short list of technology standards and a simple infrastructure picture.',
          'Going too deep, too early.',
          T,
          "VO: Fourth, the Technology layer. The question: what does all of this run on? Networks, "
          "hosting, the identity platform, the security controls, the data-exchange backbone. The "
          "deliverable is a short list of technology standards and a simple infrastructure picture — "
          "not a server-by-server audit. The mistake here is the opposite of the others. First-time "
          "architects go too deep, too early. They document every server and switch before they "
          "understand a single capability. At the architecture level, the Technology layer answers a "
          "few questions: what are we standardising on, where are the single points of failure, and "
          "what must be running for anything else to work. The detailed audit belongs to the "
          "operations team, not the architect.\n\nThe contrast with the other three mistakes is the "
          "memorable beat — name it out loud.")

trace_stack('The layers connect downward — a layer that floats free is your first gap',
            [('SERVICE', 'What a citizen or another body receives'),
             ('CAPABILITY', 'Business — what the body can do'),
             ('APPLICATION', 'The software that supports the capability'),
             ('DATA DOMAIN', 'The information it uses, one owner each'),
             ('TECHNOLOGY', 'What all of it runs on')],
            'Every element connects up and down. A layer that floats free is your first gap.',
            T,
            "VO: The four layers are not four separate documents. They connect. Every service traces "
            "to a capability in the Business layer. Every capability is served by one or more "
            "applications. Every application uses one or more data domains. Every data domain and "
            "application runs on the technology layer. When you can trace a citizen-facing service "
            "all the way down to the infrastructure it depends on, you are reading the government as "
            "a system — which is the whole job. When a layer floats free, with no connection up or "
            "down, that is the first gap you have found.\n\nCentrepiece of 2.1 — reveal one band at "
            "a time, top to bottom, then hold the closing line.")

sources_slide(prs, T, [
    'PAERA v1.0 — Annex 2 (Metamodel)',
    'PAERA v1.0 — §2.3 (Role of Enterprise Architecture)',
])


# ================================================================ 2.2
T = '2.2 · The shared vocabulary that makes re-use possible'
section('2.2', 'The shared vocabulary that makes re-use possible',
        "The metamodel is the small set of entities — Capability, Service, Application, Data Domain, "
        "Technology Component — and the relationships between them that PAERA already defines. Adopt "
        "it, and two ministries' architectures can be compared, connected and re-used. Skip it, and "
        "every team draws a different picture that no one else can read.",
        '~5 minutes',
        "VO, slide 1: Here is a problem you will hit in your first month as an architect. Two "
        "ministries each hand you an architecture. One calls a thing a 'service'. The other calls "
        "the same thing a 'function'. One ministry's 'application' is another's 'system' is a "
        "third's 'platform'. You cannot compare them. You cannot connect them. You cannot even tell "
        "whether they are doing the same work twice. The pictures exist, but they do not fit together.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'A metamodel is a small shared dictionary')
rows_slide(s, [
    ('Capability', 'Something a public body can do — register a learner, run an examination.'),
    ('Service', 'How a capability reaches a citizen or another body.'),
    ('Application', 'The software that supports a capability.'),
    ('Data Domain', 'A kind of information, with one owner.'),
    ('Technology Component', 'The infrastructure underneath.'),
    ('Organisation', 'Owns each of these.'),
], numbered=False, head_size=18, sub_size=15)
footer(s, T)
notes(s, "VO: The fix is a metamodel. A metamodel is a small, shared dictionary. It names the kinds "
         "of boxes everyone is allowed to draw, and it defines each one in plain words. PAERA "
         "defines them for you. A Capability — something a public body can do. A Service — how that "
         "capability reaches a citizen or another body. An Application — the software that supports "
         "a capability. A Data Domain — a kind of information, with one owner. A Technology "
         "Component — the infrastructure underneath. And the Organisation that owns each of these. "
         "Five or six entity types. That is most of the metamodel.\n\nThe core reference slide of "
         "the video — keep it on screen while the definitions are read.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'The relationships matter as much as the entities')
rows_slide(s, [
    ('Capability  →  delivered by  →  Service', ''),
    ('Capability  →  supported by  →  Application', ''),
    ('Application  →  uses  →  Data Domain', ''),
    ('Everything  →  runs on  →  Technology Component', ''),
    ('Organisation  →  owns  →  each of them', ''),
], numbered=False, head_size=20, top=1.6, bottom=6.15)
tb = box(s, 0.72, 6.4, 11.9, 0.5)
set_text(tb.text_frame, [[('These relationships are fixed. You do not invent them per ministry.',
                           15.5, True, ITU_BLUE_DARK, False)]])
footer(s, T)
notes(s, "VO: The entities are half of it. The relationships are the other half. A Capability is "
         "delivered by a Service. A Capability is supported by an Application. An Application uses a "
         "Data Domain. Everything runs on a Technology Component. An Organisation owns each of "
         "these. These relationships are fixed. You do not invent them per ministry. Once every team "
         "uses the same entities and the same relationships, every ministry's picture can be laid "
         "over every other ministry's picture — and they line up.")

block_slide('Re-use is not a matter of good intentions',
            ["Everyone says they want re-use — one identity system used by many ministries, instead "
             "of five. But before the agriculture ministry can consume the identity authority's "
             "building block, both must describe that building block the same way: same entity "
             "type, same relationships, same data domain.",
             "If each ministry models in its own private language, re-use is impossible even when "
             "everyone wants it."],
            'The metamodel is the precondition for re-use.',
            T,
            "VO: This is the part that matters most. Everyone says they want re-use — one identity "
            "system, used by many ministries, instead of five. But re-use is not a matter of good "
            "intentions. Before the agriculture ministry can consume the identity authority's "
            "building block, both must describe that building block the same way — same entity type, "
            "same relationships, same data domain. If each ministry models in its own private "
            "language, re-use is impossible even when everyone wants it. The metamodel is the "
            "precondition for re-use. It is the quiet, unglamorous thing that makes the "
            "whole-of-government saving — first ministry builds it, the rest consume it — actually "
            "achievable, instead of just hoped for.\n\nThe pivotal slide of 2.2. Hold it.",
            punch_fill=ITU_BLUE, punch_ink=WHITE)

two_panel('The same few entities let business and IT understand each other',
          ('The head of policy says', ['“We want to register every learner once.”']),
          ('The architect answers', ['“That is one Capability,',
                                     'one Service,',
                                     'and one Data Domain owned by the learner registry.”']),
          'Same few entities — both sides understand them. That is a policy goal translated into an architecture.',
          T,
          "VO: The metamodel does a second job. The work of redesigning how a ministry serves "
          "citizens needs the business side — the minister, the director-general, the head of policy "
          "— and the IT side — you and your engineers — to decide together. They do not share a "
          "language. The metamodel gives them one. When the head of policy says 'we want to register "
          "every learner once', and you answer 'that is one Capability, one Service, and one Data "
          "Domain owned by the learner registry', you have translated a policy goal into an "
          "architecture — in words both sides can hold. The metamodel is the bridge. The same few "
          "entities that let two ministries connect also let business and IT understand each other.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, "You adopt the metamodel — you don't design it")
rows_slide(s, [
    ('PAERA Annex 2 publishes it',
     'The entities, the definitions, the relationships — already worked out and debated across many countries.'),
    ('Use it from day one',
     'Every model you draw connects to everyone else’s instead of standing alone.'),
    ('Extend it only where your country genuinely needs a new entity',
     'Deliberately, written down, shared — never one private dictionary per team.'),
], numbered=True)
footer(s, T)
notes(s, "VO: One more thing. You do not design the metamodel. PAERA publishes it in Annex 2 — the "
         "entities, the definitions, the relationships, already worked out and debated across many "
         "countries. Your job is to adopt it and use it from day one. If your country genuinely needs "
         "an entity PAERA does not have, you extend the metamodel deliberately, and you write the "
         "extension down so every team shares it. What you never do is let each team invent its own. "
         "The value of a shared dictionary disappears the moment two teams keep private ones.")

big_slide(prs,
          "The metamodel is the small shared dictionary that lets two ministries' architectures fit "
          "together — which is what turns re-use from a wish into a plan.",
          T,
          "VO: So the metamodel is not paperwork. It is the small shared dictionary that lets two "
          "ministries' architectures fit together, and lets business and IT understand each other. "
          "It is what turns re-use from a wish into a plan. Learn PAERA's entities and relationships, "
          "use them on every model you draw, and your architecture will connect to everyone else's "
          "instead of standing alone.")

sources_slide(prs, T, [
    'PAERA v1.0 — Annex 2 (Metamodel — entities and relationships)',
])


# ================================================================ 2.3
T = "2.3 · Adopt your principles, don't draft them"
section('2.3', "Adopt your principles, don't draft them",
        'PAERA publishes ten architectural principles, already debated across many countries. Your '
        'job is to adopt them, tailor the wording to your context, and use them to settle design '
        'arguments — not to spend your first year drafting principles from scratch.',
        '~4 minutes',
        "VO, slide 1: Every architecture team faces the same temptation early on. Someone says: let "
        "us write our country's architectural principles. A workshop is booked. Three months later "
        "there are forty draft principles, half of them contradicting each other, and no agreement. "
        "Meanwhile, no ministry has been modelled. There is a faster way. The principles already exist.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'A principle is a short rule that settles a design argument before it starts')
ex = ['Reuse before you buy; buy before you build.',
      'Data has one owner.',
      'Ask the citizen only once.']
gx, gw, gap = 0.72, 3.82, 0.22
for i, e in enumerate(ex):
    x = gx + i * (gw + gap)
    panel(s, x, 1.8, gw, 1.9, LIGHT)
    panel_text(s, x, 1.8, gw, 1.9, [[(e, 18, True, ITU_BLUE_DARK, False)]])
tb = box(s, 0.72, 4.1, 11.9, 1.6)
set_text(tb.text_frame, [
    [('A decision you make once, so your team does not re-argue it on every project.', 21, False, INK, False)],
], space_after=Pt(10))
panel(s, 0.72, 5.35, 11.9, 1.1, PANEL_GREY)
panel_text(s, 0.72, 5.35, 11.9, 1.1,
           [[('If a principle never changes a decision, it is not a principle — it is a slogan.',
              20, True, INK, False)]])
footer(s, T)
notes(s, "VO: First, be clear what a principle is for. An architectural principle is a short rule "
         "that settles a design argument before it starts. 'Reuse before you buy; buy before you "
         "build.' 'Data has one owner.' 'The citizen is asked for information only once.' A good "
         "principle is a decision you make once, so your team does not re-argue it on every project. "
         "Principles are working tools, not wall decorations. If a principle never changes a "
         "decision, it is not a principle — it is a slogan.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'PAERA gives you ten — already debated across many countries')
tens = ['Rule of law', 'Intrinsic security and privacy', 'Openness and transparency',
        'Interoperability by default', 'Once-only', 'User-centred and inclusive',
        'Reuse and sharing', 'Data as a managed asset — one source of truth',
        'Technology neutrality — avoid lock-in', 'Sustainability']
top, rh = 1.62, 1.02
for i, p_ in enumerate(tens):
    col, row = i // 5, i % 5
    x = 0.72 + col * 6.15
    y = top + row * rh
    num_chip(s, x, y + 0.06, i + 1)
    tb = box(s, x + 0.5, y, 5.4, rh - 0.12)
    set_text(tb.text_frame, [[(p_, 17, True, INK, False)]])
    if row < 4:
        hline(s, x, y + rh - 0.1, 5.75)
footer(s, T)
notes(s, "VO: PAERA publishes ten architectural principles in section 5.2. They cover the ground "
         "every public-sector architecture needs. Rule of law, so every system has a legal basis. "
         "Intrinsic security and privacy, built in, not added later. Openness and transparency. "
         "Interoperability by default. Once-only — the state asks a citizen for the same information "
         "only once. User-centred and inclusive design. Reuse and sharing before building new. Data "
         "as a managed asset, with one source of truth. Technology neutrality, to avoid lock-in. And "
         "sustainability, so what you build can be maintained. Ten principles, already debated across "
         "many countries, ready to adopt.\n\nReveal cumulatively, one row at a time, in the order "
         "spoken. This is the reference payload of the video.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Adopting them is three steps, not a three-month workshop')
rows_slide(s, [
    ('Adopt the ten as your baseline', 'You inherit the thinking instead of starting from a blank page.'),
    ('Tailor the wording to your country',
     'Point each principle at your own laws, your data-protection act, your procurement rules, so it has teeth in your context.'),
    ('Add at most a few country-specific principles',
     'Only where PAERA genuinely does not cover something — deliberately, with a reason written down.'),
], numbered=True)
footer(s, T)
notes(s, "VO: Adopting them is three steps, not a three-month workshop. Step one: take the ten as "
         "your baseline. Step two: tailor the wording to your country — point each principle at your "
         "own laws, your data-protection act, your procurement rules, so it has teeth in your "
         "context. Step three: add at most a few principles your country genuinely needs that PAERA "
         "does not cover, and add them deliberately, with a reason written down. That is the whole "
         "job. You inherit the thinking; you localise the wording. You do not start from a blank page.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Make a principle bite — statement, reason, implication')
bands = [
    ('STATEMENT', 'We consume an existing building block before we build a new one.', LIGHT, INK, ITU_BLUE_DARK),
    ('REASON', 'The country pays once, not many times.', LIGHT, INK, ITU_BLUE_DARK),
    ('IMPLICATION', 'Any project proposing to build its own identity must first show why the shared '
                    'one cannot be used — and the EA Board can say no.', ITU_BLUE, WHITE, WHITE),
]
y = 1.75
for label, text, fill, ink, lab_ink in bands:
    h = 1.6 if label == 'IMPLICATION' else 1.35
    panel(s, 0.72, y, 11.9, h, fill)
    tb = box(s, 0.95, y + 0.15, 2.4, 0.4)
    set_text(tb.text_frame, [[(label, 11, True, lab_ink, False)]])
    tb = box(s, 3.4, y + 0.12, 8.95, h - 0.25)
    set_text(tb.text_frame, [[(text, 19 if label == 'IMPLICATION' else 18, label == 'IMPLICATION', ink, False)]])
    y += h + 0.17
tb = box(s, 0.72, 6.45, 11.9, 0.4)
set_text(tb.text_frame, [[('A principle with a written implication can settle an argument. A principle without one cannot.',
                           14.5, False, GREY, False)]])
footer(s, T)
notes(s, "VO: One discipline makes principles useful instead of decorative. For each one, write "
         "three things: the statement, the reason, and the implication — what the principle forces "
         "you to do, or to refuse. Take 'reuse before build'. Statement: we consume an existing "
         "building block before we build a new one. Reason: the country pays once instead of many "
         "times. Implication: any project proposing to build its own identity or payment function "
         "must first show why the shared one cannot be used — and the EA Board can say no. A "
         "principle with a written implication can settle an argument. A principle without one "
         "cannot.\n\nThe implication band is the operative beat of 2.3 — the only full-colour block "
         "in the video. Hold it.")

big_slide(prs,
          'Principles are decisions you make once so your team does not re-argue them on every '
          'project — and PAERA already made the first ten for you.',
          T,
          "VO: So do not spend your first year writing principles. Adopt PAERA's ten, point each one "
          "at your country's laws, give each a written implication so it can settle a real argument, "
          "and add your own only where there is a genuine gap. Principles are decisions you make "
          "once so your team does not re-argue them on every project. PAERA already made the first "
          "ten for you.")

sources_slide(prs, T, [
    'PAERA v1.0 — §5.2 (Principles)',
    'PAERA v1.0 — §3.3 (Digital Infrastructure principles)',
])


# ================================================================ 2.4
T = '2.4 · Classify any public body before you model it'
section('2.4', 'Classify any public body before you model it',
        'PAERA publishes a taxonomy of public bodies — policy unit, regulatory agency, '
        'service-delivery authority, plus supporting elements like state registries. Classify a body '
        'first, and you already know what capabilities, data and governance to expect from it — '
        'before you interview anyone.',
        '~4 minutes',
        "VO, slide 1: Before you model a government body, you should know what kind of body it is. "
        "Because the kind tells you, in advance, roughly what to expect — what it does, what data it "
        "owns, how it is governed. PAERA publishes a taxonomy that sorts public bodies into a few "
        "types. Learn it, and you walk into the first interview already knowing what questions to ask.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Three types cover most of government')
types_ = [
    ('Policy unit', ['Sets policy, owns the rules.',
                     'Does not run services at scale — it decides what the services should be.',
                     'Usually a ministry’s core.']),
    ('Regulatory agency', ['Licenses, supervises, enforces.',
                           'Holds registers of the things it regulates and records of its decisions.',
                           'Expect an appeals process.']),
    ('Service-delivery authority', ['Runs services to citizens at scale.',
                                    'Has the queues, the case files, the front-line systems.',
                                    'Expect high transaction volumes.']),
]
cw, ch, cy = 3.9, 4.5, 1.7
for i, (h_, lines) in enumerate(types_):
    x = 0.72 + i * (cw + 0.2)
    panel(s, x, cy, cw, ch, LIGHT)
    panel_text(s, x, cy, cw, ch,
               [[(h_, 19, True, ITU_BLUE_DARK, False)]] + [[(l, 15, False, INK, False)] for l in lines])
tb = box(s, 0.72, 6.35, 11.9, 0.5)
set_text(tb.text_frame, [[('Tell me which of the three a body is, and I can already guess its main '
                           'capabilities and data domains before I meet anyone.', 15, True, ITU_BLUE_DARK, False)]])
footer(s, T)
notes(s, "VO: Three main types cover most of government. A policy unit — usually a ministry's core — "
         "sets policy and owns the rules. It does not run services at scale; it decides what the "
         "services should be. A regulatory agency licenses, supervises and enforces. It holds "
         "registers of the things it regulates and records of its decisions. A service-delivery "
         "authority runs services to citizens at scale — it has the queues, the case files, the "
         "front-line systems. Tell me which of the three a body is, and I can already guess its main "
         "capabilities and its main data domains before I meet anyone.")

two_panel('Around those three sit the foundations everyone else stands on',
          ('State registries', ['The authoritative single source for a kind of thing:',
                                'person · business · land · learner.',
                                'A registry’s whole job is to be the one place the truth lives.']),
          ('Shared platforms', ['Identity, payments, data exchange.',
                                'Used across many bodies.',
                                'Not policy units, not regulators, not service authorities.']),
          'Treat them as shared — not as the private property of whichever ministry happens to host them.',
          T,
          "VO: Around those three sit the supporting elements. State registries — the authoritative "
          "single source of truth for a kind of thing: the population register, the business "
          "register, the land register, the learner registry. A registry's whole job is to be the one "
          "place the truth lives. And shared platforms — identity, payments, data exchange — used "
          "across many bodies. These are not policy units or regulators or service authorities. They "
          "are the foundations the others stand on. The taxonomy names them so you treat them as "
          "shared, not as the private property of whichever ministry happens to host them.",
          right_fill=LIGHT)

block_slide('Each type comes with an expected profile',
            ["Name a body a regulatory agency, and you expect a licensing capability, a register of "
             "the regulated, an enforcement record, and an appeals process. You walk in confirming "
             "or correcting a profile — not building one from nothing.",
             "When a body that should be a neutral registry is acting like a policy unit and shaping "
             "rules to suit itself, the taxonomy is what lets you see it."],
            'Misclassification is itself a finding.',
            T,
            "VO: Why does classifying first save you time? Because each type comes with an expected "
            "profile. Name a body a regulatory agency, and you expect a licensing capability, a "
            "register of the regulated, an enforcement record, and an appeals process. You walk in "
            "looking for those, and you spend the interview confirming or correcting a profile — not "
            "building one from nothing. Misclassification is itself a finding. When a body that "
            "should be a neutral registry is acting like a policy unit and shaping rules to suit "
            "itself, the taxonomy is what lets you see it. Classification is not bureaucracy. It is "
            "the architect's fastest way to orient.",
            punch_fill=ITU_BLUE, punch_ink=WHITE)

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'It is already published — adopt it')
rows_slide(s, [
    ('PAERA §4.6, with the detail in Annex A1.2',
     'You do not invent a classification scheme for your country.'),
    ('Extend it only where your country has a body type it genuinely does not cover',
     'When you say “service-delivery authority”, every other architect in the country pictures the same thing.'),
], numbered=False, top=1.7, bottom=4.4)
footer(s, T)
notes(s, "VO: Like the metamodel and the principles, the taxonomy is already published — PAERA "
         "section 4.6, with the detail in Annex A1.2. You do not invent a classification scheme for "
         "your country. You adopt this one, and you extend it only where your country has a body type "
         "it genuinely does not cover. A shared taxonomy means that when you say 'service-delivery "
         "authority', every other architect in the country pictures the same thing. That shared "
         "picture is worth more than a bespoke scheme that is perfectly tuned to your country but "
         "understood by no one else.")

big_slide(prs,
          'Classify a body first, and you already know most of what to expect from it — so the '
          'interview confirms a profile instead of starting from a blank page.',
          T,
          "VO: So before you model, classify. Policy unit, regulatory agency, service-delivery "
          "authority, plus registries and shared platforms. Each type carries an expected profile of "
          "capabilities, data and governance. Classify first, and the modelling work starts from a "
          "head start instead of a blank page — using a taxonomy every other architect in your "
          "country already shares.")

sources_slide(prs, T, [
    'PAERA v1.0 — §4.6 (Organisational taxonomy)',
    'PAERA v1.0 — Annex A1.2 (taxonomy detail)',
])


# ================================================================ 2.5
T = '2.5 · BDAT on a real ministry — the Progressa walkthrough'
section('2.5', 'BDAT on a real ministry — the Progressa walkthrough',
        'Watch the four layers and the shared entities applied to one real education system — '
        "Progressa's ministry, learner registry, examination authority and identity authority — and "
        'the abstract method becomes a concrete picture you can reproduce on your own sector.',
        '~5 minutes',
        "VO, slide 1: Let us put the four layers on a real sector. Progressa is a demonstration "
        "country with an education system like many across the continent. It has a Ministry of "
        "Education, Youth and Sport. A national examination authority. A learner registry. A national "
        "identity authority. And a digital government authority that runs shared platforms. We will "
        "read this system in four layers, using the shared entities, exactly as you would on your own "
        "sector.\n\nOn screen: Progressa is a fictional demonstration country — say so once, here.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Classify the bodies first — that is where the head start comes from')
rows_slide(s, [
    ('Ministry of Education, Youth and Sport (MoEYS)', 'Policy unit — sets education policy, owns the rules.'),
    ('National Examination Authority (PNEA)', 'Service-delivery authority with a regulatory edge — runs examinations at scale, certifies results.'),
    ('Learner Registry (PLR)', 'State registry — the authoritative single source for who is a learner.'),
    ('National Identity Authority (PNIA)', 'State registry and shared platform — owns the person identity every sector reuses.'),
    ('Digital Government Authority (PDGA)', 'Shared-platform provider — data exchange and payments.'),
], numbered=False, head_size=17, sub_size=14.5)
footer(s, T)
notes(s, "VO: Start by classifying, because classification gives us the head start. The Ministry of "
         "Education is a policy unit — it sets education policy and owns the rules. The National "
         "Examination Authority runs examinations at scale and certifies results — a "
         "service-delivery authority with a regulatory edge. The Learner Registry is a state registry "
         "— the authoritative single source for who is a learner. The National Identity Authority is "
         "a registry and a shared platform — it owns the person identity every sector reuses. The "
         "Digital Government Authority runs shared platforms: data exchange and payments. Five "
         "bodies, classified in a minute, and we already expect what each one does.\n\nAll "
         "institutions are fictional — no emblems, no real agency names on screen.")

two_panel('Business layer — capabilities belong to bodies, services sit on top',
          ('Capabilities', ['Register a learner — PLR',
                            'Run an examination, certify a result — PNEA',
                            'Prove identity — PNIA',
                            'Set policy, fund schools — MoEYS']),
          ('Services', ['Enrol a child',
                        'Sit an examination',
                        'Receive a certificate',
                        'Transfer schools']),
          'Each service traces down to a capability owned by exactly one body. Two bodies claiming one capability is your first gap.',
          T,
          "VO: Now the Business layer — capabilities and services. Register a learner: owned by the "
          "Learner Registry. Run an examination and certify a result: the Examination Authority. "
          "Prove who a learner is: the Identity Authority. Set policy and fund schools: the Ministry. "
          "Notice that we describe what each body does — its capabilities — not how it is organised "
          "inside. The citizen-facing services sit on top: enrol a child, sit an examination, receive "
          "a certificate, transfer between schools. Each service traces down to a capability owned by "
          "exactly one body. When two bodies both claim the same capability — say, each keeping its "
          "own list of learners — you have found your first gap, just by drawing the Business layer.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Data layer — one domain, one owner, one authoritative copy')
rows_slide(s, [
    ('Person', 'Owned by PNIA. Everyone else uses the person identity; no one else mints it.'),
    ('Learner', 'Owned by PLR.'),
    ('Examination result', 'Owned by PNEA.'),
    ('School', 'Owned by MoEYS.'),
], numbered=False, top=1.6, bottom=5.2)
panel(s, 0.72, 5.4, 11.9, 1.35, ITU_BLUE)
panel_text(s, 0.72, 5.4, 11.9, 1.35, [
    [('Once-only made concrete: others consume, they do not copy.', 20, True, WHITE, False)],
    [('An examination authority keeping its own learner list is a duplicate registry — write it down.', 15, False, WHITE, False)],
])
footer(s, T)
notes(s, "VO: The Data layer names the domains and their owners. The Person domain — owned by the "
         "Identity Authority; everyone else uses the person identity, no one else mints it. The "
         "Learner domain — owned by the Learner Registry. Examination results — owned by the "
         "Examination Authority. Schools — owned by the Ministry. One domain, one owner, one "
         "authoritative copy. The once-only principle now becomes concrete: when the Examination "
         "Authority needs to know who a learner is, it consumes the Learner Registry and the Identity "
         "Authority — it does not keep its own private copy of the learner that drifts out of date. "
         "If you find the Examination Authority maintaining its own learner list, that is a duplicate "
         "registry, and you write it down.\n\nThe duplicate-registry example previews the Assess gaps "
         "in 2.6 — name the link.")

two_panel('Application and Technology — every system points up and across',
          ('Applications', ['Enrolment system → register-a-learner (PLR)',
                            'Exam-management system → run-an-examination (PNEA)',
                            'Identity-verification service → prove-identity (PNIA)']),
          ('Technology', ['Shared identity platform',
                          'Data-exchange backbone (PDGA)',
                          'Hosting']),
          'A short list of shared standards — not a server audit.',
          T,
          "VO: The Application layer maps software to those capabilities. An enrolment system supports "
          "the register-a-learner capability and uses the Learner and Person domains. An "
          "examination-management system supports run-an-examination. An identity-verification "
          "service supports prove-identity. Each application points up to a capability and across to "
          "the data domains it uses — so you can see, at a glance, which systems would break if the "
          "Identity Authority changed, and whether two systems are quietly doing the same job. "
          "Underneath, the Technology layer: the shared identity platform, the data-exchange backbone "
          "run by the Digital Government Authority, the hosting. A short list of shared standards — "
          "not a server audit.\n\nApplication names stay generic and descriptive — no vendor or "
          "product names on screen.",
          right_fill=LIGHT)

trace_stack('Trace one service all the way down',
            [('SERVICE', 'Sit an examination, get a certificate'),
             ('CAPABILITY', 'Run-an-examination, certify-a-result — PNEA'),
             ('APPLICATION', 'Exam-management system'),
             ('DATA', 'Learner · person · examination result — three owners'),
             ('TECHNOLOGY', 'Data exchange · identity platform')],
            'That single thread, from a citizen service down to the infrastructure, is a complete reading of the system.',
            T,
            "VO: Now trace one service all the way down. Sit an examination and get a certificate. "
            "The service is delivered by the Examination Authority's run-an-examination and "
            "certify-a-result capabilities. Those are supported by the examination-management "
            "application. That application uses three data domains — the learner, the person "
            "identity, and the examination result — each owned by a different body and reached over "
            "the data-exchange backbone, which runs on the shared technology layer. That single "
            "thread, from a citizen service down to the infrastructure, is a complete reading of the "
            "system in four layers. Reproduce that on your own sector — classify the bodies, draw the "
            "capabilities, name the data owners, map the applications, list the shared technology — "
            "and you have an architecture, not an inventory.\n\nCentrepiece of 2.5. Reveal one band "
            "at a time and say the owner aloud at the DATA band.")

big_slide(prs,
          'Classify the bodies, draw the capabilities, name one owner per data domain, map '
          'applications to capabilities, list shared technology — and any sector becomes a connected '
          'picture.',
          T,
          "VO: That is BDAT on a real ministry. The method is the same on health, on agriculture, on "
          "social protection — only the bodies and the domains change. Classify, then read the four "
          "layers using the shared entities, and any sector in your country becomes a connected "
          "picture you can plan against.")

sources_slide(prs, T, [
    'PAERA v1.0 — Annex 2 (Metamodel)',
    'PAERA v1.0 — §5.2 (Principles — once-only)',
])


# ================================================================ 2.6
T = '2.6 · Run a Phase 2 Assess'
section('2.6', 'Run a Phase 2 Assess',
        'A good current-state picture is judged by a few quality tests per layer, not by its length. '
        'Learn the tests, learn the gaps you will always find, and you can run a Phase 2 Assess that '
        'names the right problems in the right order.',
        '~5 minutes',
        "VO, slide 1: The Assess phase produces the current-state picture and the gap analysis your "
        "country's roadmap is built on. The hard part is not writing a lot. It is writing a "
        "description good enough to make decisions from. So you need to know what a good current-state "
        "picture looks like — the quality tests, layer by layer — and the gaps you will almost always "
        "find. That is what lets you assess, instead of just document.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Three tests apply to every layer')
rows_slide(s, [
    ('Complete enough to decide — not exhaustive',
     'A description covering the capabilities and systems that matter, with the rest noted, beats a five-hundred-page audit no one reads.'),
    ('Owned — every element has a named owner',
     'An element with no owner is a finding, not a detail.'),
    ('Traceable — every element connects up and down',
     'A service that traces to no capability, an application that traces to no data, a capability no application supports: each broken link is a gap.'),
], numbered=True, bottom=6.4)
tb = box(s, 0.72, 6.5, 11.9, 0.4)
set_text(tb.text_frame, [[('If your picture fails these three tests, more pages will not fix it.', 14.5, False, GREY, False)]])
footer(s, T)
notes(s, "VO: Three tests apply to every layer. First, complete enough to decide — not exhaustive. A "
         "description that covers the capabilities and systems that matter, with the rest noted, "
         "beats a five-hundred-page audit no one reads. Second, owned — every capability, every data "
         "domain, every application has a named owner. An element with no owner is a finding, not a "
         "detail. Third, traceable — every element connects up and down. A service that traces to no "
         "capability, an application that traces to no data, a capability no application supports: "
         "each broken link is a gap. If your current-state picture passes these three tests, it is "
         "good enough to assess against. If it fails them, more pages will not fix it.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Then hold each layer to its own test')
rows_slide(s, [
    ('Business — capabilities, not org boxes',
     'Good when you could swap two ministers and the capability map would not change.'),
    ('Data — one owner, one authoritative copy per domain', ''),
    ('Application — every system maps to a capability and a data domain', 'No orphan systems.'),
    ('Technology — standards in use and single points of failure named', 'Not a list of every server.'),
], numbered=False)
footer(s, T)
notes(s, "VO: Within that, each layer has its own test. The Business layer is good when it describes "
         "capabilities, not organisation boxes — when you could swap two ministers and the capability "
         "map would not change. The Data layer is good when every domain has exactly one owner and "
         "one authoritative copy named. The Application layer is good when every system maps to a "
         "capability and a data domain — no orphan systems. The Technology layer is good when it "
         "names the standards in use and the single points of failure — not when it lists every "
         "server. Hold each layer to its own test, and the quality of the whole picture takes care of "
         "itself.\n\nThe 'swap two ministers' line is the memorable beat — land it.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Four gaps show up in almost every first assessment — look for them on purpose')
rows_slide(s, [
    ('Duplicate registries',
     'Several bodies each keeping their own copy of the same domain, none agreeing.'),
    ('Orphan systems',
     'Applications that map to no current capability, often left over from a project that ended.'),
    ('Point-to-point spaghetti',
     'Every system connected to every other by its own custom link, with no shared data exchange.'),
    ('No clear owner',
     'A capability or a domain that everyone uses and no one owns — where accountability quietly disappears.'),
], numbered=True)
footer(s, T)
notes(s, "RETRIEVAL MOMENT: before revealing, ask the viewer to predict which gaps they would expect "
         "to find in their own sector. Hold two seconds. Then reveal the four, one at a time — the "
         "answer is delivered on this slide.\n\nVO: Now the gaps. Four of them show up in almost every "
         "first assessment, so look for them on purpose. One: duplicate registries — several bodies "
         "each keeping their own copy of the same domain, none agreeing. Two: orphan systems — "
         "applications that map to no current capability, often left over from a project that ended. "
         "Three: point-to-point spaghetti — every system connected to every other by its own custom "
         "link, with no shared data exchange. Four: no clear owner — a capability or a domain that "
         "everyone uses and no one owns, which is where accountability quietly disappears. You will "
         "find these. Naming them is most of the Assess.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'A list of gaps is not an assessment — the priority order is')
rows_slide(s, [
    ('For each gap, judge how much it hurts',
     'The cost to the country, the burden on the citizen, the risk to the minister’s programme.'),
    ('And judge how hard it is to close',
     'A per-layer maturity scorecard helps — how far each layer is from where it needs to be.'),
], numbered=False, top=1.6, bottom=4.6)
panel(s, 0.72, 4.85, 11.9, 1.5, LIGHT)
panel_text(s, 0.72, 4.9, 11.9, 1.5, [
    [('Priority = high impact, where movement is possible.', 22, True, ITU_BLUE_DARK, False)],
    [('What the EA Board signs off is not “here is everything wrong”. It is “here are the right problems, in the right order, with the reasons”.', 15, False, INK, False)],
])
footer(s, T)
notes(s, "VO: A list of gaps is not an assessment. The assessment is the priority order. For each "
         "gap, judge two things: how much it hurts — the cost to the country, the burden on the "
         "citizen, the risk to the minister's programme — and how hard it is to close. The gaps you "
         "put at the top are the ones that hurt the most where movement is actually possible. A "
         "maturity scorecard per layer helps — a simple rating of how far each layer is from where it "
         "needs to be. The output your EA Board signs off is not 'here is everything wrong'. It is "
         "'here are the right problems, in the right order, with the reasons'.")

block_slide('The sign-off has one quality test of its own — honesty',
            ["Assess ends when the senior decision-maker confirms the gap analysis reflects ground "
             "truth. Your job is to name the right problems in the right order — including the "
             "politically uncomfortable ones — in language the decision-maker can act on."],
            'An assessment that softens a problem because a powerful ministry owns one of the copies '
            'fails quietly, a year later.',
            T,
            "VO: The Assess phase ends with a sign-off: the senior decision-maker confirms the gap "
            "analysis reflects ground truth. That sign-off has one quality test of its own — honesty. "
            "An assessment that flatters the current state, that softens the duplicate-registry "
            "problem because a powerful ministry owns one of the copies, fails — quietly, and "
            "expensively, a year later. Your job in Assess is to name the right problems in the right "
            "order, including the politically uncomfortable ones, in language the decision-maker can "
            "act on. Get that sign-off honestly, and the roadmap that follows stands on solid "
            "ground.\n\nThe emotional peak of the module — the only full-colour block in 2.6. Hold it "
            "a beat longer.",
            punch_fill=ITU_BLUE, punch_ink=WHITE, lead_size=21)

big_slide(prs,
          'A good current-state picture is owned, traceable, and complete enough to decide from — and '
          'the Assess is finished when it names the right problems in the right order.',
          T,
          "VO: So a Phase 2 Assess is not about volume. It is a current-state picture that is owned, "
          "traceable and complete enough to decide from, plus a gap analysis that scores and orders "
          "the problems honestly. Hold each layer to its test, look for the four gaps on purpose, and "
          "rank by impact where movement is possible. Do that, and you can run the Assess — which is "
          "the work the whole roadmap depends on.")

sources_slide(prs, T, [
    'PAERA v1.0 — §3.1.3 (Readiness Assessment)',
    'PAERA v1.0 — §5.4 (Organisational Assessment & Roadmap)',
])


# ================================================================ 2.7
T = '2.7 · The two traps to catch at Assess'
section('2.7', 'The two traps to catch at Assess — bespoke and vendor-driven',
        'Two traps recur in every assessment: the bespoke trap, where each project builds its own '
        "version of a shared function, and the vendor-driven trap, where a supplier's product quietly "
        'becomes the architecture. Learn to spot both at Assess, and you protect the country from '
        'paying many times for one thing.',
        '~4 minutes',
        "VO, slide 1: Two traps catch governments again and again. As the architect at the Assess "
        "phase, you are the one positioned to spot them early — before they are built, while they are "
        "still a line in a project plan. Learn to recognise both, and you save your country years and "
        "a great deal of money.")

block_slide('Trap one — building your own is rational for a project and ruinous for a country',
            ["A new project needs to identify citizens. Reusing the national platform means learning "
             "it, negotiating with the body that owns it, and accepting their timelines. Building its "
             "own small identity function is faster — for this project. So it builds its own.",
             "This is not laziness. Inside the project, building is the rational choice, every time. "
             "Multiply it across ten projects and the country has ten identity functions, ten learner "
             "lists, ten payment integrations."],
            'The country has paid ten times for what it should have built once.',
            T,
            "VO: The first is the bespoke trap. A new project needs to identify citizens. Reusing the "
            "national identity platform means learning it, negotiating with the body that owns it, "
            "and accepting their timelines. Building its own small identity function is faster — for "
            "this project. So the project builds its own. This is not laziness. Inside the project, "
            "building is the rational choice, every time. But multiply it across ten projects and the "
            "country has ten identity functions, ten learner lists, ten payment integrations — and "
            "has paid ten times for what it should have built once. The math that makes reuse worth "
            "it only exists at the level of the whole government. No single project can see it. You "
            "can. At Assess, every time a project proposes to build a function that already exists as "
            "a shared building block, you flag it.",
            punch_fill=ITU_BLUE, punch_ink=WHITE)

two_panel('Procurement rules cannot catch it — that is why this is architecture work',
          ('What a rule can do', ['Require open standards.',
                                  'Set the intention.']),
          ('What only architecture can do', ['Show the whole-of-government view where the reuse math exists.',
                                             'Give a Board the authority to say: not this one, reuse the shared platform.']),
          'A rule cannot make reusing another body’s platform cheaper or faster than building fresh — so the project, optimising for its deadline, still builds.',
          T,
          "VO: You might think procurement rules already prevent this. They do not. A rule can require "
          "a project to use open standards. It cannot make reusing another body's platform cheaper or "
          "faster than building fresh — so the project, optimising for its deadline, still builds. The "
          "only thing that catches the bespoke trap is the whole-of-government view that an "
          "architecture gives you, plus a governance board with the authority to say: not this one, "
          "reuse the shared platform. The rule sets the intention. The architecture and the Board "
          "enforce it. That is why catching this trap is architecture work, not procurement work.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Trap two — a product becomes the architecture one step at a time')
steps = ['A ministry buys a product to solve one problem. The product works.',
         'Slowly, more processes are bent to fit it.',
         'Data is stored the way the product wants.',
         'Other systems integrate to the product, not to a standard.']
top, rh = 1.6, 0.98
for i, st in enumerate(steps):
    y = top + i * rh
    num_chip(s, 0.72, y + 0.13, i + 1)
    tb = box(s, 1.26, y + 0.05, 11.2, rh - 0.15)
    set_text(tb.text_frame, [[(st, 18, False, INK, False)]])
    if i < len(steps) - 1:
        hline(s, 0.72, y + rh - 0.06, 11.9)
panel(s, 0.72, 5.65, 11.9, 1.3, ITU_BLUE)
panel_text(s, 0.72, 5.68, 11.9, 1.3, [
    [('Five years on, the product is the architecture — and only the vendor understands it.', 21, True, WHITE, False)],
])
footer(s, T)
notes(s, "VO: The second is the vendor-driven trap. A ministry buys a product to solve one problem. "
         "The product works. Slowly, more processes are bent to fit it. Data is stored the way the "
         "product wants. Other systems integrate to the product, not to a standard. Five years on, "
         "the product is not a system the government owns — it is the architecture, and the only "
         "people who understand it work for the vendor. Changing anything means calling them and "
         "paying what they ask. The government has lost the ability to leave, and the price reflects "
         "it. The trap is not buying from a vendor — sometimes buying is right. The trap is letting "
         "the vendor's product, rather than your architecture, decide how your government is "
         "shaped.\n\nReveal the four steps one at a time, then the closing block. The 'only the "
         "vendor understands it' line is the operative beat.")

s = add_slide(prs, LAYOUT_WHITE)
title(s, 'Four questions catch both — asked at Assess, while both are still cheap to change')
rows_slide(s, [
    ('Does a shared building block already exist for this?', ''),
    ("Is the data stored to an open standard, or to one vendor's format?", ''),
    ('If this supplier doubled their price, could we replace them within two years — and if not, why not?', ''),
    ('Is the sourcing choice deliberate — build, buy, share, sandbox — or the path of least resistance?', ''),
], numbered=True, head_size=18)
footer(s, T)
notes(s, "VO: You catch both with a few questions, asked at Assess. Does a shared building block "
         "already exist for what this project wants to build? Is the data stored to an open standard, "
         "or to one vendor's format? If this supplier doubled their price, could we replace them "
         "within two years — and if not, why not? And for anything new: is the sourcing choice "
         "deliberate — build, buy, share another country's, or test in a sandbox first — or is it "
         "just the path of least resistance? These questions turn both traps from things you discover "
         "too late into things you flag while they are still cheap to change.")

big_slide(prs,
          "Build-your-own is rational for a project and ruinous for a country; a vendor's product "
          'should fit your architecture, not become it — and Assess is where you catch both.',
          T,
          "VO: So watch for the two traps. The bespoke trap — rational for a project, ruinous for a "
          "country, caught only by the whole-of-government view and a board that can say no. And the "
          "vendor-driven trap — where a product quietly becomes the architecture and the exit door "
          "closes. Both are cheap to fix at Assess and expensive to fix later. Spotting them early is "
          "one of the most valuable things you do as an architect.")

sources_slide(prs, T, [
    'PAERA v1.0 — §1.3 (GovStack Vision)',
    'PAERA v1.0 — §3.3 (Digital Infrastructure principles)',
    'PAERA v1.0 — §5.2 (Principles)',
    'PAERA v1.0 — §5.6 (Sourcing Strategy — build / buy / share / sandbox)',
])


# ================================================================ Thank you
s = add_slide(prs, LAYOUT_THANKS)
notes(s, 'Closing slide for the combined deck. Individual videos end on their sources slide instead.')

# Self-check: the split spec's slide ranges depend on this count, and a helper that
# silently stops drawing shows up first as a slide with no voice-over.
assert len(prs.slides._sldIdLst) == 55, 'slide count changed — update decks/split_spec.json'
assert all(sl.has_notes_slide and sl.notes_slide.notes_text_frame.text.strip() for sl in prs.slides), \
    'every slide carries its voice-over in the notes'

OUT = os.environ.get('OUT_PATH') or os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    'videos', 'module_2', 'en', 'decks', 'KP1_M2_Deck_v0.1.pptx')
prs.save(OUT)
print('slides:', len(prs.slides._sldIdLst))
print('saved', OUT)
